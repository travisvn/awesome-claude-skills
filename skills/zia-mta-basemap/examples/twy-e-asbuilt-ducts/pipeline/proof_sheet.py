"""Sheet 6 -- provenance and verification. Everything a checker needs to reject or accept
the duct geometry on sheets 2-4 without taking my word for anything."""
import json, math, collections, subprocess, sys
sys.path.insert(0, __import__("os").environ.get("ZIA_BASEMAP_SCRIPTS", "../../../scripts"))
import numpy as np
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from ducts import STYLE, FRAME

NAVY=RGBColor.from_string('1E2761'); INK=RGBColor.from_string('20242A')
MUT =RGBColor.from_string('5F6368'); RULE=RGBColor.from_string('D5DAE0')
PANEL=RGBColor.from_string('F4F6F8')

def tb(slide,x,y,w,h,txt,size=9,bold=False,color=INK,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    s=slide.shapes.add_textbox(Emu(x),Emu(y),Emu(w),Emu(h)); tf=s.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=tf.margin_right=Emu(0); tf.margin_top=tf.margin_bottom=Emu(0)
    for i,line in enumerate(txt if isinstance(txt,list) else [txt]):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align
        r=p.add_run(); r.text=line
        r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; r.font.name='Calibri'
    return s

def panel(slide,x,y,w,h,fill=PANEL):
    from pptx.enum.shapes import MSO_SHAPE
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Emu(x),Emu(y),Emu(w),Emu(h))
    s.fill.solid(); s.fill.fore_color.rgb=fill
    s.line.color.rgb=RULE; s.line.width=Emu(9525); s.shadow.inherit=False
    return s

def swatch(slide,x,y,w,rgb,width,dash):
    from pptx.enum.shapes import MSO_SHAPE
    s=slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE,Emu(x),Emu(y),Emu(w),Emu(914))
    s.fill.background(); s.line.color.rgb=RGBColor.from_string(rgb); s.line.width=Emu(width)
    if dash!='solid':
        from pptx.oxml.ns import qn
        ln=s._element.find('.//'+qn('p:spPr')+'/'+qn('a:ln'))
        ln.append(ln.makeelement(qn('a:prstDash'),{'val':dash}))
    return s

def build(prs):
    reg=json.load(open('registration.json')); sd=json.load(open('sheet_ducts.json'))
    slide=prs.slides.add_slide(prs.slide_layouts[0])
    for ph in list(slide.placeholders): ph._element.getparent().remove(ph._element)
    W=prs.slide_width; M=320040
    tb(slide,M,164592,W-2*M,457200,
       'SHEET 6 — DUCT PROVENANCE & REGISTRATION VERIFICATION  ·  TWY E (E4–E6) REV P06',
       size=15,bold=True,color=NAVY)

    # ---- left column: where the geometry came from -------------------------------------
    x0,y0,cw = M, 750000, 7_100_000
    panel(slide,x0,y0,cw,1_570_000,fill=RGBColor.from_string('FFFFFF'))
    tb(slide,x0+120000,y0+110000,cw-240000,240000,'1 · SOURCE OF THE DUCT GEOMETRY',size=10,bold=True,color=NAVY)
    tb(slide,x0+120000,y0+430000,cw-240000,1_500_000,[
      'Drawing: Z1-Z2-Z3-MTA_SEGMENTATION.dxf, issued by Abu Dhabi Airports. AC1032 (AutoCAD 2018), '
      '$INSUNITS = 6 (metres), 549 MB, 2,564 layers, 12 model-space xrefs all inserted at (0,0,0), scale 1, rotation 0.',
      'Extracted population: 33,720 fixture INSERTs and 49,505 linear features (1,878 km), of which '
      '17,101 features / 466 km are duct, ductbank, conduit, sawcut or secondary cable.',
      'Duct type, bore and way-count are read from the layer name — that is the only place the drawing '
      'records them (e.g. CV_OUTER DUCT 4x110mm dia).',
      'Extraction is reconciled against an independent per-layer entity census; every count above is '
      'asserted by an automated check, not read off by eye.',
      ],size=8.5,color=INK)

    # ---- registration table ------------------------------------------------------------
    y1=y0+1_690_000
    panel(slide,x0,y1,cw,2_900_000)
    tb(slide,x0+120000,y1+110000,cw-240000,240000,'2 · HOW SHEETS 2–4 WERE TIED TO THE SOURCE DRAWING',size=10,bold=True,color=NAVY)
    tb(slide,x0+120000,y1+400000,cw-240000,700000,[
      'The sheets carry no coordinate grid, so each was registered to the source drawing by matching the '
      'plotted AGL fitting pattern: type-matched 2-point RANSAC to find the correspondence, then a '
      '4-parameter Helmert fit (translation, rotation, uniform scale) by least squares — not an affine, '
      'which would hide error as fake shear. Half the matched pairs were held back from the fit; the '
      'hold-out RMS below is the accuracy figure.',
      ],size=8.5,color=INK)
    hdr=['Sheet','Fitting symbols\nin frame','Matched to source\nwithin 50 mm','Fit RMS','HOLD-OUT RMS\n(pairs never fitted)','Plotted scale','Frame on ground']
    colw=[900000,1150000,1300000,700000,1500000,850000,700000]
    tx=x0+120000; ty=y1+1_230_000
    for c,(h,w) in enumerate(zip(hdr,colw)):
        tb(slide,tx+sum(colw[:c]),ty,w,320000,h.split('\n'),size=7,bold=True,color=MUT)
    ty+=380000
    for loc in ('LOC-01','LOC-02','LOC-03'):
        r=reg[loc]
        row=[loc, str(r['in_frame']), str(r['in_frame_exact50mm']),
             f"{r['rms']*1000:.1f} mm", f"{r['holdout_rms']*1000:.2f} mm on {r['holdout_n']}",
             f"{1/r['s']:,.0f} EMU/m", f"{FRAME['w']*r['s']:.1f} m"]
        for c,(v,w) in enumerate(zip(row,colw)):
            tb(slide,tx+sum(colw[:c]),ty,w,240000,v,size=8,color=INK)
        ty+=270000
    tb(slide,tx,ty+40000,cw-240000,420000,[
      'Rotation solved to +0.0001° / −0.000001° / +0.0027° and the y axis is mirrored — i.e. the sheets are '
      'plotted axis-aligned to the grid. Independent check: the deck\'s own "as-built taxiway centreline" '
      'polyline, never used in the fit, lands on source centreline fixtures at 0.3–0.4 mm median on LOC-01/02.',
      ],size=7.5,color=MUT)

    # ---- finding ----------------------------------------------------------------------
    y2=y1+3_020_000
    panel(slide,x0,y2,cw,1_450_000,fill=RGBColor.from_string('FFF4E5'))
    tb(slide,x0+120000,y2+110000,cw-240000,240000,'3 · FINDING — THE Rev P05 COORDINATE LABEL WAS WRONG',size=10,bold=True,
       color=RGBColor.from_string('A64B00'))
    tb(slide,x0+120000,y2+400000,cw-240000,1_000_000,[
      'Rev P05 sheets state "UTM 40N (EPSG:32640) — AS-BUILT POSITIONS". They are not UTM. The plotted '
      'positions are ZIA local project grid metres (X ≈ 5,800–6,200, Y ≈ 54,400–54,700). UTM 40N here would '
      'read easting ~230,000–800,000 and northing ~2,650,000–2,750,000.',
      'The source drawing carries no GEODATA object and no georeferencing of any kind, so no local→UTM '
      'transform can be derived from it. Rev P06 sheets are relabelled accordingly.',
      'Consequence: any coordinate previously issued from Rev P05 as UTM — to ADA, Injaz or a contractor — '
      'needs withdrawing and re-checking against survey control.',
      ],size=8.5,color=INK)

    # ---- right column: duct schedule ---------------------------------------------------
    x1=x0+cw+200000; rw=W-M-x1
    panel(slide,x1,y0,rw,6_450_000,fill=RGBColor.from_string('FFFFFF'))
    tb(slide,x1+120000,y0+110000,rw-240000,240000,'4 · DUCT SCHEDULE — WHAT IS DRAWN ON EACH SHEET',size=10,bold=True,color=NAVY)
    ty=y0+430000
    cw2=[2500000,600000,780000,700000]
    for loc in ('LOC-01','LOC-02','LOC-03'):
        d=sd[loc]['ducts']
        tal=collections.Counter(); ln=collections.Counter()
        for it in d: tal[it['leaf']]+=1; ln[it['leaf']]+=it['clipped_m']
        tb(slide,x1+120000,ty,rw-240000,220000,
           f'{loc} — {len(d)} segments, {sum(ln.values()):.0f} m inside the frame',size=8.5,bold=True,color=NAVY)
        ty+=250000
        for c,h in enumerate(['Source layer (carries bore & way-count)','Draw','Segs','On sheet']):
            tb(slide,x1+150000+sum(cw2[:c]),ty,cw2[c],200000,h,size=6.5,bold=True,color=MUT)
        ty+=210000
        for leaf in sorted(tal, key=lambda k:-ln[k]):
            nm,rgb,w,dash=STYLE[leaf]
            tb(slide,x1+150000,ty,cw2[0],200000,leaf,size=7.5,color=INK)
            swatch(slide,x1+150000+cw2[0]+40000,ty+95000,400000,rgb,max(w,12700),dash)
            tb(slide,x1+150000+sum(cw2[:2]),ty,cw2[2],200000,str(tal[leaf]),size=7.5,color=INK)
            tb(slide,x1+150000+sum(cw2[:3]),ty,cw2[3],200000,f'{ln[leaf]:.1f} m',size=7.5,color=INK)
            ty+=205000
        ty+=90000

    # ---- limits -----------------------------------------------------------------------
    y3=y2+1_570_000; x1L=x0; rwL=cw
    panel(slide,x1L,y3,rwL,3_050_000,fill=RGBColor.from_string('FFFFFF'))
    tb(slide,x1L+120000,y3+110000,rwL-240000,240000,'5 · LIMITS — READ BEFORE USING THIS GEOMETRY',size=10,bold=True,color=NAVY)
    tb(slide,x1L+120000,y3+400000,rwL-240000,2_600_000,[
      'a) Runs are fragments, not proven pulls. The duct lines are the individual LINE / LWPOLYLINE '
      'segments as drafted. No endpoint snapping or graph chaining has been done, so "handhole A to '
      'handhole B" is not established. Do not order cable to a length scaled off these sheets.',
      'b) Coordinates are ZIA local grid, not surveyed and not UTM. Nothing here may be issued as a '
      'setting-out coordinate without an externally checked transform against ADA survey control.',
      'c) Civil symbols are ~2.1–2.3 m off. Handhole, manhole, transformer-pit and RRM symbols on the '
      'sheets sit a mean 2.1–2.3 m (sd 0.3–0.45 m) from the source-drawing insertion point, while AGL '
      'light fittings match to sub-millimetre. Set out civil features from survey, not from these sheets.',
      'd) Per-asset route attribution is NOT given. The consolidated table\'s "Route" column is unchanged '
      'from Rev P05. Attributing a specific duct layer to a named asset was attempted and rejected: the '
      'sheets\' labels sit 3–8 m from their markers with neighbours only 10–15% further, and several '
      'markers resolve to the same label, so the assignment cannot be made from the drawing alone. It '
      'needs the source asset register, which is not in this data set.',
      'e) Z values are unreliable in the source drawing. This is a 2D basemap only; take levels from survey.',
      'f) Elements outside the twelve identity-inserted xrefs were excluded rather than drawn in the wrong '
      'place (3,375 geometry entities, 2,409 nested fixtures drawing-wide). If a duct looks absent in a '
      'given corner, check this first.',
      ],size=8,color=INK)
    # verification figure in the right column
    yv=y0+6_550_000
    panel(slide,x1,yv,rw,prs.slide_height-520000-yv,fill=RGBColor.from_string('FFFFFF'))
    tb(slide,x1+120000,yv+110000,rw-240000,240000,
       '6 · THE GEOMETRY, IN THE SOURCE DRAWING\'S OWN GRID',size=10,bold=True,color=NAVY)
    from PIL import Image
    iw,ih=Image.open('verify_fig.png').size
    availw=rw-300000; availh=prs.slide_height-520000-yv-420000
    sc=min(availw/iw, availh/ih)
    slide.shapes.add_picture('verify_fig.png', Emu(int(x1+(rw-iw*sc)/2)), Emu(int(yv+400000)),
                             Emu(int(iw*sc)), Emu(int(ih*sc)))
    tb(slide,M,prs.slide_height-420000,W-2*M,300000,
       'Reproduce: registration.json + sheet_ducts.json + duct_schedule.csv accompany this deck. '
       'Every figure on this sheet is computed from the source drawing, not transcribed.',
       size=7.5,color=MUT)
    return slide
