"""Rebuild each location sheet's legend as two columns:
  col 1 -- the Rev P05/P06 works-action and linework rows, harvested and REPOSITIONED so
           their styling is preserved byte-for-byte rather than recreated
  col 2 -- the as-built AGL asset types actually present on that sheet, with counts
"""
import collections
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from agl_symbols import SYM

PANEL_X, PANEL_W = 10104120, 4709160
LEGEND_TOP = 7047281
BOTTOM = 10600000
PAD = 137160
GUTTER = 60000

SHORT = {
 'TRANSFORMER HANDHOLE (SECONDARY FEED HUB)': 'TRANSFORMER HH (SEC. FEED HUB)',
 'TRANSFORMER MANHOLE (SEC. FEED)': 'TRANSFORMER MH (SEC. FEED)',
 'RUNWAY GUARD LIGHT / RRM': 'RUNWAY GUARD LIGHT / RRM',
}

# the swatch now carries the colour cue, so drop the "outer red / inner blue" tails and
# trim the linework rows to fit a half-width column
TRIM = [
 ('CORE OUT + NEW CABLE (DUCT)',        'CORE OUT + NEW CABLE (DUCT)'),
 ('CORE OUT + NEW CABLE (SAWCUT)',      'CORE OUT + NEW CABLE (SAWCUT)'),
 ('NEW SEC. CABLE ONLY (DUCT)',         'NEW SEC. CABLE ONLY (DUCT)'),
 ('NEW SEC. CABLE ONLY (SAWCUT)',       'NEW SEC. CABLE ONLY (SAWCUT)'),
 ('FIELD VERIFIED',                     'FIELD VERIFIED — NOT AFFECTED'),
 ('RRM — REMOVE',                       'RRM — REMOVE / PROTECT / RE-FIX'),
 ('AS-BUILT DUCT / DUCTBANK',           'AS-BUILT DUCT / DUCTBANK 4×110 · 6×110'),
 ('AS-BUILT SECONDARY',                 'AS-BUILT SECONDARY CONDUIT'),
 ('AGL WORKS AREA',                     'AGL WORKS AREA (GOVERNING)'),
 ('TAXIWAY CENTERLINE',                 'TAXIWAY CENTRELINE (AS-BUILT)'),
 ('STOP BAR',                           'STOP BAR (AS-BUILT)'),
 ('TWY EDGE',                           'TWY EDGE / PAVEMENT (INDICATIVE 23 m)'),
]

def _trim(sp):
    t=sp.text_frame.text.strip()
    for pre,short in TRIM:
        if t.startswith(pre):
            tf=sp.text_frame; p0=tf.paragraphs[0]
            for r in list(p0.runs)[1:]: r._r.getparent().remove(r._r)
            if p0.runs: p0.runs[0].text=short
            return

def _legend_top(slide):
    """the legend panel's own top -- it is NOT at a fixed y across the sheets, because the
    scope and notes panels above it are different heights on each one."""
    head=[sp for sp in slide.shapes
          if sp.has_text_frame and sp.text_frame.text.strip().startswith('LEGEND')]
    if not head: return LEGEND_TOP
    h=head[0]
    bgs=[sp for sp in slide.shapes
         if sp.width and sp.width>4_500_000 and sp.left is not None and sp.left>=PANEL_X
         and sp.top is not None and sp.height and sp.top<=h.top<=sp.top+sp.height]
    return min(bgs, key=lambda s:s.top).top if bgs else h.top-73152

def _rows(slide, top):
    """cluster the existing legend shapes into rows, top-down"""
    items=[sp for sp in slide.shapes
           if sp.left is not None and sp.left>=PANEL_X and sp.top is not None and sp.top>=top]
    bg=[sp for sp in items if sp.width and sp.width>4_500_000]
    head=[sp for sp in items if sp.has_text_frame and sp.text_frame.text.strip().startswith('LEGEND')]
    body=[sp for sp in items if sp not in bg and sp not in head]
    body.sort(key=lambda s:s.top)
    rows=[]
    for sp in body:
        if rows and sp.top - rows[-1][0].top < 160000: rows[-1].append(sp)
        else: rows.append([sp])
    return bg, head, rows

def _draw_symbol(slide, cx, cy, asset_type):
    txt, shape, fill, line, size, _ = SYM[asset_type]
    if shape=='ellipse':   ms=MSO_SHAPE.OVAL
    elif shape=='rect':    ms=MSO_SHAPE.RECTANGLE
    else:                  ms=MSO_SHAPE.ISOSCELES_TRIANGLE
    s=slide.shapes.add_shape(ms, Emu(int(cx-size/2)), Emu(int(cy-size/2)), Emu(size), Emu(size))
    if fill: s.fill.solid(); s.fill.fore_color.rgb=RGBColor.from_string(fill)
    else:    s.fill.background()
    s.line.color.rgb=RGBColor.from_string(line); s.line.width=Emu(9525)
    if asset_type.startswith('Existing '):
        ln=s._element.find('.//'+qn('p:spPr')+'/'+qn('a:ln'))
        ln.append(ln.makeelement(qn('a:prstDash'),{'val':'dash'}))
    s.shadow.inherit=False
    s.name=f'ASBUILT ASSET · {asset_type}'
    return s

def rebuild(slide, census, drop_texts=('EXISTING AGL ASSET','AGL FEED MANHOLE')):
    top=_legend_top(slide)
    bg, head, rows = _rows(slide, top)
    # drop the rows the typed asset block supersedes
    keep=[]
    for row in rows:
        t=''.join(sp.text_frame.text for sp in row if sp.has_text_frame).strip()
        if any(t.startswith(d) for d in drop_texts):
            for sp in row: sp._element.getparent().remove(sp._element)
        else:
            keep.append((row, t))
    present=collections.Counter(a['type'] for a in census)
    types=sorted(present, key=lambda k: SYM[k][5])

    colw=(PANEL_W-2*PAD-GUTTER)//2
    c1x, c2x = PANEL_X+PAD, PANEL_X+PAD+colw+GUTTER
    y0=top+560000
    n=max(len(keep), len(types))
    pitch=min(274320, int((BOTTOM-y0)/max(n,1)))

    # column 1: reposition harvested rows, narrow their text to the column
    for i,(row,_) in enumerate(keep):
        rtop=min(sp.top for sp in row); rleft=min(sp.left for sp in row)
        newtop=y0+i*pitch
        for sp in row:
            dy=sp.top-rtop; dx=sp.left-rleft
            sp.top=newtop+dy
            sp.left=c1x+dx
            if sp.has_text_frame and sp.text_frame.text.strip():
                _trim(sp)
                sp.width=Emu(colw-dx)
                for p in sp.text_frame.paragraphs:
                    for r in p.runs: r.font.size=Pt(7)

    # column 2: as-built asset types on this sheet
    for cx,label in ((c1x,'WORKS ACTIONS & LINEWORK'), (c2x,'AS-BUILT AGL ASSETS')):
        hd=slide.shapes.add_textbox(Emu(cx), Emu(top+330000), Emu(colw), Emu(200000))
        hd.text_frame.word_wrap=True
        tfm=hd.text_frame
        tfm.margin_left=tfm.margin_right=tfm.margin_top=tfm.margin_bottom=Emu(0)
        r=tfm.paragraphs[0].add_run(); r.text=label
        r.font.size=Pt(7.5); r.font.bold=True; r.font.color.rgb=RGBColor.from_string('5F6368')
    for i,t in enumerate(types):
        cy=y0+i*pitch+110000
        _draw_symbol(slide, c2x+110000, cy, t)
        tb=slide.shapes.add_textbox(Emu(c2x+250000), Emu(int(cy-115000)), Emu(colw-250000), Emu(230000))
        tf=tb.text_frame; tf.word_wrap=True
        tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=Emu(0)
        label=SHORT.get(SYM[t][0], SYM[t][0])
        run=tf.paragraphs[0].add_run(); run.text=f'{label}  ({present[t]})'
        run.font.size=Pt(7); run.font.color.rgb=RGBColor.from_string('20242A')
    for sp in bg:
        h=BOTTOM+120000-sp.top
        if 0 < h: sp.height=Emu(h)
    for sp in head:
        for p in sp.text_frame.paragraphs:
            for r in p.runs:
                if 'LEGEND' in r.text: r.text='LEGEND — THIS SHEET'
        sp.width=Emu(colw)
    return len(keep), len(types)
