"""Rebuild the TWY E shop-drawing deck with REAL as-built duct geometry.

Replaces the Rev P05 'INDICATIVE DUCT' construction lines (light-to-nearest-MH/HH
straight lines) with the actual duct, ductbank, conduit and sawcut polylines from the
ADA-issued source drawing, registered onto each sheet, and corrects the coordinate
system statement. Everything stays native editable PowerPoint geometry.
"""
import copy, json, math, collections
import numpy as np
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from ducts import STYLE, FRAME, inv, fwd, sheet_window

SRC='src.pptx'; OUT='TWY-E-AGL-SHOPDWG-ASBUILT-DUCTS_RevP06.pptx'
SHEETS=[('LOC-01',1,'1001'),('LOC-02',2,'1002'),('LOC-03',3,'1003')]
INDICATIVE={'C9CDD2','BB00BB'}          # per the Rev P05 legend: the two indicative duct layers

reg=json.load(open('registration.json'))
sd =json.load(open('sheet_ducts.json'))

FRAME_R = FRAME['x']+FRAME['w']

def line_shapes_to_drop(slide):
    """indicative duct construction lines: prstGeom 'line', indicative colour, INSIDE the
    plot frame -- the identically coloured legend swatches must survive."""
    out=[]
    for sp in slide.shapes:
        el=sp._element
        pg=el.find('.//'+qn('a:prstGeom'))
        if pg is None or pg.get('prst')!='line': continue
        if sp.left is None or sp.left > FRAME_R: continue
        ln=el.find('.//'+qn('p:spPr')+'/'+qn('a:ln'))
        if ln is None: continue
        clr=ln.find('.//'+qn('a:srgbClr'))
        if clr is not None and clr.get('val') in INDICATIVE: out.append(sp)
    return out

def find(slide, prefix):
    """locate a shape by the start of its text -- shape ids are not stable across slides"""
    for sp in slide.shapes:
        if sp.has_text_frame and sp.text_frame.text.strip().startswith(prefix): return sp
    raise KeyError(prefix)

def swatch_left_of(slide, textshape, colour):
    """the legend line swatch on the same row as a legend label"""
    y=textshape.top + textshape.height/2
    best=None
    for sp in slide.shapes:
        el=sp._element
        pg=el.find('.//'+qn('a:prstGeom'))
        if pg is None or pg.get('prst')!='line': continue
        if sp.left is None or sp.left < FRAME_R: continue
        if abs(sp.top - y) > 300000: continue
        ln=el.find('.//'+qn('p:spPr')+'/'+qn('a:ln'))
        clr=ln.find('.//'+qn('a:srgbClr')) if ln is not None else None
        if clr is not None and clr.get('val')==colour: best=sp
    return best

def add_polyline(slide, pts, rgb, width, dash, name):
    fb=slide.shapes.build_freeform(pts[0][0], pts[0][1], scale=1.0)
    fb.add_line_segments(pts[1:], close=False)
    sh=fb.convert_to_shape()
    sh.name=name
    sh.fill.background()
    ln=sh.line; ln.color.rgb=RGBColor.from_string(rgb); ln.width=Emu(width)
    if dash!='solid':
        lnEl=sh._element.find('.//'+qn('p:spPr')+'/'+qn('a:ln'))
        d=lnEl.makeelement(qn('a:prstDash'),{'val':dash}); lnEl.append(d)
    return sh

def restack_after(slide, shapes_, anchor_idx):
    """move newly appended shapes to just after anchor_idx so labels stay on top"""
    tree=slide.shapes._spTree
    kids=[c for c in tree if c.tag in (qn('p:sp'),qn('p:grpSp'),qn('p:pic'),qn('p:cxnSp'))]
    anchor=kids[anchor_idx]
    for sh in shapes_:
        el=sh._element; tree.remove(el); anchor.addnext(el); anchor=el

def set_text(shape, lines, size=None):
    tf=shape.text_frame
    p0=tf.paragraphs[0]
    tmpl=p0.runs[0]._r if p0.runs else None
    for p in list(tf.paragraphs)[1:]: p._p.getparent().remove(p._p)
    for r in list(p0.runs): r._r.getparent().remove(r._r)
    for i,txt in enumerate(lines):
        para = p0 if i==0 else tf.add_paragraph()
        if i>0 and tmpl is not None:
            pPr=tf.paragraphs[0]._p.find(qn('a:pPr'))
            if pPr is not None: para._p.insert(0, copy.deepcopy(pPr))
        if tmpl is not None:
            r=copy.deepcopy(tmpl); para._p.append(r)
            for t in r.findall(qn('a:t')): t.text=txt
        else:
            para.text=txt
        if size:
            for r in para.runs: r.font.size=Pt(size)

LEGEND_TOP=7047281

def respace_legend(slide, slide_h, bottom_margin=90000):
    """Rev P05 spaced the legend at 274320 EMU and the last two rows fell off the sheet.
    Re-space to whatever pitch makes every row land above the bottom edge."""
    items=[sp for sp in slide.shapes
           if sp.left is not None and sp.left>=FRAME_R and sp.top is not None and sp.top>=LEGEND_TOP]
    bg=[sp for sp in items if sp.width and sp.width>4_500_000]
    body=[sp for sp in items if sp not in bg
          and not (sp.has_text_frame and sp.text_frame.text.strip().startswith('LEGEND'))]
    body.sort(key=lambda s: s.top)
    rows=[]                                    # cluster shapes into rows by vertical gap
    for sp in body:
        if rows and sp.top - rows[-1][0].top < 160000: rows[-1].append(sp)
        else: rows.append([sp])
    keys=list(range(len(rows)))
    if len(keys)<2: return 0
    tops=[min(s.top for s in rows[k]) for k in keys]
    y0=tops[0]
    avail=slide_h-bottom_margin-y0-240000
    pitch=min(274320, int(avail/(len(keys)-1)))
    for i,k in enumerate(keys):
        shift=(y0+i*pitch)-tops[i]
        for sp in rows[k]: sp.top=sp.top+shift
    for sp in bg:
        h=slide_h-bottom_margin-sp.top
        if 0 < h < sp.height: sp.height=h
    return len(keys)

prs=Presentation(SRC)
report={}; asset_census={}
for loc, sidx, dwg in SHEETS:
    slide=prs.slides[sidx]
    r=reg[loc]; T=dict(s=r['s'],theta=r['theta'],t=np.array(r['t']),reflect=r['reflect'])

    # 1. remove the indicative duct construction lines
    drop=line_shapes_to_drop(slide)
    for sp in drop: sp._element.getparent().remove(sp._element)

    # 2. draw the real duct geometry
    added=[]; tally=collections.Counter(); length=collections.Counter()
    for it in sd[loc]['ducts']:
        nm,rgb,w,dash = STYLE[it['leaf']]
        pts=[(Emu(a),Emu(b)) for a,b in it['emu']]
        if len(pts)<2: continue
        added.append(add_polyline(slide, pts, rgb, w, dash,
                                  f"ASBUILT DUCT · {it['leaf']} · {it['clipped_m']:.1f} m"))
        tally[it['leaf']]+=1; length[it['leaf']]+=it['clipped_m']
    restack_after(slide, added, 1)          # just above the plot frame, below everything else

    # 2b. typed as-built asset layer: drop the uniform 'existing AGL asset' dots and draw
    #     each asset with its own symbol. Works-action markers are left untouched.
    import agl_symbols, legend as legend_mod
    from pptx.enum.shapes import MSO_SHAPE
    dropped_dots=0
    for sp in list(slide.shapes):
        el=sp._element
        pg=el.find('.//'+qn('a:prstGeom'))
        if pg is None or pg.get('prst')!='ellipse': continue
        if sp.left is None or sp.left>FRAME_R: continue
        sf=el.find('./p:spPr/a:solidFill/a:srgbClr',
                   {'a':'http://schemas.openxmlformats.org/drawingml/2006/main',
                    'p':'http://schemas.openxmlformats.org/presentationml/2006/main'})
        if sf is not None and sf.get('val')=='12A5B8':
            el.getparent().remove(el); dropped_dots+=1
    census=agl_symbols.sheet_assets(loc, T, f'pptx_in/ppt/slides/slide{sidx+1}.xml')
    sym_shapes=[]
    for a in census:
        if a['works_claimed']: continue        # its action symbology already marks it
        sym_shapes.append(legend_mod._draw_symbol(slide, a['x'], a['y'], a['type']))
    restack_after(slide, sym_shapes, 1+len(added))
    nlegend = legend_mod.rebuild(slide, census)
    asset_census[loc]=census

    # 3. correct the coordinate-system strip under the frame
    span=FRAME['w']*T['s']
    set_text(find(slide,'UTM 40N (EPSG:32640)'), [f'ZIA LOCAL PROJECT GRID (m) — NOT UTM, NOT SURVEYED  ·  '
                       f'FRAME = {span:.1f} m ACROSS  ·  NOT TO SCALE WHEN EDITED'])
    # 4. title-block scale line
    set_text(find(slide,'Scale 1:250'), [f'Frame {span:.1f} m across · ZIA local grid (m) · duct geometry from '
                         f'Z1-Z2-Z3-MTA_SEGMENTATION source dwg · Rev P06 · Prepared: Mohammed, '
                         f'AGL Team Leader — ADB SAFEGATE'])
    # 5. general notes
    notes=[
     '1. Positions plotted are ZIA LOCAL PROJECT GRID metres — NOT UTM 40N. The Rev P05 '
     '"EPSG:32640" label was wrong; no surveyed local-to-UTM transform exists. Do not issue a '
     'UTM coordinate from this sheet.',
     '2. AGL works area (red) per field condition. Coring at Location 1 only.',
     '3. Duct / ductbank / conduit / sawcut lines are REAL as-built geometry from ADA drawing '
     'Z1-Z2-Z3-MTA_SEGMENTATION.dxf — not light-to-nearest-pit construction lines. Layer name '
     'carries size and way-count; see sheet 6.',
     '4. Ducts are drawn as the separate segments they are in the source. Continuity NOT proven — '
     'do not take a continuous pull length off this sheet.',
     '5. Isolate, lock out & prove dead all listed circuits at CCR under permit before cutting, '
     'coring or excavation.',
     '6. HH / MH / transformer-pit / RRM symbols sit ~2.1-2.3 m off the source insertion point '
     '(sheet 6) — set out civil features from survey.',
     '7. Centerline & stop bar from as-built TCC / SBC fittings. Edge indicative 23 m — confirm '
     'on site.',
     '8. Every asset in frame carries its own as-built symbol, legended right. Type comes from '
     'the source layer, not the label. Assets with a works-action marker keep that marker.',
    ]
    nb=find(slide,'1. All asset positions as-built')
    set_text(nb, notes, size=7.5)
    nb.text_frame.word_wrap=True
    from pptx.enum.text import MSO_ANCHOR
    nb.text_frame.vertical_anchor=MSO_ANCHOR.TOP

    # 6. legend: the two freed indicative rows now describe the real duct families
    present=set(tally)
    bank=[k for k in present if 'DUCT' in k and 'SEC' not in k and 'GRND' not in k]
    sec =[k for k in present if 'SEC' in k]
    lab1=find(slide,'INDICATIVE DUCT'); sw1=swatch_left_of(slide,lab1,'C9CDD2')
    lab2=find(slide,'SECONDARY DUCT RUN CROSSING CUT'); sw2=swatch_left_of(slide,lab2,'BB00BB')
    if sw1 is not None:
        sw1.line.color.rgb=RGBColor.from_string('0B3D91'); sw1.line.width=Emu(22225)
        d=sw1._element.find('.//'+qn('a:ln')+'/'+qn('a:prstDash'))
        if d is not None: d.set('val','solid')
    set_text(lab1, ['AS-BUILT DUCT / DUCTBANK 4×110 · 6×110'])
    if sw2 is not None:
        sw2.line.color.rgb=RGBColor.from_string('00838F'); sw2.line.width=Emu(15875)
    set_text(lab2, ['AS-BUILT SECONDARY CONDUIT'])
    print(f'{loc}: dropped {len(drop)} indicative lines + {dropped_dots} uniform asset dots; '
          f'drew {len(added)} duct segments ({sum(length.values()):.1f} m) and '
          f'{len(sym_shapes)} typed as-built asset symbols; '
          f'legend {nlegend[0]} action/linework rows + {nlegend[1]} asset types')
    report[loc]=dict(dropped=len(drop), drawn=len(added), span_m=span,
                     layers={k:(tally[k],round(length[k],1)) for k in sorted(tally)})


# ---- correct the LOC-02 scope statement the as-built contradicts -----------------------
def fix_loc02_scope(slide):
    """Rev P05 described a chained secondary run. topology.py shows five separate home runs
    from one transformer handhole. Correct the statement at source, not just on sheet 7."""
    hits=0
    for sp in slide.shapes:
        if not sp.has_text_frame: continue
        for para in sp.text_frame.paragraphs:
            joined=''.join(r.text for r in para.runs)
            if 'CHAINED RUN' in joined:
                for r in list(para.runs)[1:]: r._r.getparent().remove(r._r)
                para.runs[0].text=('SECONDARY DUCT (AS-BUILT): 5 SEPARATE HOME RUNS FROM ONE '
                    'TRANSFORMER HANDHOLE — NOT A CHAINED RUN. 20.5 / 21.1 / 29.1 / 30.3 / 41.5 m '
                    '= 142.5 m TOTAL. RE-PRICE — SEE SHEET 7.')
                hits+=1
            elif 'SPUR FROM HH.E.055' in joined:
                for r in list(para.runs)[1:]: r._r.getparent().remove(r._r)
                para.runs[0].text=('  (Rev P05 chained-run assumption superseded by the source '
                    'drawing — see sheet 7)')
                hits+=1
    return hits

print('LOC-02 scope statements corrected:', fix_loc02_scope(prs.slides[2]))

# ---- rev strings and title-slide basis -------------------------------------------------
def retitle(slide):
    for sp in slide.shapes:
        if not sp.has_text_frame: continue
        for para in sp.text_frame.paragraphs:
            for r in para.runs:
                if 'P05' in r.text: r.text=r.text.replace('REV P05','REV P06').replace('P05','P06')
for s in prs.slides: retitle(s)

t=prs.slides[0]
for sp in t.shapes:
    if not sp.has_text_frame: continue
    txt=sp.text_frame.text
    if txt.startswith('As-built AGL asset survey'):
        set_text(sp,['As-built AGL fitting positions (per assets_20260726120533.xlsx). AS PLOTTED they are '
                     'ZIA LOCAL PROJECT GRID (m), matching the ADA source drawing to sub-millimetre — NOT '
                     'UTM 40N. The Rev P05 EPSG:32640 label was incorrect — see sheet 6.'])
    elif txt.startswith('Field verification sheets'):
        set_text(sp,['Field verification sheets: Document_3 (LOC-01) · Second_milling_area rev _1 (LOC-02/03)  ·  '
                     'Duct geometry: ADA drawing Z1-Z2-Z3-MTA_SEGMENTATION.dxf'])
    elif txt.startswith('EDITABLE VERSION'):
        set_text(sp,['REV P06 — the "indicative duct" construction lines of Rev P05 (straight lines from each '
                     'light to its nearest pit) are removed and replaced with the real as-built duct, ductbank, '
                     'conduit and sawcut geometry from the ADA source drawing, registered onto each sheet to '
                     'sub-millimetre. Every element remains a native editable PowerPoint shape. '
                     'Provenance, registration residuals and limits: sheet 6.'])

# footnote on the consolidated table
for sp in prs.slides[4].shapes:
    if sp.has_text_frame and sp.text_frame.text.startswith('Totals:'):
        set_text(sp,['Totals: 11 core-outs (Location 1 only) · 31 secondary cable runs · 3 RRMs · 1 field-verified '
                     'not affected.  ·  The "Route" column is carried over from Rev P05 and is NOT derived from '
                     'the as-built duct data: per-asset attribution could not be established from the sheets '
                     '(see sheet 6, limit d). Sheet-level duct schedule is on sheet 6.'], size=9)

json.dump(asset_census, open('asset_census.json','w'))
import proof_sheet, sheet7, sheet8
proof_sheet.build(prs)
sheet7.build(prs)
sheet8.build(prs)

prs.save(OUT)
json.dump(report, open('build_report.json','w'), indent=1)
print('\nwrote', OUT)
