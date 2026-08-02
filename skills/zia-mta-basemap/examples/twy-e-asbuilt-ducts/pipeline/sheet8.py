"""Sheet 8 -- master as-built AGL asset symbol schedule.

The per-sheet legends show only what is on that sheet. This is the full symbol set with
the source layer each type comes from and the count per location, so a checker can tie
any symbol on any sheet back to a layer in the source drawing.
"""
import json, collections
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from proof_sheet import tb, panel, NAVY, INK, MUT
from agl_symbols import SYM
import legend as legend_mod

LOCS=('LOC-01','LOC-02','LOC-03')

def build(prs):
    census=json.load(open('asset_census.json'))
    slide=prs.slides.add_slide(prs.slide_layouts[0])
    for ph in list(slide.placeholders): ph._element.getparent().remove(ph._element)
    W=prs.slide_height and prs.slide_width; M=320040; FW=W-2*M
    tb(slide,M,164592,FW,457200,
       'SHEET 8 — AS-BUILT AGL ASSET SYMBOL SCHEDULE  ·  TWY E (E4–E6) REV P06',
       size=15,bold=True,color=NAVY)

    y0=750000
    panel(slide,M,y0,FW,1_150_000,fill=RGBColor.from_string('FFFFFF'))
    tb(slide,M+120000,y0+110000,FW-240000,240000,'HOW TO READ THIS',size=10,bold=True,color=NAVY)
    tb(slide,M+120000,y0+420000,FW-240000,650000,[
      'Circle = light · square = civil chamber · triangle = earthing · filled rectangle = sign '
      'foundation. Lights are coloured to the light\'s own showing: taxiway centreline green, '
      'taxiway edge blue, stop bar red, runway guard amber. A dashed outline marks the source '
      'drawing\'s "existing" layer for that family.',
      'Asset type is taken from the source-drawing LAYER, never from the label on the sheet. '
      'Assets that carry a Rev P05 works-action marker keep that marker instead of a type symbol — '
      'their type is in the ID label and in this schedule. Counts are assets inside each sheet frame.',
      ],size=8.5,color=INK)

    # table
    y1=y0+1_270_000
    ph=prs.slide_height-520000-y1
    panel(slide,M,y1,FW,ph,fill=RGBColor.from_string('FFFFFF'))
    tb(slide,M+120000,y1+110000,FW-240000,240000,
       'SYMBOL  ·  ASSET TYPE  ·  SOURCE LAYER  ·  COUNT PER SHEET',size=10,bold=True,color=NAVY)
    cnt={loc:collections.Counter(a['type'] for a in census[loc]) for loc in LOCS}
    leaves={}
    for loc in LOCS:
        for a in census[loc]: leaves.setdefault(a['type'],set()).add(a['leaf'])
    types=[t for t in sorted(SYM, key=lambda k: SYM[k][5]) if any(cnt[l][t] for l in LOCS)]

    cw=[520000, 3_900_000, 4_400_000, 780000, 780000, 780000, 900000]
    hdr=['Symbol','Asset type','Source drawing layer(s)','LOC-01','LOC-02','LOC-03','Total in frame']
    tx=M+140000; ty=y1+480000
    for c,h in enumerate(hdr):
        tb(slide,tx+sum(cw[:c]),ty,cw[c],220000,h,size=7.5,bold=True,color=MUT)
    ty+=260000
    pitch=min(300000, int((ph-480000-260000-260000)/max(len(types),1)))
    tot=collections.Counter()
    for t in types:
        legend_mod._draw_symbol(slide, tx+180000, ty+pitch//2-20000, t)
        vals=[cnt[l][t] for l in LOCS]; s=sum(vals); tot.update({t:s})
        row=[None, SYM[t][0], ' · '.join(sorted(leaves.get(t,[]))),
             *[str(v) if v else '—' for v in vals], str(s)]
        for c,v in enumerate(row):
            if v is None: continue
            tb(slide,tx+sum(cw[:c]),ty+pitch//2-110000,cw[c],210000,v,size=8,color=INK)
        ty+=pitch
    grand=[sum(cnt[l][t] for t in types) for l in LOCS]
    tb(slide,tx+sum(cw[:1]),ty+30000,cw[1],220000,
       f'TOTAL — {len(types)} asset types',size=8,bold=True,color=NAVY)
    for c,v in enumerate(grand):
        tb(slide,tx+sum(cw[:3+c]),ty+30000,cw[3+c],220000,str(v),size=8,bold=True,color=NAVY)
    tb(slide,tx+sum(cw[:6]),ty+30000,cw[6],220000,str(sum(grand)),size=8,bold=True,color=NAVY)

    notdrawn={loc: sum(1 for a in census[loc] if not a['drawn_in_p05']) for loc in LOCS}
    tb(slide,M,prs.slide_height-420000,FW,300000,
       f'Of these, Rev P05 drew no symbol at all for {notdrawn["LOC-01"]} assets on LOC-01, '
       f'{notdrawn["LOC-02"]} on LOC-02 and {notdrawn["LOC-03"]} on LOC-03 — they are added in '
       f'Rev P06 from the as-built and are NOT field-verified. asset_census.json accompanies this deck.',
       size=7.5,color=MUT)
    return slide
