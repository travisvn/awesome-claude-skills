"""For each affected asset named in the consolidated table, find the nearest REAL as-built
duct/conduit run and report its source layer.

Affected assets are the concentric works-action markers on the location sheets (an unfilled
outer ring plus a filled inner dot). Marker -> asset ID comes from the nearest label text,
and is accepted only when that label is unambiguously closest AND the ID appears in the
table. Anything ambiguous is reported as such rather than guessed.
"""
import json, math, sys
sys.path.insert(0, __import__("os").environ.get("ZIA_BASEMAP_SCRIPTS", "../../../scripts"))
import numpy as np
from shapely.geometry import LineString, Point
from parse_deck import shapes
from ducts import STYLE, fwd, FRAME

RING={'CC0000','0055CC','E8710A','9AA0A6'}
FRAME_R=FRAME['x']+FRAME['w']

def markers_and_ids(slide_xml, table_ids):
    sh=shapes(slide_xml)
    rings=[s for s in sh if s['prst']=='ellipse' and s['fill'] is None
           and s['line'] in RING and s['x']<FRAME_R]
    labels=[s for s in sh if s['text'] and s['prst']=='rect' and s['x']<FRAME_R
            and s['text'].strip() in table_ids]
    out=[]
    for r in rings:
        cx,cy=r['x']+r['cx']/2, r['y']+r['cy']/2
        d=sorted(((math.dist((cx,cy),(l['x']+l['cx']/2,l['y']+l['cy']/2)), l['text'].strip())
                  for l in labels))
        if not d: continue
        best=d[0]; second=d[1] if len(d)>1 else (9e18,'')
        unambiguous = second[0] > best[0]*1.6 and best[0] < 1_200_000
        out.append(dict(px=cx,py=cy,asset=best[1],ring=r['line'],
                        unambiguous=bool(unambiguous), d1=best[0], d2=second[0]))
    return out

if __name__=='__main__':
    from pptx import Presentation
    prs=Presentation('src.pptx'); tbl=[s for s in prs.slides[4].shapes if getattr(s,'has_table',False)][0].table
    rows=[[c.text.strip() for c in r.cells] for r in list(tbl.rows)[1:]]
    ids_by_loc={}
    for loc,asset,*_ in rows: ids_by_loc.setdefault(loc,set()).add(asset)
    reg=json.load(open('registration.json')); sd=json.load(open('sheet_ducts.json'))
    result={}
    for loc,sl in (('LOC-01',2),('LOC-02',3),('LOC-03',4)):
        r=reg[loc]; T=dict(s=r['s'],theta=r['theta'],t=np.array(r['t']),reflect=r['reflect'])
        segs=[(it['leaf'], LineString(it['coords'])) for it in sd[loc]['ducts']]
        mk=markers_and_ids(f'pptx_in/ppt/slides/slide{sl}.xml', ids_by_loc.get(loc,set()))
        for m in mk:
            P=Point(fwd(T,[[m['px'],m['py']]])[0])
            if segs:
                dmin,leaf=min(((g.distance(P), lf) for lf,g in segs))
                m['route_leaf']=leaf; m['route_d']=float(dmin)
            else:
                m['route_leaf']=None; m['route_d']=None
        result[loc]=mk
        ok=[m for m in mk if m['unambiguous']]
        print(f'{loc}: {len(mk)} works-action markers, {len(ok)} mapped unambiguously to a table asset')
        for m in sorted(ok,key=lambda z:z['asset'])[:60]:
            print(f"   {m['asset']:16s} -> {m['route_leaf']:30s} at {m['route_d']:6.2f} m")
    json.dump(result, open('asset_routes.json','w'), indent=1)
