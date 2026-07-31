"""
TWY E (E4-E6) AGL shop-drawing deck - Rev P08 final-submission pass.

Checks and corrects internal consistency, removes internal/working notes,
re-words comparisons against other parties diplomatically, aligns the layout,
and applies the ADB SAFEGATE logo.
"""
import copy
import os
import sys

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.shapes.autoshape import Shape as _Shape

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
sys.path.insert(0, HERE)
import textfit

SRC = os.path.join(ROOT, "input", "TWYEAGLSHOPDWG_RevP08.pptx")
OUT = os.path.join(ROOT, "out",
                   "AUH-SK-AGL-TWYE-001 - TWY E Shop Drawings - ZIA Rev P08 FINAL.pptx")
LOGO = os.path.join(HERE, "assets", "logo-adb.png")
LOGO_REV = os.path.join(HERE, "assets", "logo-adb-reverse.png")

SW, SH = 15124113, 10688638
prs = Presentation(SRC)
S = prs.slides
changes = []


def log(msg):
    changes.append(msg)


# ----------------------------------------------------------------- helpers
def shp(slide, name, nth=0):
    hits = [s for s in slide.shapes if s.name == name]
    return hits[nth] if hits else None


def set_para(para, text, like=None):
    runs = para.runs
    if not runs:
        r = para.add_run()
        src = like or _last_run_like(para)
        if src is not None:
            r.font.size, r.font.name, r.font.bold = src.font.size, src.font.name, src.font.bold
            try:
                if src.font.color and src.font.color.type is not None:
                    r.font.color.rgb = src.font.color.rgb
            except Exception:
                pass
        runs = para.runs
    runs[0].text = text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)
    for br in para._p.findall(qn("a:br")):        # stray soft breaks from earlier revisions
        br.getparent().remove(br)


def _last_run_like(para):
    """Find a run in a sibling paragraph to copy formatting from."""
    body = para._p.getparent()
    for sib in body.iterchildren():
        for r in sib.findall(qn("a:r")):
            from pptx.text.text import _Run
            return _Run(r, None)
    return None


def set_lines(shape, lines):
    paras = shape.text_frame.paragraphs
    assert len(paras) == len(lines), (shape.name, len(paras), len(lines))
    for p, t in zip(paras, lines):
        set_para(p, t)


def sub_all(shape, pairs):
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            t = r.text
            for old, new in pairs:
                if old in t:
                    t = t.replace(old, new)
            r.text = t


def clone_para(shape, src_idx, text):
    tf = shape.text_frame
    src = tf.paragraphs[src_idx]._p
    new = copy.deepcopy(src)
    src.addnext(new)
    from pptx.text.text import _Paragraph
    para = _Paragraph(new, tf)
    set_para(para, text)
    return para


def del_shape(shape):
    shape._element.getparent().remove(shape._element)


def para_specs(shape):
    out = []
    for p in shape.text_frame.paragraphs:
        runs = p.runs
        pt, bold = 10.0, False
        if runs:
            sz = runs[0].font.size
            pt = sz.pt if sz else 10.0
            bold = bool(runs[0].font.bold)
        out.append((p.text, pt, bold, p.space_after.emu if p.space_after else 0))
    return out


def fit_height(shape):
    return textfit.frame_height(para_specs(shape), shape.width)


def place(shape, left=None, top=None, width=None, height=None):
    if left is not None:
        shape.left = Emu(int(left))
    if top is not None:
        shape.top = Emu(int(top))
    if width is not None:
        shape.width = Emu(int(width))
    if height is not None:
        shape.height = Emu(int(height))


def add_logo(slide, right_edge, top, height_emu, reverse=False):
    w = int(height_emu * 788 / 272)
    pic = slide.shapes.add_picture(LOGO_REV if reverse else LOGO,
                                   Emu(int(right_edge - w)), Emu(int(top)),
                                   Emu(w), Emu(int(height_emu)))
    pic.name = "ADB SAFEGATE LOGO"
    return pic


FOOTER_DOC = ("AUH-SK-AGL-TWYE-001 REV P08 (FINAL SCOPE)  ·  Prepared: Mohammed Faheem, "
              "AGL Team Leader — ADB SAFEGATE  ·  Approved: Ragesh Menon, AGL Manager")


def add_footer(slide, text, left, width, top, size=7.5):
    tb = slide.shapes.add_textbox(Emu(int(left)), Emu(int(top)), Emu(int(width)), Emu(190500))
    tb.name = "SHEET FOOTER"
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.name = "Arial"
    r.font.color.rgb = RGBColor(0x5F, 0x63, 0x68)
    return tb


def fit_font(shape, max_w_emu, start_pt, min_pt=9.0):
    """Shrink a single-line title until it fits max_w_emu."""
    txt = shape.text_frame.text
    pt = start_pt
    while pt > min_pt and textfit.text_w_pt(txt, pt, True) * textfit.EMU_PT > max_w_emu:
        pt -= 0.25
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(pt)
    return pt


# =========================================================== SLIDE 1 - cover
s1 = S[0]

set_para(shp(s1, "Text 5").text_frame.paragraphs[0],
         "Final scope of work issued 30.07.2026, superseding the Rev P06 / P07 field-sheet "
         "quantities. Quantities and locations are taken from the joint field verification "
         "of the three milling locations.")
log("S1  Basis of scope — internal working-file names removed; wording tidied.")

t7 = shp(s1, "Text 7")
r = t7.text_frame.paragraphs[0].add_run()
r.text = "ADA as-built drawing Z1-Z2-Z3-MTA_SEGMENTATION  ·  ZIA local project grid (metres)"
src_run = shp(s1, "Text 5").text_frame.paragraphs[0].runs[0]
r.font.size = src_run.font.size
r.font.name = src_run.font.name or "Arial"
r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
log("S1  'Base layer' row was labelled but blank — value added.")

set_para(shp(s1, "Text 9").text_frame.paragraphs[0],
         "Mohammed Faheem — AGL Team Leader, ADB SAFEGATE  ·  Issued 28.07.2026  ·  "
         "Rev P08 issued 30.07.2026")
log("S1  Preparer given in full, consistent with the sheet footers.")

ROW_DY = 5413248 - 4937760
for src, text, nm in ((shp(s1, "Text 8"), "Approved by", "Text 8b"),
                      (shp(s1, "Text 9"), "Ragesh Menon — AGL Manager, ADB SAFEGATE", "Text 9b")):
    new = copy.deepcopy(src._element)
    src._element.addnext(new)
    sh = _Shape(new, src._parent)
    sh.name = nm
    sh.top = Emu(src.top + ROW_DY)
    set_para(sh.text_frame.paragraphs[0], text)
log("S1  'Approved by' row added — the sheet footers name an approver, the cover did not.")

set_para(shp(s1, "Text 10").text_frame.paragraphs[0],
         "REV P08 IS THE FINAL SCOPE OF WORK.  All secondary routes at the three locations are "
         "laid by SAW CUT; no duct route remains in the AGL scope. New side-entry shallow bases: "
         "Location 1 — 11 No. 8\" (9 No. bases cored out: 3 No. 12\" + 6 No. 8\"; 2 No. 12\" bases "
         "removed earlier); Location 2 — 1 No. 12\" at TCCECH-03/008; Location 3 — 1 No. 12\" at "
         "TCCECH-03/003, with the remaining 3 No. TCC taking new secondary cable into their "
         "existing 12\" bases. A further 3 No. fittings at Location 1 (SBC102-02/027, TCCECH-03/035, TCCECH-03/018) take new "
         "secondary cable only — no coring — and are protected by a dummy plate during "
         "milling. The affected-fitting list moves from 30 No. to 19 No. A saw cut and "
         "side-entry base detail sheet is included and is held pending issue of the detail "
         "drawing. The technical queries on the existing 4 x 19 mm secondary duct are closed by "
         "agreement between the AGL team, the civil team and ADA AGL: saw cut is adopted as an "
         "interim arrangement, with permanent duct provision to follow under the South "
         "Rehabilitation works. As-built duct geometry, sheet registration and the stated limits "
         "of Rev P06 are carried forward unchanged.")
log("S1  Cover summary restated so its quantities read the same as sheets 6 and 7.")

add_logo(s1, 914400 + 13258800, 838200, 700000, reverse=True)

# ==================================================== SLIDE 2 - scope & sequence
s2 = S[1]

sub_all(shp(s2, "TextBox 4"), [
    ("2 No. 12\" bases removed prior", "2 No. 12\" bases removed earlier"),
    ("No asphalt to be laid over an unwitnessed duct.",
     "Asphalt laying to follow the witnessed mandrel test."),
    ("At Rev P08 EVERY secondary route at all three locations is saw cut — no duct route "
     "remains in the AGL scope.",
     "At Rev P08 every secondary route at the three locations is saw cut, and no duct route "
     "remains in the AGL scope."),
    ("6.  Testing and commissioning of all affected circuits.", "6.  Testing and commissioning at circuit level — covering the affected fittings together with the fittings on the same circuits that fall outside these works, since the series circuits are proved end to end before handover."),
    ("4.  Laying of new secondary cable through the saw cut. No joints — full manhole-to-light "
     "replacement.",
     "4.  Laying of new secondary cable through the saw cut — 1 No. 2-core 4 sq.mm cable to each "
     "fitting, SBC and TCC alike (7 No. SBC, 12 No. TCC). No joints — full manhole-to-light "
     "replacement."),
    ("R3.  Dummy plates installed on open bases for protection during milling and civil works.",
     "R3.  Dummy plates installed on the open bases, and on the 3 No. Location 1 fittings that "
     "take secondary cable only (SBC102-02/027, TCCECH-03/035, TCCECH-03/018), for protection during milling and civil "
     "works."),
    ("Balance 3 No. TCC at Location 3: existing 12\" shallow base maintained, new side-entry "
     "cable only.",
     "Balance 3 No. TCC at Location 3 and 3 No. at Location 1: existing shallow base maintained, "
     "new side-entry cable only, no coring."),
])
sub_all(shp(s2, "TextBox 4"), [
    ("7.  Final functionality check, then handover to Operations for final inspection and "
     "return of the area to operational service.",
     "8.  Final functionality check, then handover to Operations for final inspection and "
     "return of the area to operational service."),
])


def insert_after(shape, prefix, text):
    """Add a paragraph directly after the one whose text starts with `prefix`."""
    tf = shape.text_frame
    for n, para in enumerate(tf.paragraphs):
        if para.text.startswith(prefix):
            return clone_para(shape, n, text)
    raise AssertionError("paragraph not found: " + prefix)


seq = shp(s2, "TextBox 4")
insert_after(seq, "R3.",
             "R4.  Direction signboards leading to E4 and E6 masked with matt black vinyl "
             "sticker before works commence, and unmasked under Phase 3 before the functionality "
             "check.")
insert_after(seq, "6.  Testing and commissioning at circuit level",
             "7.  Remove the matt black vinyl masking from the direction signboards leading to "
             "E4 and E6 (applied under Phase 1, R4).")
sub_all(shp(s2, "TextBox 7"), [("(Phase 3, Item 7)", "(Phase 3, Item 8)")])
log("S2  Signboard masking added — Phase 1 R4 masks the direction signboards leading to E4 and "
    "E6, Phase 3 item 7 unmasks them before the functionality check (which becomes item 8, with "
    "hold point H4 re-referenced).")
log("S2  Phase 2 / Phase 3 wording softened without changing the requirement.")

hp = shp(s2, "TextBox 7")
sub_all(hp, [("H2 ", "H1 "), ("H3 ", "H2 "), ("H4 ", "H3 "), ("H5 ", "H4 "), ("H6 ", "H5 ")])
sub_all(hp, [("No saw cutting before the saw cut and side-entry base detail drawing is issued "
              "and accepted (Rev P08).",
              "Saw cutting to start once the saw cut and side-entry base detail drawing is "
              "issued and accepted (Rev P08).")])
log("S2  Hold points renumbered H1–H5 — the list started at H2 with no H1.")

set_para(shp(s2, "TextBox 9").text_frame.paragraphs[0],
         "TECHNICAL QUERY — CLOSED BY AGREEMENT (AGL TEAM · CIVIL TEAM · ADA AGL)")
set_lines(shp(s2, "TextBox 10"), [
    "Subject:  existing 4 x 19 mm secondary duct — cover, reuse and continuity",
    "Field condition:  the existing 4 x 19 mm secondary duct bank was exposed at 50 mm milling "
    "depth.",
    "Q1.  Whether new duct may be re-laid on the existing duct alignment.   CLOSED — no duct is "
    "re-laid on the existing alignment under this scope; the new secondary cable is laid by saw "
    "cut.",
    "Q2.  Suitability of the existing duct for reuse.   CLOSED — the existing 4 x 19 mm "
    "cross-section does not accommodate the new secondary cable, and the existing cables are "
    "obsolete and are to be withdrawn. Reuse is not adopted; the duct bank is left in place.",
    "Q3.  Transition into the duct in the non-construction area.   CLOSED — the civil scope of "
    "work covers the milling area only, so a full-stretch duct cannot be delivered under this "
    "work. Saw cut is adopted for the full affected length, which removes the transition.",
])
log("S2  Technical query updated: duct size corrected to 4 x 19 mm and Q1–Q3 recorded as closed "
    "by agreement rather than open.")

set_para(shp(s2, "TextBox 12").text_frame.paragraphs[0],
         "AGREED WAY FORWARD — SAW CUT AS AN INTERIM ARRANGEMENT")
set_lines(shp(s2, "TextBox 13"), [
    "The AGL team's original proposal was a full-stretch new secondary duct sized for the new "
    "secondary cable, giving a continuous route end to end and minimising saw cutting in the "
    "finished pavement.",
    "Agreed position:  the matter was raised with the civil team and with ADA AGL. It is agreed "
    "that the existing 4 x 19 mm duct cannot accommodate the new secondary cable, and that the "
    "civil scope of work covers the milling area only, so a full-stretch duct cannot be delivered "
    "under this work. SAW CUT is therefore adopted for every milling-affected asset — agreed by "
    "the AGL team (ADB SAFEGATE), the civil team and ADA AGL.",
    "Status:  the saw cut route is an INTERIM arrangement. Permanent duct provision is to follow "
    "under the South Rehabilitation works, where the saw cut routes are to be replaced with duct "
    "to the standard design. Q1, Q2 and Q3 are closed on that basis.",
])
log("S2  Proposal panel becomes the record of the agreed way forward — saw cut adopted as an "
    "interim arrangement, permanent duct under the South Rehabilitation works.")

sub_all(shp(s2, "TextBox 16"), [
    ("All secondary routes saw cut at all three locations; no duct route remains in the AGL scope.",
     "All secondary routes saw cut at the three locations; no duct route remains in the AGL scope."),
    ("Affected fittings reduce from 30 to 16", "Affected fittings move from 30 No. to 19 No."),
    ("LOC-01 11, LOC-02 1 (TCCECH-03/008), LOC-03 4 (TCC only — 03/003 cored, balance 3 "
     "secondary cable only).",
     "LOC-01 14 (11 with new bases, 3 secondary cable only), LOC-02 1 (TCCECH-03/008), LOC-03 4 "
     "(TCC only — 03/003 cored, balance 3 secondary cable only)."),
    ("2 No. 12\" removed prior", "2 No. 12\" removed earlier"),
    ("Coring applies at all three locations: LOC-01 mixed 12\"/8\" (9 cored); LOC-02 and LOC-03 "
     "one 12\" core each.",
     "Coring: 11 No. in total — LOC-01 9 No. (3 No. 12\" + 6 No. 8\"); LOC-02 1 No. 12\"; "
     "LOC-03 1 No. 12\"."),
])
wc = shp(s2, "TextBox 16")
clone_para(wc, len(wc.text_frame.paragraphs) - 1,
           "Technical queries Q1–Q3 on the existing 4 x 19 mm secondary duct are closed by "
           "agreement between the AGL team, the civil team and ADA AGL: saw cut is adopted as an "
           "interim arrangement, with permanent duct provision under the South Rehabilitation "
           "works.")
log("S2  'What changed' quantities restated so they reconcile with the quantities and schedule "
    "sheets, and the closure of Q1–Q3 by agreement recorded.")

set_para(shp(s2, "TextBox 17").text_frame.paragraphs[0],
         "Sequence and hold points apply to the three milling locations (sheets 1001, 1002, 1003). "
         "Field condition governs. Where this sequence and a civil programme instruction differ, "
         "the AGL Team Leader is to be informed so that the two can be aligned before work "
         "proceeds.     ·     " + FOOTER_DOC)
log("S2  Footer reworded; preparer / approver naming made consistent across the deck.")

COL_L, COL_R, COL_W = 457200, 7699248, 6958584
HDR_H, PAD_TOP, PAD_BOT, GAP = 237744, 109728, 137160, 146304
BOT = 10020300


def _scale_font(shape, factor, base):
    """Set every run to `factor` x its original size (captured in `base`)."""
    for pi, para in enumerate(shape.text_frame.paragraphs):
        for ri, run in enumerate(para.runs):
            orig = base.setdefault((pi, ri), run.font.size.pt if run.font.size else 10.0)
            run.font.size = Pt(round(orig * factor * 4) / 4)


def lay_column(panels, top, bottom):
    avail = (bottom - top) - GAP * (len(panels) - 1)
    bases = [{} for _ in panels]
    fitted = []
    for factor in (1.0, 0.96, 0.92, 0.88, 0.84, 0.80):
        fitted = []
        for (pan, hdr, body), base in zip(panels, bases):
            body.width = Emu(COL_W - 2 * (body.left - pan.left))
            if factor != 1.0 or base:
                _scale_font(body, factor, base)
            fitted.append(PAD_TOP + HDR_H + int(0.55 * HDR_H) + fit_height(body) + PAD_BOT)
        if sum(fitted) <= avail:
            break
    slack = avail - sum(fitted)
    if slack > 0:                       # spread the spare room evenly
        share = slack // len(fitted)
        fitted = [h + share for h in fitted]
        fitted[-1] += slack - share * len(fitted)
    y = top
    for (pan, hdr, body), h in zip(panels, fitted):
        pan.top, pan.height = Emu(y), Emu(h)
        hdr.top, hdr.height = Emu(y + PAD_TOP), Emu(HDR_H)
        body.top = Emu(y + PAD_TOP + HDR_H + int(0.55 * HDR_H))
        body.height = Emu(h - (PAD_TOP + HDR_H + int(0.55 * HDR_H) + PAD_BOT))
        y += h + GAP


lay_column([(shp(s2, "P08 PANEL · SCOPE OF WORK — SEQUENCE OF "), shp(s2, "TextBox 3"), shp(s2, "TextBox 4")),
            (shp(s2, "P08 PANEL · HOLD POINTS & WITNESS REQUIR"), shp(s2, "TextBox 6"), shp(s2, "TextBox 7"))],
           914400, BOT)
lay_column([(shp(s2, "P08 PANEL · TECHNICAL QUERY — CONFIRMATI"), shp(s2, "TextBox 9"), shp(s2, "TextBox 10")),
            (shp(s2, "P08 PANEL · AGL TEAM PROPOSAL"), shp(s2, "TextBox 12"), shp(s2, "TextBox 13")),
            (shp(s2, "P08 PANEL · REV P08 — WHAT CHANGED"), shp(s2, "TextBox 15"), shp(s2, "TextBox 16"))],
           914400, BOT)
place(shp(s2, "TextBox 17"), top=10149840)
place(shp(s2, "TextBox 1"), width=COL_R + COL_W - COL_L - 1700000)
add_logo(s2, COL_R + COL_W, 265430, 434340)
log("S2  Panels re-fitted to their text; the two columns now finish on the same line.")

# ============================================ SLIDES 3-5 - location sheets
TXT_X, TXT_W = 10241280, 4434840
PANEL_X, PANEL_W = 10104120, 4709160
SHEET_TOP, FRAME = 777240, 9555480

NOTES = [
    "1.  Positions plotted are ZIA LOCAL PROJECT GRID metres. The EPSG:32640 reference carried "
    "on Rev P06 is withdrawn — no surveyed local-to-UTM transform is held, so no UTM coordinate "
    "is to be taken from this sheet.",
    "2.  AGL works area (red) is the Rev P05 field-condition milling extent, unchanged at Rev "
    "P08 — it is the civil milling limit and is not derived from the affected-asset positions.",
    "3.  Duct / ductbank / conduit / saw cut lines are as-built geometry lifted from ADA drawing "
    "Z1-Z2-Z3-MTA_SEGMENTATION.dxf. The source layer name carries the duct size and way-count.",
    "4.  Ducts are drawn as the separate segments they are in the source. Continuity is not "
    "proven — a continuous pull length is not to be taken from this sheet.",
    "5.  Isolate, lock out and prove dead all listed circuits at the CCR under permit before "
    "cutting, coring or excavation.",
    "6.  HH / MH / transformer-pit / RRM symbols sit approximately 2.1–2.3 m off the source "
    "insertion point; set out civil features from survey.",
    "7.  Centreline and stop bar are taken from the as-built TCC / SBC fittings. Taxiway edge is "
    "indicative at 23 m — to be confirmed on site.",
    "8.  Every asset in frame carries its own as-built symbol, legended right. Type is taken from "
    "the source layer, not from the label.",
    "9.  AGL assets within the closed area pass into the civil team's care for the duration of "
    "the closure. Any damage to an AGL asset inside the closed area is to be made good under the "
    "civil scope, and is to be reported to the AGL Team Leader on discovery so that it is "
    "recorded and rectified before handover.",
]

SCOPE_TEXT = {
3: [
    "FINAL SCOPE (REV P08) — 14 No. FITTINGS AFFECTED, ALL ROUTES SAW CUT",
    "CORE OUT 3 No. 12\" + 6 No. 8\" (9 No. cored); 2 No. 12\" bases removed earlier — INSTALL "
    "11 No. 8\" SIDE-ENTRY SHALLOW BASE + NEW CABLE.",
    "   • SBC (6 No.): SBC102-02/024, 01/027, 02/025, 01/028, 02/026, 01/029",
    "   • TCC (5 No.): TCCECH-04/034, 03/036, 04/035, 03/037, 04/036",
    "ROUTE: SAW CUT throughout — the Rev P06/P07 via-duct routes become saw cut secondary cable "
    "laying at Rev P08. Saw cut detail per the SAW CUT & SIDE-ENTRY SHALLOW BASE DETAIL sheet; "
    "detail drawing to be issued.",
    "SECONDARY CABLE ONLY — 3 No.: SBC102-02/027, TCCECH-03/035, TCCECH-03/018. New saw cut, NO CORING; existing "
    "base kept and protected by a dummy plate during milling.",
    "FIELD VERIFIED NOT AFFECTED: SBC102-01/026 (no works)",
    "TCC103 FITTINGS NOT IN THE FIELD SCOPE — SHOWN AS EXISTING ONLY",
    "ISOLATE CIRCUITS: SBC102.01/.02, TCCECH.03/.04",
    "SITE RECORD — 23.07.2026: 9 No. bases cored out, 2 No. removed. 30.07.2026: preparing "
    "installation, all 11 No. new bases 8\".",
    "Rev P08 supersedes the Rev P06/P07 quantities. Superseded = no works under this revision; "
    "it is not a field verification that the asset is unaffected.",
    "RED AGL WORKS AREA = the Rev P05 field-condition milling extent, unchanged at Rev P08 — "
    "the civil milling limit.",
],
4: [
    "FINAL SCOPE (REV P08) — 1 No. FITTING AFFECTED, SAW CUT",
    "CORE OUT 12\" EXISTING SHALLOW BASE + INSTALL NEW 12\" SIDE-ENTRY SHALLOW BASE + NEW CABLE: "
    "1 No.",
    "   • TCCECH-03/008 — 12\" coring of the existing shallow base; new 12\" side-entry base "
    "installed; secondary cable via SAW CUT. Saw cut detail per the SAW CUT & SIDE-ENTRY SHALLOW "
    "BASE DETAIL sheet; detail drawing to be issued.",
    "DUCTS: the remaining ducts in frame are not affected during milling — no duct works at "
    "Location 2.",
    "NOT IN THE REV P08 SCOPE — no works: TCCECH-03/007, 04/007, 04/008, 03/009",
    "   Rev P06/P07 carried all 5 No. as new-secondary-cable-only via duct; the final scope names "
    "03/008 only.",
    "RRM.555 (0.41 m from the cut): REMOVE / PROTECT BEFORE SAW CUT, RE-FIX AFTER PAVING",
    "ISOLATE CIRCUITS: TCCECH.03, TCCECH.04",
    "Rev P08 supersedes the Rev P06/P07 quantities. Superseded = no works under this revision; "
    "it is not a field verification that the asset is unaffected.",
    "RED AGL WORKS AREA = the Rev P05 field-condition milling extent, unchanged at Rev P08 — "
    "the civil milling limit.",
],
5: [
    "FINAL SCOPE (REV P08) — 4 No. TCC, SAW CUT + NEW CABLE; CORE OUT 12\" AT 1 No. ONLY",
    "   • TCCECH-03/003 — CORE OUT 12\" + INSTALL NEW 12\" SIDE-ENTRY SHALLOW BASE (no side-entry "
    "provision in the existing base).",
    "   • TCCECH-03/002, 04/002, 04/003 — SECONDARY CABLE ONLY; the existing 12\" shallow base "
    "accommodates side entry — NO CORING.",
    "ROUTE: SAW CUT — these 4 No. were carried as via-duct on Rev P06/P07. Saw cut detail per the "
    "SAW CUT & SIDE-ENTRY SHALLOW BASE DETAIL sheet; detail drawing to be issued.",
    "NOT IN THE REV P08 SCOPE — no works, 8 No. SBC:",
    "   • SBC102-01/038, 01/039, 02/035, 02/036 (EP7 stop bar — was saw cut, base protect)",
    "   • SBC102-01/040, 01/041, 02/037, 02/038 (was new cable only, via duct)",
    "RRM.557 (0.07 m — treat as within the works) and RRM.670: REMOVE / PROTECT / RE-FIX",
    "ISOLATE CIRCUITS: TCCECH.03/.04 — SBC102.01/.02 carries no AGL works here at Rev P08; "
    "confirm isolation against the saw cut alignment before works.",
    "Rev P08 supersedes the Rev P06/P07 quantities. Superseded = no works under this revision; "
    "it is not a field verification that the asset is unaffected.",
    "RED AGL WORKS AREA = the Rev P05 field-condition milling extent, unchanged at Rev P08 — "
    "the civil milling limit.",
],
}

SCOPE_SHAPE = {3: "Text 313", 4: "Text 110", 5: "Text 305"}
NOTES_SHAPE = {3: "Text 316", 4: "Text 113", 5: "Text 308"}
SCOPE_HDR = {3: "Text 312", 4: "Text 109", 5: "Text 304"}
NOTES_HDR = {3: "Text 315", 4: "Text 112", 5: "Text 307"}
LEG_HDR = {3: "Text 318", 4: "Text 115", 5: "Text 310"}
SCOPE_PANEL = {3: "Shape 311", 4: "Shape 108", 5: "Shape 303"}
NOTES_PANEL = {3: "Shape 314", 4: "Shape 111", 5: "Shape 306"}
LEG_PANEL = {3: "Shape 317", 4: "Shape 114", 5: "Shape 309"}

# --- scope + notes text -----------------------------------------------------
for i in (3, 4, 5):
    sl = S[i - 1]
    body = shp(sl, SCOPE_SHAPE[i])
    tf = body.text_frame
    want = SCOPE_TEXT[i]
    while len(tf.paragraphs) > len(want):
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)
    while len(tf.paragraphs) < len(want):
        clone_para(body, len(tf.paragraphs) - 1, "x")
    for p, t in zip(tf.paragraphs, want):
        set_para(p, t)
        p.space_after = Pt(0)
        p.space_before = Pt(0)
    notes_shape = shp(sl, NOTES_SHAPE[i])
    ntf = notes_shape.text_frame
    while len(ntf.paragraphs) > len(NOTES):
        p = ntf.paragraphs[-1]._p
        p.getparent().remove(p)
    while len(ntf.paragraphs) < len(NOTES):
        clone_para(notes_shape, len(ntf.paragraphs) - 1, "x")
    for p, t in zip(ntf.paragraphs, NOTES):
        set_para(p, t)
log("S3–S5  Scope panels and general notes rewritten: EPSG note reworded, the stale 'see sheet 6' "
    "cross-reference removed, internal shorthand replaced with plain wording.")

# --- legend: keep only the symbols that appear on the sheet -----------------
LEGEND_ROWS = {
    3: [("ring", "CORE OUT 12\" — NEW 8\" SIDE-ENTRY BASE (SAW CUT)"),
        ("ring8", "CORE OUT 8\" — NEW 8\" SIDE-ENTRY BASE (SAW CUT)"),
        ("green", "SECONDARY CABLE ONLY — NO CORING (DUMMY PLATE)"),
        ("line_sc", "SECONDARY CABLE — SAW CUT ROUTE (REV P08)"),
        ("line_new", "NEW SAW CUT — SECONDARY CABLE (REV P08)"),
        ("line_area", "AGL WORKS AREA (GOVERNING)")],
    4: [("ring", "CORE OUT 12\" — NEW 12\" SIDE-ENTRY BASE (SAW CUT)"),
        ("rrm", "RRM — REMOVE / PROTECT / RE-FIX"),
        ("line_sc", "SECONDARY CABLE — SAW CUT ROUTE (REV P08)"),
        ("line_new", "NEW SAW CUT — SECONDARY CABLE (REV P08)"),
        ("line_area", "AGL WORKS AREA (GOVERNING)")],
    5: [("ring", "CORE OUT 12\" — NEW 12\" SIDE-ENTRY BASE (SAW CUT)"),
        ("green", "SECONDARY CABLE ONLY — EXISTING BASE (SAW CUT)"),
        ("rrm", "RRM — REMOVE / PROTECT / RE-FIX"),
        ("line_sc", "SECONDARY CABLE — SAW CUT ROUTE (REV P08)"),
        ("line_new", "NEW SAW CUT — SECONDARY CABLE (REV P08)"),
        ("line_area", "AGL WORKS AREA (GOVERNING)")],
}
KIND_OF_TEXT = [
    ("CORE OUT 12\"", "ring"),
    ("CORE OUT 8\"", "ring8"),
    ("ONLY SECONDARY CABLE AFFECTED", "green"),
    ("RRM — REMOVE", "rrm"),
    ("DUMMY PLATE", "dummy"),
    ("SECONDARY CABLE — SAW CUT ROUTE", "line_sc"),
    ("NEW SAW CUT — SECONDARY CABLE", "line_new"),
    ("AGL WORKS AREA (GOVERNING)", "line_area"),
]


def legend_kind(txt):
    for needle, kind in KIND_OF_TEXT:
        if needle in txt:
            return kind
    return None


LEGEND_STATE = {}
for i in (3, 4, 5):
    sl = S[i - 1]
    keep_kinds = [k for k, _ in LEGEND_ROWS[i]]
    rows = {}
    for s in list(sl.shapes):
        if not s.has_text_frame or s.left is None or s.top is None:
            continue
        if s.left < 10400000 or s.top < 7000000:
            continue
        k = legend_kind(s.text_frame.text)
        if k:
            rows[k] = {"text": s, "glyphs": []}
    for s in list(sl.shapes):
        if s.left is None or s.top is None:
            continue
        if not (10241280 <= s.left <= 10620000 and s.top > 7000000):
            continue
        if s.has_text_frame and legend_kind(s.text_frame.text):
            continue
        cy = s.top + (s.height or 0) / 2
        best, bd = None, None
        for k, row in rows.items():
            t = row["text"]
            d = abs(cy - (t.top + t.height / 2))
            if bd is None or d < bd:
                best, bd = k, d
        if best is not None and bd < 200000:
            rows[best]["glyphs"].append(s)
    for k in list(rows):
        if k not in keep_kinds:
            del_shape(rows[k]["text"])
            for g in rows[k]["glyphs"]:
                del_shape(g)
            del rows[k]
            log(f"S{i}  Legend entry '{k}' removed — that symbol is not used on this sheet.")
    # remember each row's glyph offsets so the block can be re-seated later
    for k, row in rows.items():
        row["dy"] = [g.top - row["text"].top for g in row["glyphs"]]
        set_para(row["text"].text_frame.paragraphs[0], dict(LEGEND_ROWS[i])[k])
    asset_shapes, sub_heads = [], []
    for s in list(sl.shapes):
        if s.left is None or s.top is None or "LOGO" in s.name:
            continue
        txt = s.text_frame.text.strip() if s.has_text_frame else ""
        if txt in ("WORKS ACTIONS & LINEWORK", "AS-BUILT AGL ASSETS"):
            sub_heads.append(s)
        elif s.left >= 12400000 and s.top > 7000000:
            asset_shapes.append(s)
    # cluster the asset column into rows by their plotted y
    clusters = []
    for s in sorted(asset_shapes, key=lambda x: x.top):
        if clusters and s.top - clusters[-1][0] < 150000:
            clusters[-1][1].append(s)
        else:
            clusters.append([s.top, [s]])
    LEGEND_STATE[i] = {"rows": rows, "subs": sub_heads,
                       "assets": [(top, [(x, x.top - top) for x in grp])
                                  for top, grp in clusters]}

# --- sheet 1001: drop the green markers on assets that carry no works -------
log("S3  The three green markers at SBC102-02/027, TCCECH-03/035, TCCECH-03/018 are kept — those fittings take new "
    "secondary cable by saw cut with no coring, and are protected by a dummy plate during "
    "milling.")

# --- sheet 1003: plot the RRM remove/protect rings the scope calls for ------
s5 = S[4]
RRM_DX, RRM_DY = -118872, 219456          # label origin -> ring centre, as on sheet 1002
for label, tag in (("RRM.557", "557"), ("RRM.670", "670")):
    lab = next((s for s in s5.shapes
                if s.has_text_frame and s.text_frame.text.strip() == label), None)
    if lab is None:
        continue
    d = 210312
    ring = s5.shapes.add_shape(MSO_SHAPE.OVAL, Emu(int(lab.left + RRM_DX - d / 2)),
                               Emu(int(lab.top + RRM_DY - d / 2)), Emu(d), Emu(d))
    ring.name = "P08 RRM RING " + tag
    ring.fill.background()
    ring.line.color.rgb = RGBColor(0xBB, 0x00, 0xBB)
    ring.line.width = Pt(1.5)
    ring.shadow.inherit = False
log("S5  RRM remove/protect rings plotted at RRM.557 and RRM.670 — both are in the scope and "
    "legended, but had no symbol on the sheet.")

# --- works area: restore the Rev P05 governing extent on sheets 1002/1003 ---
def add_quad(slide, pts, name):
    b = slide.shapes.build_freeform(Emu(int(pts[0][0])), Emu(int(pts[0][1])))
    b.add_line_segments([(Emu(int(x)), Emu(int(y))) for x, y in pts[1:]], close=True)
    sh = b.convert_to_shape()
    sh.name = name
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor(0xCC, 0x00, 0x00)
    srgb = sh.fill._xPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
    alpha = srgb.makeelement(qn("a:alpha"), {"val": "34000"})
    srgb.append(alpha)
    sh.line.color.rgb = RGBColor(0xCC, 0x00, 0x00)
    sh.line.width = Pt(1.5)
    sh.shadow.inherit = False
    return sh


AREAS = {
    4: dict(pts=[(3017137, 4158882), (7397996, 7676881), (7177721, 7951172), (2796862, 4433261)],
            old="AGL WORKS AREA FILL 1"),
    5: dict(pts=[(3087402, 6322353), (5933660, 7044156), (6504497, 4793428), (3658139, 4071626)],
            old="AGL WORKS AREA FILL 2"),
}
for i, spec in AREAS.items():
    sl = S[i - 1]
    old = shp(sl, spec["old"])
    z = list(sl.shapes._spTree).index(old._element)
    del_shape(old)
    new = add_quad(sl, spec["pts"], "AGL WORKS AREA FILL")
    sl.shapes._spTree.remove(new._element)
    sl.shapes._spTree.insert(z, new._element)
log("S4/S5  Red AGL works area redrawn to the Rev P05 governing milling extent (257.1 m² and "
    "676.7 m²) that every sheet note describes — the plotted patch had drifted to 99.7 m² and "
    "292.9 m².")

for i, m2, nm in ((3, "425.3", "Text 163"), (5, "676.7", "Text 150")):
    set_para(shp(S[i - 1], nm).text_frame.paragraphs[0],
             f"AGL WORKS AREA — FIELD CONDITION ({m2} m² TOTAL)")
src_cap = shp(S[4], "Text 150")
new_cap = copy.deepcopy(src_cap._element)
S[3].shapes._spTree.append(new_cap)
cap4 = _Shape(new_cap, S[3].shapes)
cap4.name = "Text AREA CAPTION"
cap4.left, cap4.top = Emu(2130000), Emu(3900000)
set_para(cap4.text_frame.paragraphs[0], "AGL WORKS AREA — FIELD CONDITION (257.1 m² TOTAL)")
for i, nm in ((3, "Text 163"), (4, "Text AREA CAPTION"), (5, "Text 150")):
    sl = S[i - 1]
    cap = shp(sl, nm)
    cap.fill.solid()
    cap.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    cap.line.fill.background()
    cap.width = Emu(int(textfit.text_w_pt(cap.text_frame.text, 9.0, True) * textfit.EMU_PT) + 137160)
    sl.shapes._spTree.remove(cap._element)        # draw above the linework
    sl.shapes._spTree.append(cap._element)
log("S3–S5  Works-area captions given an opaque backing so they read clear of the linework.")
log("S4  Works-area caption added — sheets 1001 and 1003 carried one, sheet 1002 did not.")

# --- titles, frames, right-hand column, footer, logo ------------------------
TITLES = {
    3: "SHOP DRAWING — ELECTRICAL / AGL SCOPE (FIELD-GOVERNED)   ·   TAXIWAY E (E4–E6) — "
       "LOCATION 1   ·   DWG NO AUH-SK-AGL-TWYE-001-1001 REV P08",
    4: "SHOP DRAWING — ELECTRICAL / AGL SCOPE (FIELD-GOVERNED)   ·   TAXIWAY E (E4–E6) — "
       "LOCATION 2   ·   DWG NO AUH-SK-AGL-TWYE-001-1002 REV P08",
    5: "SHOP DRAWING — ELECTRICAL / AGL SCOPE (FIELD-GOVERNED)   ·   TAXIWAY E (E4–E6) — "
       "LOCATION 3   ·   DWG NO AUH-SK-AGL-TWYE-001-1003 REV P08",
}
LOGO_W = int(434340 * 788 / 272)
for i in (3, 4, 5):
    sl = S[i - 1]
    t = shp(sl, "Text 0")
    place(t, left=320040, top=196000, width=13100000, height=411480)
    set_para(t.text_frame.paragraphs[0], TITLES[i])
    t.text_frame.word_wrap = False
    fit_font(t, 14813280 - LOGO_W - 320040 - 274320, 14.0, 9.5)
    add_logo(sl, 14813280, 175260, 434340)

    fr = shp(sl, "Shape 1")
    place(fr, left=320040, top=SHEET_TOP, width=FRAME, height=FRAME)
    if i == 3:
        fr.fill.solid()
        fr.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        fr.line.color.rgb = RGBColor(0x5F, 0x63, 0x68)
        fr.line.width = Pt(0.75)

    banner = next(s for s in sl.shapes
                  if s.left == 10104120 and s.top == 777240 and s.width == PANEL_W)
    GAP2 = 137160
    y = banner.top + banner.height + GAP2
    for pnm, hnm, bnm in ((SCOPE_PANEL[i], SCOPE_HDR[i], SCOPE_SHAPE[i]),
                          (NOTES_PANEL[i], NOTES_HDR[i], NOTES_SHAPE[i])):
        panel, hdr, body = shp(sl, pnm), shp(sl, hnm), shp(sl, bnm)
        place(body, left=TXT_X, width=TXT_W)
        need = fit_height(body)
        h = 137160 + 219456 + 128016 + need + 137160
        place(panel, top=y, height=h)
        place(hdr, top=y + 137160)
        place(body, top=y + 137160 + 219456 + 128016, height=need)
        y += h + GAP2

    st = LEGEND_STATE[i]
    # each legend row is as tall as its own wrapped label needs
    LEG_MIN, LEG_LINE = 274320, 152400
    row_h = []
    for k, label in LEGEND_ROWS[i]:
        w = st["rows"][k]["text"].width
        lines = textfit.wrap_lines(label, (w - 91440) / textfit.EMU_PT, 8.0)
        row_h.append(max(LEG_MIN, lines * LEG_LINE + 91440))
    LEG_CHROME = 137160 + 219456 + 128016 + 137160
    rows_h = max(sum(row_h), 274320 * len(st["assets"]))
    avail = 10058400 - y - LEG_CHROME                  # keep clear of the sheet footer
    if rows_h > avail:                                  # scale back rather than overrun
        k = avail / rows_h
        row_h = [max(228600, int(h * k)) for h in row_h]
        rows_h = max(sum(row_h), 274320 * len(st["assets"]))
    leg_h = rows_h + LEG_CHROME
    place(shp(sl, LEG_PANEL[i]), top=y, height=leg_h)
    place(shp(sl, LEG_HDR[i]), top=y + 137160)
    sub_y = y + 137160 + 256848
    row0 = sub_y + 292608
    for s_ in st["subs"]:
        s_.top = Emu(sub_y)
    # works-action / linework column
    yy = row0
    for n_i, (k, _) in enumerate(LEGEND_ROWS[i]):
        row = st["rows"][k]
        for g, dy in zip(row["glyphs"], row["dy"]):
            g.top = Emu(yy + dy)
        row["text"].top = Emu(yy)
        row["text"].height = Emu(row_h[n_i])
        yy += row_h[n_i]
    # as-built asset column on the same row pitch
    for n_i, (_, grp) in enumerate(st["assets"]):
        new_top = row0 + n_i * 274320
        for s_, dy in grp:
            s_.top = Emu(new_top + dy)
    add_footer(sl, "Field condition governs.  Quantities are per the Rev P08 final scope of work "
                   "(30.07.2026) and are for check by the Engineer.     ·     " + FOOTER_DOC,
               320040, 14493240, 10390000, size=7.0)
log("S3–S5  Titles fitted on one line clear of the logo, drawing frames aligned to their content, "
    "right-hand panels fitted to their text and legends re-seated on an even pitch, sheet footers "
    "added.")

# ================================================= SLIDE 6 - consolidated scope
s6 = S[5]
tbl_shape = [s for s in s6.shapes if s.has_table][0]
tbl = tbl_shape.table
ACTION = {
    "CORE OUT 8\" + NEW SIDE-ENTRY SHALLOW BASE + NEW CABLE":
        "CORE OUT 8\" → NEW 8\" SIDE-ENTRY BASE + NEW CABLE",
    "CORE OUT 12\" / REMOVED PRIOR — NEW 8\" BASE + NEW CABLE":
        "CORE OUT 12\" / REMOVED EARLIER → NEW 8\" BASE + CABLE",
    "CORE OUT EXISTING BASE + NEW 12\" SIDE-ENTRY SHALLOW BASE + NEW CABLE":
        "CORE OUT EXISTING → NEW 12\" SIDE-ENTRY BASE + CABLE",
    "SECONDARY CABLE ONLY — EXISTING 12\" BASE MAINTAINED":
        "SECONDARY CABLE ONLY — EXISTING 12\" BASE KEPT",
    "CORE OUT 12\" + NEW 12\" SIDE-ENTRY BASE + NEW CABLE":
        "CORE OUT 12\" → NEW 12\" SIDE-ENTRY BASE + CABLE",
    "NOT IN REV P08 SCOPE": "NOT IN THE REV P08 SCOPE",
    "RRM REMOVE / PROTECT": "RRM — REMOVE / PROTECT / RE-FIX",
}
REMARK = {
    "8\" core out; new side-entry shallow base; full manhole-to-light cable replacement via saw cut":
        "8\" core out; new side-entry base; full manhole-to-light cable replacement via saw cut",
    "12\" base cored (3 No.) or removed prior (2 No.); new 8\" base; cable replaced via saw cut":
        "12\" base cored (3 No.) or removed earlier (2 No.); new 8\" base; cable via saw cut",
    "In the Rev P06/P07 scope; not named in the final scope of work 30.07.2026. No works. Not a "
    "field verification of 'not affected'":
        "Carried in Rev P06/P07; not named in the final scope of work 30.07.2026 — no works at "
        "this revision",
    "Existing shallow base cored out; new 12\" side-entry shallow base; full manhole-to-light "
    "cable replacement via saw cut":
        "Existing shallow base cored out; new 12\" side-entry base; full cable replacement via "
        "saw cut",
    "Existing 12\" base accepts new side-entry cable — no coring; cable via saw cut":
        "Existing 12\" base accepts the new side-entry cable — no coring; cable via saw cut",
    "No side-entry provision in existing base; 12\" core; new 12\" base; cable via saw cut":
        "No side-entry provision in the existing base; 12\" core; new 12\" base; cable via saw cut",
    "Remove / protect before saw cut, re-fix after paving":
        "Remove and protect before saw cut; re-fix after paving",
}
CABLE_ONLY_L1 = ("SBC102-02/027", "TCCECH-03/035", "TCCECH-03/018")
for ri in range(1, len(tbl.rows)):
    if tbl.cell(ri, 1).text in CABLE_ONLY_L1:
        set_para(tbl.cell(ri, 2).text_frame.paragraphs[0],
                 "SECONDARY CABLE ONLY — NEW SAW CUT, NO CORING")
        set_para(tbl.cell(ri, 3).text_frame.paragraphs[0], "Sawcut")
        set_para(tbl.cell(ri, 4).text_frame.paragraphs[0],
                 "Existing base kept, no coring; protected by dummy plate during milling; "
                 "new cable via saw cut")
        continue
    for ci, mapping in ((2, ACTION), (4, REMARK)):
        cell = tbl.cell(ri, ci)
        if cell.text_frame.text in mapping:
            set_para(cell.text_frame.paragraphs[0], mapping[cell.text_frame.text])
for ci, w in enumerate((960120, 1737360, 4023360, 914400, 6537960)):
    tbl.columns[ci].width = Emu(w)
for r in tbl.rows:
    r.height = Emu(240030)
place(tbl_shape, top=960120, height=240030 * len(tbl.rows))
log("S6  Action and remark wording shortened and softened; columns re-proportioned so every row "
    "is a single line — the table used to run underneath the totals line.")

set_para(shp(s6, "Text 1").text_frame.paragraphs[0],
         "Totals (Rev P08 final scope): 19 No. fittings — LOC-01 14 No. (9 No. cored: 3 No. 12\" + "
         "6 No. 8\"; 2 No. 12\" removed earlier; 11 No. new 8\" bases; balance 3 No. secondary "
         "cable only, no coring); LOC-02 1 No. (12\" core, new 12\" base); LOC-03 4 No. (12\" "
         "core at 03/003 only, balance 3 No. secondary cable only) — all new cable via SAW CUT."
         "   ·   Coring 11 No. (6 @ 8\" · 5 @ 12\")   ·   New side-entry bases 13 No.   ·   "
         "12 No. not in the Rev P08 scope   ·   3 No. RRM remove / protect / re-fix   ·   1 No. "
         "field verified not affected.")
place(shp(s6, "Text 1"), top=960120 + 240030 * len(tbl.rows) + 182880, width=14173200)
place(shp(s6, "Text 0"), width=14173200 - 1700000)
add_logo(s6, 457200 + 14173200, 265430, 434340)
add_footer(s6, FOOTER_DOC, 457200, 14173200, 10240000)
log("S6  Totals line restated with the coring and new-base totals, and moved clear of the table.")

# ==================================================== SLIDE 7 - detail hold sheet
s7 = S[6]
set_lines(shp(s7, "TextBox 10"), [
    "SAW CUT DETAIL DRAWING TO BE ISSUED",
    "The final scope of work records that the saw cut detail drawing will be provided. It had not "
    "been issued at the date of this revision.",
    "Until it is issued and accepted, saw cutting is not to commence — hold point H5.",
])
log("S7  Internal 'paste the issued detail into this frame' instruction removed; hold-point "
    "reference follows the H1–H5 renumbering.")

sub_all(shp(s7, "TextBox 14"), [
    ("How the saw cut route transitions into the existing duct in the NON-CONSTRUCTION area — "
     "this is open technical query Q3 and it governs the ends of every run.",
     "How the saw cut route terminates at the edge of the milling area, where the existing "
     "4 x 19 mm duct continues into the non-construction area. Saw cut for the full affected "
     "length is agreed, so no duct transition is made; the detail is to show the termination and "
     "the seal at that point."),
    ("Treatment where the saw cut crosses the existing 4 x 110 mm secondary duct bank exposed at "
     "50 mm milling depth.",
     "Treatment where the saw cut crosses the existing 4 x 19 mm secondary duct bank exposed at "
     "50 mm milling depth, which is left in place."),
    ("Items 8 and 9 are the two that are not purely detailing. Item 8 is open technical query Q3 "
     "— the transition into the non-construction area was never answered, and every run has two "
     "ends. Item 9 is a physical clash with a duct bank the field already found at 50 mm. Neither "
     "can be closed out on site.",
     "Items 8 and 9 are the two that go beyond plain detailing. Item 8 follows from the agreed "
     "saw cut arrangement and governs the ends of every run. Item 9 is a physical clash with a "
     "duct bank already found in the field at 50 mm. Both are to be resolved in the detail rather "
     "than on site."),
])
sub_all(shp(s7, "TextBox 8"), [
    ("12.  At Location 2, the method for coring out the EXISTING shallow base and making good "
     "before the new 12\" base is set.",
     "12.  At Location 2, the method for coring out the existing shallow base and making good "
     "before the new 12\" base is set."),
])
log("S7  Wording on items 8, 9 and 12 reworded.")

sched_tbl_shape = [s for s in s7.shapes if s.has_table][0]
tt = sched_tbl_shape.table
ROWS = [
    ("LOC-01", "14 No.  (7 SBC · 7 TCC)", "9 No. — 3 @ 12\" · 6 @ 8\"",
     "11 No. 8\" new; 3 No. existing kept", "Saw cut", "1001"),
    ("LOC-02", "1 No.  (TCCECH-03/008)", "1 No. 12\" (existing base)",
     "1 No. 12\" side-entry base", "Saw cut", "1002"),
    ("LOC-03", "4 No.  (TCC only)", "1 No. 12\" (TCCECH-03/003)",
     "1 No. 12\" new; 3 No. existing kept", "Saw cut", "1003"),
    ("TOTAL", "19 No.", "11 No. — 6 @ 8\" · 5 @ 12\"",
     "13 No. new side-entry bases", "Saw cut", "—"),
]
for ri, row in enumerate(ROWS, start=1):
    for ci, val in enumerate(row):
        set_para(tt.cell(ri, ci).text_frame.paragraphs[0], val)
S7_AVAIL, S7_PAD = 6665976, 91440
need = []
for ci in range(len(tt.columns)):
    w = max(textfit.text_w_pt(tt.cell(ri, ci).text_frame.text, 9.0 if ri == 0 else 8.0, ri == 0)
            for ri in range(len(tt.rows)))
    need.append(int(w * textfit.EMU_PT) + S7_PAD + 45720)
spare = S7_AVAIL - sum(need)
S7_COLS = tuple(w + spare // len(need) for w in need)
S7_COLS = S7_COLS[:-1] + (S7_AVAIL - sum(S7_COLS[:-1]),)
for ci, w in enumerate(S7_COLS):
    tt.columns[ci].width = Emu(w)
place(sched_tbl_shape, left=603504, width=S7_AVAIL)
for r in tt.rows:
    r.height = Emu(219456)
for ri in range(len(tt.rows)):
    for ci in range(len(tt.columns)):
        c = tt.cell(ri, ci)
        c.margin_left = c.margin_right = Emu(45720)
# every cell must sit on one line, or the table outgrows its panel
S7_ROW_H = 219456
for ri in range(len(tt.rows)):
    for ci in range(len(tt.columns)):
        txt = tt.cell(ri, ci).text_frame.text
        pt = 9.0 if ri == 0 else 8.0
        inner = (S7_COLS[ci] - 91440) / textfit.EMU_PT
        if textfit.text_w_pt(txt, pt, ri == 0) > inner:
            print(f"   ! s7 table wrap risk r{ri} c{ci}: {txt!r}")
log("S7  Schedule corrected: coring 11 No. (6 @ 8\" · 5 @ 12\") and 13 No. new side-entry bases. "
    "The TOTAL row previously read '11 @ 8\" · 5 @ 12\"' and '16 side-entry shallow bases', "
    "neither of which reconciled with the per-location rows.")

set_para(shp(s7, "TextBox 15").text_frame.paragraphs[0],
         "This sheet is a HOLD. It states what the awaited detail has to resolve and the "
         "quantities that depend on it; it does not show a section, because none has been issued "
         "yet. Quantities are per the Rev P08 final scope of work (30.07.2026) and are for check "
         "by the Engineer.     ·     " + FOOTER_DOC)

# move the 'what the detail must also show' block into the right column so the
# two columns carry a similar amount and finish level
pan_hold = shp(s7, "P08 PANEL · DETAIL — HELD FOR ISSUE")
pan_sched = shp(s7, "P08 PANEL · SAW CUT & CORING SCHEDULE — ")
pan_also = shp(s7, "P08 PANEL · SIDE-ENTRY SHALLOW BASE — WH")
pan_ans = shp(s7, "P08 PANEL · WHAT THE SAW CUT DETAIL MUST")
hold_hdr, hold_body = shp(s7, "TextBox 3"), shp(s7, "TextBox 10")
hold_frame = shp(s7, "P08 HOLD · SAW CUT DETAIL FRAME")
sched_hdr = shp(s7, "TextBox 5")
sched_tbl = sched_tbl_shape
also_hdr, also_body = shp(s7, "TextBox 7"), shp(s7, "TextBox 8")
ans_hdr, ans_body = shp(s7, "TextBox 13"), shp(s7, "TextBox 14")

R_X = 7699248
dx = R_X - pan_also.left
for s in (pan_also, also_hdr, also_body):
    s.left = Emu(s.left + dx)

# right column first: it sets the common bottom line
TOP7 = 914400
y = TOP7
ans_h = 109728 + HDR_H + 128016 + fit_height(ans_body) + 137160
place(pan_ans, top=y, height=ans_h)
place(ans_hdr, top=y + 109728)
place(ans_body, top=y + 109728 + HDR_H + 128016)
y += ans_h + GAP
also_h = 109728 + HDR_H + 128016 + fit_height(also_body) + 137160
place(pan_also, top=y, height=also_h)
place(also_hdr, top=y + 109728)
place(also_body, top=y + 109728 + HDR_H + 128016)
right_bot = y + also_h

# left column: hold frame, then the schedule - levelled to the same bottom
sched_h = 109728 + HDR_H + 128016 + S7_ROW_H * len(tt.rows) + 137160
hold_h = right_bot - TOP7 - GAP - sched_h
y = TOP7
place(pan_hold, top=y, height=hold_h)
place(hold_hdr, top=y + 109728)
fr_top = y + 109728 + HDR_H + 128016
fr_h = hold_h - (109728 + HDR_H + 128016) - 137160
place(hold_frame, top=fr_top, height=fr_h)
place(hold_body, top=fr_top + int((fr_h - hold_body.height) / 2))
y += hold_h + GAP
place(pan_sched, top=y, height=sched_h)
place(sched_hdr, top=y + 109728)
place(sched_tbl, top=y + 109728 + HDR_H + 128016, height=S7_ROW_H * len(tt.rows))
log("S7  'What the detail must also show' moved into the right column and both columns levelled — "
    "the right column previously ended well above the left.")

place(shp(s7, "TextBox 1"), width=14173200 - 1700000)
add_logo(s7, 457200 + 14173200, 265430, 434340)

# =============================================== NEW SHEET - 5-day AGL installation schedule
NAVY = RGBColor(0x1E, 0x27, 0x61)
INK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x5F, 0x63, 0x68)
RULE = RGBColor(0xD3, 0xD8, 0xDE)
CELL_RULE = RGBColor(0xC8, 0xCD, 0xD3)
TINT = RGBColor(0xF4, 0xF6, 0xF8)


def _run(para, text, pt, bold=False, color=INK, italic=False):
    r = para.add_run()
    r.text = text
    r.font.size = Pt(pt)
    r.font.name = "Arial"
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r


def box(slide, x, y, w, h, fill=RGBColor(0xFF, 0xFF, 0xFF), line=RULE, name="BOX"):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)),
                                Emu(int(w)), Emu(int(h)))
    sh.name = name
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.text_frame.word_wrap = True
    return sh


def textbox(slide, x, y, w, h, lines, name="TXT"):
    """lines: list of (text, pt, bold, colour) - one paragraph each."""
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tb.name = name
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for n, (txt, pt, bold, col) in enumerate(lines):
        para = tf.paragraphs[0] if n == 0 else tf.add_paragraph()
        _run(para, txt, pt, bold, col)
        para.space_after = Pt(4)
    return tb


def cell_borders(cell, colour=CELL_RULE, w=6350):
    tc = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for old in tc.findall(qn(tag)):
            tc.remove(old)
    for tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        ln = tc.makeelement(qn(tag), {"w": str(w), "cap": "flat", "cmpd": "sng", "algn": "ctr"})
        fill = ln.makeelement(qn("a:solidFill"), {})
        clr = fill.makeelement(qn("a:srgbClr"), {"val": str(colour)})
        fill.append(clr)
        ln.append(fill)
        tc.insert(0, ln)


sched = prs.slides.add_slide(prs.slide_layouts[0])
L, W = 457200, 14173200

textbox(sched, L, 274320, W - 1700000, 502920,
        [("AGL INSTALLATION — 5-DAY WORK SCHEDULE  ·  TWY E (E4–E6)  ·  REV P08 (FINAL SCOPE)",
          20, True, NAVY)], "TITLE")
add_logo(sched, L + W, 265430, 434340)

# ---- the five-day strip ----------------------------------------------------
DAYS = [
    ("DAY 1", "SETTING-OUT & SAW CUT START", "LOCATION 1  ·  SHEET 1001"),
    ("DAY 2", "SAW CUT COMPLETE · NEW BASES", "LOCATION 1  ·  SHEET 1001"),
    ("DAY 3", "CABLE, TEST & REINSTATEMENT", "LOCATION 1  ·  SHEET 1001"),
    ("DAY 4", "FULL CYCLE — CUT TO SEALED", "LOCATION 3  ·  SHEET 1003"),
    ("DAY 5", "FULL CYCLE + T&C & HANDOVER", "LOCATION 2  ·  SHEET 1002"),
]
GAP5, BAR_H, DAY_H = 137160, 320040, 941070
BOX_W = (W - 4 * GAP5) // 5
for n, (day, head, loc) in enumerate(DAYS):
    x = L + n * (BOX_W + GAP5)
    box(sched, x, 914400, BOX_W, DAY_H, name=f"DAY BOX {n+1}")
    bar = box(sched, x, 914400, BOX_W, BAR_H, fill=NAVY, line=NAVY, name=f"DAY BAR {n+1}")
    p0 = bar.text_frame.paragraphs[0]
    _run(p0, day, 11, True, RGBColor(0xFF, 0xFF, 0xFF))
    textbox(sched, x + 137160, 914400 + BAR_H + 128016, BOX_W - 274320, 640080,
            [(head, 9.5, True, NAVY), (loc, 8.5, False, MUTED)], f"DAY TXT {n+1}")
    if textfit.text_w_pt(head, 9.5, True) * textfit.EMU_PT > BOX_W - 274320:
        print(f"   ! day box {n+1} headline wraps: {head!r}")

# ---- lower panels measured first, so the table can take up the slack -------
LOWER = (
    (L, "HOLD POINTS GOVERNING THE PROGRAMME", [
        ("H1   Curing complete — Civil.  Written confirmation of the asphalt curing period. "
         "Required before Day 1 starts.", 9.5),
        ("H5   Saw cut detail — AGL / Engineer.  Detail drawing issued and accepted. Required "
         "before any cutting on Day 1.", 9.5),
        ("H2   Setting-out — Civil survey.  Coordinates issued and field points confirmed in "
         "the field. Day 1.", 9.5),
        ("H4   Functionality check — AGL / Operations.  Witnessed before handover. Day 5.", 9.5),
        ("H3 (mandrel test) sits in Phase 2 civil works and falls outside these five days.",
         9.0),
    ], "P08 PANEL · HOLD POINTS GOVERNING"),
    (7699248, "BASIS, ASSUMPTIONS & INDICATIVE RESOURCES", [
        ("Covers Phase 3 AGL installation only. Phase 1 asset removal, coring and the masking "
         "of the direction signboards leading to E4 and E6 with matt black vinyl sticker sit "
         "before this window; Phase 2 civil "
         "attendance sits outside it. Unmasking the signboards is Day 5, ahead of the "
         "functionality check.", 9.0),
        ("The three locations are worked one front at a time — LOC-01 over Days 1–3, LOC-03 on "
         "Day 4, LOC-02 on Day 5 — so each is cut, cabled, tested and sealed before the next is "
         "opened. Days are allocated to the work each location carries: 14 No. fittings and "
         "approx. 300 m of cut at LOC-01, 4 No. and approx. 95 m at LOC-03, 1 No. and approx. "
         "24 m at LOC-02.", 9.0),
        ("Only two items span all three locations: the setting-out on Day 1, taken in one survey "
         "visit, and the circuit testing and commissioning on Day 5, which cannot complete until "
         "every location is cabled because the circuits run through all three.", 9.0),
        ("One day = one approved working shift under the AWAN / PTW. The sequence holds for a "
         "day or a night closure; shift hours to be set against the approved airside window.",
         9.0),
        ("The five days run consecutively from the release of H1 and H5. Neither is released at "
         "the date of this revision.", 9.0),
        ("Saw cut lengths scaled from sheets 1001–1003 — approx. 300 m, 24 m and 95 m, approx. "
         "420 m in total. For programme only; confirm on site.", 9.0),
        ("Sealant cure time before the area returns to operational service is per the awaited "
         "saw cut detail. If it exceeds the Day 5 shift, handover moves to the following shift.",
         9.0),
        ("Indicative resourcing per shift: 1 No. AGL supervisor, 4 No. AGL technicians, 2 No. "
         "saw cut operatives, 1 No. safety banksman; survey attendance Day 1, Operations "
         "attendance Day 5. To be confirmed against the approved AWAN.", 9.0),
    ], "P08 PANEL · BASIS AND ASSUMPTIONS"),
)

PAN_W, CHROME = 6958584, 109728 + 237744 + 128016 + 137160
LOWER_BODY_W = PAN_W - 292608
TBL_Y = 914400 + DAY_H + 274320

# ---- the schedule table ----------------------------------------------------
SCH_ROWS = [
    ("Day", "Phase 3 activity", "Location", "Quantity", "Hold point / witness",
     "Output at the end of the shift"),
    ("1", "Setting-out from the civil survey points — all three locations in one visit",
     "LOC-01 / 02 / 03", "3 No. locations", "H2 — civil survey",
     "Field points confirmed and marked at all three"),
    ("1", "Saw cutting — first part of the secondary cable route", "LOC-01 · 1001",
     "approx. 150 m", "H5 — detail accepted", "Half the LOC-01 route cut and covered"),
    ("2", "Saw cutting — balance of the route", "LOC-01 · 1001", "approx. 150 m",
     "H5 — detail accepted", "LOC-01 route cut in full"),
    ("2", "Set the new side-entry shallow bases in the cored positions", "LOC-01 · 1001",
     "11 No. 8\"", "—", "11 No. bases set and levelled"),
    ("3", "Lay new secondary cable through the saw cut — no joints, manhole to light",
     "LOC-01 · 1001", "14 No. runs", "—", "LOC-01 cabled"),
    ("3", "Termination, insulation resistance and continuity testing", "LOC-01 · 1001",
     "14 No. fittings", "—", "LOC-01 test records issued"),
    ("3", "Backfill, sealant and pavement reinstatement of the saw cut", "LOC-01 · 1001",
     "approx. 300 m", "Per the awaited detail", "LOC-01 cut sealed; cure period started"),
    ("4", "Saw cutting — full route", "LOC-03 · 1003", "approx. 95 m",
     "H5 — detail accepted", "LOC-03 route cut"),
    ("4", "Set 1 No. new base; the balance takes side entry into its existing base",
     "LOC-03 · 1003", "1 No. 12\" + 3 No.", "—", "Bases ready to receive cable"),
    ("4", "Cable, termination, insulation resistance and continuity testing",
     "LOC-03 · 1003", "4 No. runs", "—", "LOC-03 test records issued"),
    ("4", "Reinstatement of the cut; re-fix RRM.557 and RRM.670", "LOC-03 · 1003",
     "approx. 95 m · 2 No. RRM", "Per the awaited detail", "LOC-03 sealed; RRMs re-fixed"),
    ("5", "Saw cut, new base, cable, termination and testing", "LOC-02 · 1002",
     "approx. 24 m · 1 No. 12\"", "H5 — detail accepted", "LOC-02 complete and tested"),
    ("5", "Reinstatement of the cut; re-fix RRM.555", "LOC-02 · 1002",
     "approx. 24 m · 1 No. RRM", "Per the awaited detail", "LOC-02 sealed; RRM re-fixed"),
    ("5", "Testing and commissioning at circuit level — affected and unaffected fittings alike",
     "LOC-01 / 02 / 03", "4 No. circuits", "—", "Circuits energised and proved end to end"),
    ("5", "Remove the matt black vinyl masking from the direction signboards",
     "E4 / E6 approaches", "Per Phase 1 R4", "—", "Signboards returned to normal display"),
    ("5", "Final functionality check, then handover to Operations", "LOC-01 / 02 / 03",
     "19 No. fittings", "H4 — AGL / Operations", "Area returned to operational service"),
]
SCH_COLS = (640080, 4297680, 1371600, 1600200, 1828800, 4142232)
LOW_BODY_H = max(
    textfit.frame_height([(t, pt, False, 50800) for t, pt in lines], LOWER_BODY_W, 0, 0, 0, 0)
    for _, _, lines, _ in LOWER)
LOW_H = LOW_BODY_H + CHROME
LOW_Y = 10020300 - LOW_H
ROW_H = max(240030, min(320040, (LOW_Y - 182880 - TBL_Y - CHROME) // len(SCH_ROWS)))
tbl_h = ROW_H * len(SCH_ROWS)
pan_h = 109728 + 237744 + 128016 + tbl_h + 137160
box(sched, L, TBL_Y, W, pan_h, name="P08 PANEL · FIVE-DAY PROGRAMME")
textbox(sched, L + 146304, TBL_Y + 109728, W - 292608, 237744,
        [("FIVE-DAY PROGRAMME — PHASE 3 AGL INSTALLATION", 9.5, True, NAVY)], "SCHED HDR")

gf = sched.shapes.add_table(len(SCH_ROWS), 6, Emu(L + 146304),
                            Emu(TBL_Y + 109728 + 237744 + 128016),
                            Emu(sum(SCH_COLS)), Emu(tbl_h))
gf.name = "Table SCHEDULE"
st = gf.table
st.first_row = False
st.horz_banding = False
for ci, w in enumerate(SCH_COLS):
    st.columns[ci].width = Emu(w)
for r in st.rows:
    r.height = Emu(ROW_H)
for ri, row in enumerate(SCH_ROWS):
    for ci, val in enumerate(row):
        cell = st.cell(ri, ci)
        cell.margin_left = cell.margin_right = Emu(91440)
        cell.margin_top = cell.margin_bottom = Emu(45720)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if ri == 0 else (
            TINT if SCH_ROWS[ri][0] in ("2", "4") else RGBColor(0xFF, 0xFF, 0xFF))
        cell_borders(cell)
        para = cell.text_frame.paragraphs[0]
        _run(para, val, 9.0 if ri == 0 else 8.0, ri == 0 or ci == 0,
             RGBColor(0xFF, 0xFF, 0xFF) if ri == 0 else INK)
        inner = (SCH_COLS[ci] - 182880) / textfit.EMU_PT
        if textfit.text_w_pt(val, 9.0 if ri == 0 else 8.0, ri == 0 or ci == 0) > inner:
            print(f"   ! schedule wrap risk r{ri} c{ci}: {val!r}")

# ---- hold points and basis -------------------------------------------------
for x, title, lines, nm in LOWER:
    box(sched, x, LOW_Y, PAN_W, LOW_H, name=nm)
    textbox(sched, x + 146304, LOW_Y + 109728, LOWER_BODY_W, 237744,
            [(title, 9.5, True, NAVY)], nm + " HDR")
    textbox(sched, x + 146304, LOW_Y + 109728 + 237744 + 128016, LOWER_BODY_W, LOW_BODY_H,
            [(t, pt, False, INK) for t, pt in lines], nm + " BODY")

textbox(sched, L, 10149840, W, 365760,
        [("Programme for the AGL installation phase. It is issued for planning and for "
          "co-ordination with the civil programme; it is not a construction instruction until "
          "the hold points above are released.     ·     " + FOOTER_DOC, 7.5, False, MUTED)],
        "SHEET FOOTER")

# place it directly after the scope & sequence sheet
sld_lst = prs.slides._sldIdLst
items = list(sld_lst)
sld_lst.remove(items[-1])
sld_lst.insert(2, items[-1])
log("New sheet added after the scope & sequence sheet: a 5-day work schedule for the Phase 3 "
    "AGL installation, tied to the deck's own quantities and hold points.")

# ===================================== NEW SHEET - works quantities summary
qty = prs.slides.add_slide(prs.slide_layouts[0])

qtitle = textbox(qty, L, 274320, W - 1700000, 502920,
                 [("WORKS QUANTITIES — CORING, SAW CUT & SECONDARY CABLE  ·  TWY E (E4–E6)  ·  "
                   "REV P08", 20, True, NAVY)], "TITLE")
qtitle.text_frame.word_wrap = False
fit_font(qtitle, W - 1700000, 20.0, 13.0)
add_logo(qty, L + W, 265430, 434340)

# ---- headline tiles --------------------------------------------------------
TILES = [
    ("19 No.", "AFFECTED FITTINGS", "14 · 1 · 4 across LOC-01 / 02 / 03"),
    ("11 No.", "EXISTING BASES CORED OUT", "6 @ 8\"  ·  5 @ 12\""),
    ("13 No.", "NEW SIDE-ENTRY BASES", "11 @ 8\"  ·  2 @ 12\""),
    ("approx. 420 m", "SAW CUTTING", "300 m  ·  24 m  ·  95 m"),
    ("19 No.", "SECONDARY CABLE RUNS", "no joints — manhole to light"),
]
TILE_H = 1005840
for n, (big, cap, sub) in enumerate(TILES):
    x = L + n * (BOX_W + GAP5)
    box(qty, x, 914400, BOX_W, TILE_H, name=f"QTY TILE {n+1}")
    textbox(qty, x + 137160, 914400 + 91440, BOX_W - 274320, 640080,
            [(big, 18, True, NAVY), (cap, 8.5, True, INK), (sub, 8, False, MUTED)],
            f"QTY TILE TXT {n+1}")

# ---- the quantities table --------------------------------------------------
QTY_Y = 914400 + TILE_H + 274320
QTY_ROWS = [
    ("Location", "Affected assets", "Core out\nexisting 8\"", "Core out\nexisting 12\"",
     "New coring\n8\" base", "New coring\n12\" base", "Saw cutting",
     "Secondary cable required\n(2c x 4 sq.mm)"),
    ("LOC-01 · sheet 1001", "14 No.", "6 No.", "3 No. †", "11 No.", "—",
     "approx. 300 m", "14 No. runs  ·  approx. 300 m"),
    ("LOC-02 · sheet 1002", "1 No.", "—", "1 No.", "—", "1 No.",
     "approx. 24 m", "1 No. run  ·  approx. 24 m"),
    ("LOC-03 · sheet 1003", "4 No.", "—", "1 No.", "—", "1 No.",
     "approx. 95 m", "4 No. runs  ·  approx. 95 m"),
    ("TOTAL", "19 No.", "6 No.", "5 No.", "11 No.", "2 No.",
     "approx. 420 m", "19 No. runs  ·  approx. 420 m"),
]
QTY_AVAIL = W - 292608
need = []
for ci in range(8):
    wmax = 0
    for ri in range(len(QTY_ROWS)):
        for part in QTY_ROWS[ri][ci].split("\n"):
            wmax = max(wmax, textfit.text_w_pt(part, 9.0 if ri == 0 else 9.5,
                                               ri == 0 or ri == len(QTY_ROWS) - 1 or ci == 0))
    need.append(int(wmax * textfit.EMU_PT) + 182880 + 91440)
spare = QTY_AVAIL - sum(need)
QTY_COLS = tuple(w + spare // len(need) for w in need)
QTY_COLS = QTY_COLS[:-1] + (QTY_AVAIL - sum(QTY_COLS[:-1]),)

QTY_ROW_H = 411480
qtbl_h = QTY_ROW_H * len(QTY_ROWS)
qpan_h = 109728 + 237744 + 128016 + qtbl_h + 237744 + 137160
box(qty, L, QTY_Y, W, qpan_h, name="P08 PANEL · WORKS QUANTITIES")
textbox(qty, L + 146304, QTY_Y + 109728, W - 292608, 237744,
        [("WORKS QUANTITIES BY LOCATION — REV P08 FINAL SCOPE", 9.5, True, NAVY)], "QTY HDR")

qgf = qty.shapes.add_table(len(QTY_ROWS), 8, Emu(L + 146304),
                           Emu(QTY_Y + 109728 + 237744 + 128016),
                           Emu(sum(QTY_COLS)), Emu(qtbl_h))
qgf.name = "Table QUANTITIES"
qt = qgf.table
qt.first_row = False
qt.horz_banding = False
for ci, w in enumerate(QTY_COLS):
    qt.columns[ci].width = Emu(w)
for r in qt.rows:
    r.height = Emu(QTY_ROW_H)
for ri, row in enumerate(QTY_ROWS):
    last = ri == len(QTY_ROWS) - 1
    for ci, val in enumerate(row):
        cell = qt.cell(ri, ci)
        cell.margin_left = cell.margin_right = Emu(91440)
        cell.margin_top = cell.margin_bottom = Emu(45720)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if ri == 0 else (TINT if last else
                                                         RGBColor(0xFF, 0xFF, 0xFF))
        cell_borders(cell)
        tf = cell.text_frame
        for n, part in enumerate(val.split("\n")):
            para = tf.paragraphs[0] if n == 0 else tf.add_paragraph()
            _run(para, part, 9.0 if ri == 0 else 9.5, ri == 0 or last or ci == 0,
                 RGBColor(0xFF, 0xFF, 0xFF) if ri == 0 else INK)
            inner = (QTY_COLS[ci] - 182880) / textfit.EMU_PT
            if textfit.text_w_pt(part, 9.0 if ri == 0 else 9.5,
                                 ri == 0 or last or ci == 0) > inner:
                print(f"   ! quantities wrap risk r{ri} c{ci}: {part!r}")

textbox(qty, L + 146304, QTY_Y + 109728 + 237744 + 128016 + qtbl_h + 45720,
        W - 292608, 192024,
        [("†  A further 2 No. 12\" bases at LOC-01 were removed earlier and need no coring "
          "out; all 11 No. LOC-01 positions still take a new 8\" coring, which is why the "
          "core-out and new-coring counts differ.", 8.5, False, MUTED)], "QTY FOOTNOTE")

# ---- breakdown by works action ---------------------------------------------
ACT_ROWS = [
    ("Location", "Works action", "Assets", "Count", "New base", "Cable route"),
    ("LOC-01", "Core out existing 8\" → new side-entry shallow base",
     "SBC102-02/024, 01/027, 02/025, 01/028, 02/026, 01/029", "6 No.", "8\"", "Saw cut"),
    ("LOC-01", "Core out existing 12\" (3 No.) or removed earlier (2 No.) → new base",
     "TCCECH-04/034, 03/036, 04/035, 03/037, 04/036", "5 No.", "8\"", "Saw cut"),
    ("LOC-02", "Core out existing shallow base → new side-entry shallow base",
     "TCCECH-03/008", "1 No.", "12\"", "Saw cut"),
    ("LOC-03", "Core out existing 12\" → new side-entry shallow base — no side entry in the "
     "existing base", "TCCECH-03/003", "1 No.", "12\"", "Saw cut"),
    ("LOC-01", "No coring — secondary cable only; dummy plate protection during milling",
     "SBC102-02/027, TCCECH-03/035, TCCECH-03/018", "3 No.", "existing", "Saw cut"),
    ("LOC-03", "No coring — the existing 12\" shallow base takes the new side entry",
     "TCCECH-03/002, 04/002, 04/003", "3 No.", "existing", "Saw cut"),
]
ACT_AVAIL = W - 292608
aneed = []
for ci in range(6):
    wmax = max(textfit.text_w_pt(ACT_ROWS[ri][ci], 9.0 if ri == 0 else 8.5, ri == 0 or ci == 0)
               for ri in range(len(ACT_ROWS)))
    aneed.append(int(wmax * textfit.EMU_PT) + 182880 + 45720)
aspare = ACT_AVAIL - sum(aneed)
ACT_COLS = tuple(w + aspare // len(aneed) for w in aneed)
ACT_COLS = ACT_COLS[:-1] + (ACT_AVAIL - sum(ACT_COLS[:-1]),)

# ---- basis and exclusions --------------------------------------------------
QLOWER = (
    (L, "BASIS OF THE QUANTITIES", [
        ("Quantities are the Rev P08 final scope of work (30.07.2026) and reconcile with the "
         "consolidated field-governed scope sheet and the saw cut & coring schedule.", 9.0),
        ("Core out = removal of the existing shallow base by coring. New coring = the core "
         "taken to receive the new side-entry shallow base. The two counts differ at LOC-01 "
         "because 2 No. 12\" bases were removed earlier, so 11 No. new 8\" bases are set "
         "against 9 No. cored positions.", 9.0),
        ("6 No. fittings take new secondary cable with no coring — 3 No. at LOC-01 "
         "(SBC102-02/027, TCCECH-03/035, TCCECH-03/018), kept under a dummy plate during milling, and 3 No. at LOC-03 "
         "whose existing 12\" shallow base takes the new side entry. The saw cut routes plotted "
         "on sheets 1001 and 1003 already serve them.", 9.0),
        ("Cable per fitting: 1 No. 2-core 4 sq.mm secondary cable to each fitting, SBC and TCC "
         "alike — 7 No. SBC and 12 No. TCC, 19 No. runs in total (LOC-01 7 SBC + 7 TCC; LOC-02 "
         "1 TCC; LOC-03 4 TCC).", 9.0),
        ("Saw cut lengths are scaled from the Rev P08 saw cut runs on sheets 1001–1003 using "
         "the per-sheet drawing scale — approx. 300 m, 24 m and 95 m. For programme and "
         "enquiry only; confirm on site before ordering.", 9.0),
    ], "P08 PANEL · BASIS OF QUANTITIES"),
    (7699248, "EXCLUSIONS & ITEMS TO CONFIRM", [
        ("Core diameter and core depth for the 8\" and 12\" corings, and the grout / bedding "
         "specification for each, are per the awaited saw cut and side-entry shallow base "
         "detail. The figures above are counts, not dimensions.", 9.0),
        ("Saw cut width and depth, and the cover to the new secondary cable, are per the same "
         "awaited detail — hold point H5.", 9.0),
        ("Secondary cable figures are route lengths taken off the saw cut alignment. An "
         "allowance for terminations, base entry and manhole tails is to be added at ordering "
         "and is not included here.", 9.0),
        ("No duct is laid under the AGL scope at Rev P08. Saw cut is adopted as an interim "
         "arrangement agreed between the AGL team, the civil team and ADA AGL; permanent duct "
         "provision follows under the South Rehabilitation works. The existing 4 x 19 mm "
         "secondary duct bank is left in place with its obsolete cables withdrawn.", 9.0),
        ("Testing and commissioning is carried out at circuit level: the affected fittings and "
         "the fittings on the same circuits that fall outside these works are all proved before "
         "handover. No material is allowed here for those unaffected fittings.", 9.0),
        ("Civil works — milling, asphalt laying and pavement reinstatement outside the AGL saw "
         "cut — are excluded.", 9.0),
    ], "P08 PANEL · QUANTITY EXCLUSIONS"),
)
QLOW_BODY_H = max(
    textfit.frame_height([(t, pt, False, 50800) for t, pt in lines], LOWER_BODY_W, 0, 0, 0, 0)
    for _, _, lines, _ in QLOWER)
QLOW_H = QLOW_BODY_H + CHROME
QLOW_Y = 10020300 - QLOW_H

ACT_Y = QTY_Y + qpan_h + 182880
ACT_CHROME = 109728 + 237744 + 128016 + 137160
ACT_ROW_H = max(228600, min(457200,
                            (QLOW_Y - 182880 - ACT_Y - ACT_CHROME) // len(ACT_ROWS)))
atbl_h = ACT_ROW_H * len(ACT_ROWS)
box(qty, L, ACT_Y, W, atbl_h + ACT_CHROME, name="P08 PANEL · WORKS ACTION BREAKDOWN")
textbox(qty, L + 146304, ACT_Y + 109728, W - 292608, 237744,
        [("AFFECTED ASSETS BEHIND THE COUNTS — BY WORKS ACTION", 9.5, True, NAVY)], "ACT HDR")
agf = qty.shapes.add_table(len(ACT_ROWS), 6, Emu(L + 146304),
                           Emu(ACT_Y + 109728 + 237744 + 128016),
                           Emu(sum(ACT_COLS)), Emu(atbl_h))
agf.name = "Table ACTIONS"
at = agf.table
at.first_row = False
at.horz_banding = False
for ci, w in enumerate(ACT_COLS):
    at.columns[ci].width = Emu(w)
for r in at.rows:
    r.height = Emu(ACT_ROW_H)
for ri, row in enumerate(ACT_ROWS):
    for ci, val in enumerate(row):
        cell = at.cell(ri, ci)
        cell.margin_left = cell.margin_right = Emu(91440)
        cell.margin_top = cell.margin_bottom = Emu(45720)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if ri == 0 else (
            TINT if ACT_ROWS[ri][0] == "LOC-03" else RGBColor(0xFF, 0xFF, 0xFF))
        cell_borders(cell)
        _run(cell.text_frame.paragraphs[0], val, 9.0 if ri == 0 else 8.5, ri == 0 or ci == 0,
             RGBColor(0xFF, 0xFF, 0xFF) if ri == 0 else INK)
        inner = (ACT_COLS[ci] - 182880) / textfit.EMU_PT
        if textfit.text_w_pt(val, 9.0 if ri == 0 else 8.5, ri == 0 or ci == 0) > inner:
            print(f"   ! action table wrap risk r{ri} c{ci}: {val!r}")

for x, title, lines, nm in QLOWER:
    box(qty, x, QLOW_Y, PAN_W, QLOW_H, name=nm)
    textbox(qty, x + 146304, QLOW_Y + 109728, LOWER_BODY_W, 237744,
            [(title, 9.5, True, NAVY)], nm + " HDR")
    textbox(qty, x + 146304, QLOW_Y + 109728 + 237744 + 128016, LOWER_BODY_W, QLOW_BODY_H,
            [(t, pt, False, INK) for t, pt in lines], nm + " BODY")

textbox(qty, L, 10149840, W, 365760,
        [("Quantities are for check by the Engineer and are stated against the Rev P08 final "
          "scope of work (30.07.2026). Field condition governs.     ·     " + FOOTER_DOC,
          7.5, False, MUTED)], "SHEET FOOTER")

# sits ahead of the asset-by-asset consolidated scope sheet
items = list(sld_lst)
sld_lst.remove(items[-1])
sld_lst.insert(6, items[-1])
log("New sheet added ahead of the consolidated scope sheet: works quantities by location — "
    "affected assets, core out and new coring at 8\" and 12\", saw cutting and secondary cable.")

prs.save(OUT)
print("saved", OUT, "\n")
for c in changes:
    print(" -", c)
