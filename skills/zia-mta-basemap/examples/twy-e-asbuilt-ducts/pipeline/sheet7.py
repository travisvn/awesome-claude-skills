"""Sheet 7 -- derived secondary feed topology, and the discrepancy it exposes in the
Rev P05 LOC-02 scope note."""
import json, collections
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from proof_sheet import tb, panel, NAVY, INK, MUT, RULE

def build(prs):
    t=json.load(open('topology.json'))
    slide=prs.slides.add_slide(prs.slide_layouts[0])
    for ph in list(slide.placeholders): ph._element.getparent().remove(ph._element)
    W=prs.slide_width; M=320040
    tb(slide,M,164592,W-2*M,457200,
       'SHEET 7 — SECONDARY FEED TOPOLOGY (DERIVED)  ·  STAR FROM TRANSFORMER HANDHOLES, NOT A CHAINED RUN',
       size=15,bold=True,color=NAVY)
    x0,y0 = M, 750000
    FW = W-2*M

    # ---- finding, full width ----------------------------------------------------------
    panel(slide,x0,y0,FW,1_150_000,fill=RGBColor.from_string('FFF4E5'))
    tb(slide,x0+120000,y0+110000,FW-240000,240000,
       'FINDING — THE Rev P05 LOC-02 SCOPE NOTE DESCRIBES THE WRONG TOPOLOGY',size=10,bold=True,
       color=RGBColor.from_string('A64B00'))
    tb(slide,x0+120000,y0+420000,FW-240000,900000,[
      'Rev P05, sheet 3 (LOC-02) states: "SECONDARY DUCT: CHAINED RUN HH.E.056 → 03/007 → 04/007 → '
      '03/008 → 04/008 → 03/009 → HH.E.054". The source drawing shows no such chain.',
      'Every affected light at LOC-02 is fed by its own individual conduit spur radiating from ONE '
      'transformer handhole — five separate home runs of 20.5, 21.1, 29.1, 30.3 and 41.5 m, 142.5 m in '
      'total, not a series loop between lights. This changes the cable quantity, the number of pulls and '
      'the isolation sequence. Re-price and re-plan LOC-02 before ordering cable.',
      ],size=8.5,color=INK)

    # ---- two narrative panels side by side --------------------------------------------
    y1=y0+1_270_000; hw=(FW-200000)//2
    panel(slide,x0,y1,hw,1_960_000,fill=RGBColor.from_string('FFFFFF'))
    tb(slide,x0+120000,y1+110000,hw-240000,240000,'WHY THIS ALSO PROVES THE GEOMETRY IS REAL',size=10,bold=True,color=NAVY)
    tb(slide,x0+120000,y1+420000,hw-240000,1_800_000,[
      'The conduit fragments were chained by snapping endpoints within 0.30 m. What comes out is not an '
      'arbitrary set of lines — it is a working secondary distribution network:',
      '·  Every hub node lands on a CV_ETRANS HH (existing transformer handhole) fixture at 0.000–0.002 m. '
      'Four hubs at LOC-01, one at LOC-02, two at LOC-03.',
      '·  40 of 45 spurs terminating at a light land within 50 mm of that light\'s insertion point; most '
      'are exactly 0.00 m.',
      '·  Spur lengths are 6–42 m, consistent with transformer-handhole-to-light home runs.',
      'Straight light-to-nearest-pit construction lines cannot produce that. Invented or inferred geometry '
      'would not terminate on fixture insertion points to the millimetre, nor converge specifically on '
      'transformer handholes rather than on any nearby pit.',
      ],size=8.5,color=INK)

    x2=x0+hw+200000
    panel(slide,x2,y1,hw,1_960_000,fill=RGBColor.from_string('FFFFFF'))
    tb(slide,x2+120000,y1+110000,hw-240000,240000,'METHOD AND ITS LIMIT',size=10,bold=True,color=NAVY)
    tb(slide,x2+120000,y1+420000,hw-240000,1_800_000,[
      'The source drawing stores these conduits as loose LINE / LWPOLYLINE segments with no connectivity. '
      'The chaining here is DERIVED, not read from the drawing: endpoints within 0.30 m were treated as one '
      'node, then each branch was walked outward from its hub until it reached a leaf or a junction.',
      'Consequence: a spur length is the drawn geometry of a branch, not a proven pull. Two conduits passing '
      'within 0.30 m without physically joining would be merged; a genuine joint drawn with a wider gap '
      'would be split. Prove each run on site before ordering.',
      'The 0.30 m tolerance sits far below the shortest real spur (6.2 m) and well above the drafting gaps '
      'observed. No spur changed hop count when the tolerance was varied between 0.10 and 0.50 m.',
      ],size=8.5,color=INK)

    # ---- schedule: one column per location --------------------------------------------
    y2=y1+2_080_000
    ph=prs.slide_height-520000-y2
    colw=(FW-2*160000)//3
    cw2=[640000, 1_620_000, 1_150_000, 900000]
    for k,loc in enumerate(('LOC-01','LOC-02','LOC-03')):
        cx=x0+k*(colw+160000)
        panel(slide,cx,y2,colw,ph,fill=RGBColor.from_string('FFFFFF'))
        r=t[loc]
        lights=[s for s in r['runs'] if 'light' in s['end_type'].lower() or 'RRM' in s['end_type']]
        tb(slide,cx+110000,y2+100000,colw-220000,240000,
           f'{loc} — DERIVED SPUR SCHEDULE',size=9.5,bold=True,color=NAVY)
        tb(slide,cx+110000,y2+330000,colw-220000,200000,
           f'{len(r["hubs"])} transformer-handhole hub(s) · {len(r["runs"])} spurs · '
           f'{sum(s["length_m"] for s in lights):.1f} m to lights',size=7.5,color=MUT)
        nrow=len(r['runs'])
        pitch=min(188000, int((ph-560000-205000-300000)/max(nrow,1)))
        fsz=7 if pitch>=170000 else 6.4
        ty=y2+560000
        for c,h in enumerate(['Length','Terminates at','Source layer','Offset']):
            tb(slide,cx+110000+sum(cw2[:c]),ty,cw2[c],200000,h,size=6.5,bold=True,color=MUT)
        ty+=205000
        for s in sorted(r['runs'], key=lambda z:-z['length_m']):
            flag = s['end_dist']>0.05
            row=[f"{s['length_m']:.2f} m", s['end_type'], s['end_fixture'],
                 f"{s['end_dist']:.2f} m" + (' !' if flag else '')]
            for c,v in enumerate(row):
                tb(slide,cx+110000+sum(cw2[:c]),ty,cw2[c],pitch-3000,v,size=fsz,
                   color=(RGBColor.from_string('A64B00') if (flag and c==3) else INK))
            ty+=pitch
        tb(slide,cx+110000,y2+ph-260000,colw-220000,200000,
           '!  offset >50 mm — confirm the termination on site',size=6.5,color=MUT)

    tb(slide,M,prs.slide_height-420000,W-2*M,300000,
       'Spur lengths are drawn-geometry branch lengths at a 0.30 m snap tolerance — not proven pulls. '
       'topology.json accompanies this deck.',size=7.5,color=MUT)
    return slide
