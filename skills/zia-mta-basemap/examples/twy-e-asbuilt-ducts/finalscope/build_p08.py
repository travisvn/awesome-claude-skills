#!/usr/bin/env python3
"""Build Rev P08 — the final scope of work — from the Rev P06 as-built deck.

    python finalscope/build_p08.py

Input : out/TWY-E-AGL-SHOPDWG-ASBUILT-DUCTS_RevP06.pptx   (8 sheets)
        finalscope/scope_final.json                       (the governing scope)
        fieldsheet/marker_positions.json                  (verified marker centres)
Output: out/TWY-E-AGL-SHOPDWG_RevP08_FINAL-SCOPE.pptx     (10 sheets)

What the final scope changes
----------------------------
Every secondary route at all three locations becomes SAW CUT — no duct route
remains in the AGL scope — and every affected fitting takes a NEW SIDE-ENTRY
SHALLOW BASE: 8" core at LOC-01, 12" at LOC-02 and LOC-03. The affected-asset
list shrinks from 30 to 16.

Assets that Rev P06/P07 carried but the final scope does not name are given a
SUPERSEDED marker rather than being deleted or re-drawn as "field verified not
affected". Deleting them would erase the asset from the sheet entirely (a works
marker replaces the base symbol, it does not sit over it), and calling them
"not affected" would assert a field verification that was never carried out.

The red AGL works area is NOT redrawn. It is the Rev P05 field-condition milling
extent — what the civil team is milling — not a hull of the affected assets, so a
reduction in AGL scope does not shrink it. Re-cutting it here would be inventing
geometry the drawing has no source for.
"""

import json
import pathlib
import sys

HERE = pathlib.Path(__file__).parent
ROOT = HERE.parent
sys.path.insert(0, str(HERE))

from pptx import Presentation                                        # noqa: E402
from pptx.util import Inches                                         # noqa: E402

import legend_p08                                                    # noqa: E402
import sheet_scope                                                   # noqa: E402
import sheet_sawcut                                                  # noqa: E402
from deckkit import (GREY, INNER_D, MAGENTA, ORANGE, RED, RING_D,     # noqa: E402
                     by_id, cell_text, disc, find_oval, fit, lines_of,
                     move_slide, ring, set_lines, sub)

SRC = ROOT / "out" / "TWY-E-AGL-SHOPDWG-ASBUILT-DUCTS_RevP06.pptx"
OUT = ROOT / "out" / "TWY-E-AGL-SHOPDWG_RevP08_FINAL-SCOPE.pptx"

# slide index -> (location, header id, rev id, source id, meta id, scope id, notes id,
#                 legend body top)
SHEET = {
    1: dict(loc="LOC-01", rev=310, source=311, meta=312, scope=315, notes=318,
            legend_top=Inches(8.25), dwg="AUH-SK-AGL-TWYE-001-1001", scope_pt=7.2),
    2: dict(loc="LOC-02", rev=107, source=108, meta=109, scope=112, notes=115,
            legend_top=Inches(7.75), dwg="AUH-SK-AGL-TWYE-001-1002", scope_pt=7.2),
    3: dict(loc="LOC-03", rev=302, source=303, meta=304, scope=307, notes=310,
            legend_top=Inches(7.75), dwg="AUH-SK-AGL-TWYE-001-1003", scope_pt=6.8),
}

SUPERSEDED_LABEL = "SUPERSEDED — NOT IN REV P08 SCOPE"

SCOPE_PANEL = {
    "LOC-01": [
        "FINAL SCOPE (REV P08) — ALL 11 FITTINGS CORED, ALL ROUTES SAW CUT",
        "CORE OUT 8\" + NEW SIDE-ENTRY SHALLOW BASE + NEW CABLE: 11 No.",
        "   • SBC (6): SBC102-02/024, 01/027, 02/025, 01/028, 02/026, 01/029",
        "   • TCC (5): TCCECH-04/034, 03/036, 04/035, 03/037, 04/036",
        "ROUTE: SAW CUT throughout — the 6 No. SBC carried as VIA DUCT on Rev P06/P07 are saw cut at P08.",
        "   Saw cut detail per sheet SAW CUT & SIDE-ENTRY SHALLOW BASE DETAIL — detail drawing to be issued.",
        "REMOVED FROM SCOPE AT REV P08 — no works: SBC102-02/027, TCCECH-03/035, TCCECH-03/018",
        "FIELD VERIFIED NOT AFFECTED: SBC102-01/026 (no works)",
        "TCC103 FITTINGS NOT IN FIELD SCOPE — SHOWN AS EXISTING ONLY",
        "ISOLATE CIRCUITS: SBC102.01/.02, TCCECH.03/.04",
        "SITE RECORD 23.07.2026: 9 cored + 2 reinstatement already worked",
    ],
    "LOC-02": [
        "FINAL SCOPE (REV P08) — ONE FITTING AFFECTED, SAW CUT",
        "CORE OUT EXISTING SHALLOW BASE + NEW 12\" SIDE-ENTRY SHALLOW BASE + NEW CABLE: 1 No.",
        "   • TCCECH-03/008 — secondary route saw cut; existing shallow base cored out; new 12\" side-entry base",
        "   Saw cut detail per sheet SAW CUT & SIDE-ENTRY SHALLOW BASE DETAIL — detail drawing to be issued.",
        "REMOVED FROM SCOPE AT REV P08 — no works: TCCECH-03/007, 04/007, 04/008, 03/009",
        "   Rev P06/P07 carried all 5 No. as new-secondary-cable-only via duct. The final scope names 03/008 only.",
        "RRM.555 (0.41 m from cut): REMOVE / PROTECT BEFORE SAWCUT, RE-FIX AFTER PAVING",
        "ISOLATE CIRCUITS: TCCECH.03, TCCECH.04",
    ],
    "LOC-03": [
        "FINAL SCOPE (REV P08) — 4 No. TCC ONLY, SAW CUT",
        "NEW 12\" SIDE-ENTRY SHALLOW BASE + NEW CABLE: 4 No.",
        "   • TCCECH-03/002, 04/002, 03/003, 04/003",
        "ROUTE: SAW CUT — these 4 No. were carried as VIA DUCT on Rev P06/P07.",
        "   Saw cut detail per sheet SAW CUT & SIDE-ENTRY SHALLOW BASE DETAIL — detail drawing to be issued.",
        "REMOVED FROM SCOPE AT REV P08 — no works, 8 No. SBC:",
        "   • SBC102-01/038, 01/039, 02/035, 02/036 (EP7 STOP BAR — was sawcut, base protect)",
        "   • SBC102-01/040, 01/041, 02/037, 02/038 (was new cable only, via duct)",
        "RRM.557 (0.07 m — treat as within works) & RRM.670: REMOVE / PROTECT / RE-FIX",
        "ISOLATE CIRCUITS: TCCECH.03/.04 — SBC102.01/.02 carries no AGL works here at P08;",
        "   confirm isolation against the saw cut alignment before works.",
    ],
}

SOURCE_LINE = {
    "LOC-01": "FINAL SCOPE OF WORK 30.07.2026 (supersedes the Document_3 quantities) — "
              "FIELD CONDITION GOVERNS",
    "LOC-02": "FINAL SCOPE OF WORK 30.07.2026 (supersedes the Second_milling_area rev _1 "
              "quantities) — FIELD CONDITION GOVERNS",
    "LOC-03": "FINAL SCOPE OF WORK 30.07.2026 (supersedes the Second_milling_area rev _1 "
              "quantities) — FIELD CONDITION GOVERNS",
}

# Rev P06 carried "Coring at Location 1 only" on all three sheets. The final scope cores at
# all three, so that note is now wrong on every sheet and is replaced rather than appended to.
NOTE_2_OLD = "2. AGL works area (red) per field condition. Coring at Location 1 only."
NOTE_2_NEW = ("2. AGL works area (red) per field condition — the Rev P05 milling extent, "
              "UNCHANGED at Rev P08. Coring at ALL THREE locations at P08: 8\" at LOC-01, "
              "12\" at LOC-02 and LOC-03.")

# Carried on the scope panel rather than the notes: the notes panel on these sheets is already
# full to its lower edge, and these two statements govern how the scope panel itself is read.
SCOPE_TAIL = [
    "REV P08 SUPERSEDES the Rev P06/P07 quantities. SUPERSEDED = no works; NOT a field "
    "verification that the asset is unaffected.",
    "RED AGL WORKS AREA = Rev P05 field-condition milling extent, UNCHANGED at P08 — the civil "
    "milling limit, not a hull of the affected assets.",
]


# --------------------------------------------------------------------- markers

def remark(slide, loc, scope, positions):
    """Recolour every works marker on one sheet to the final scope."""
    done = []
    for light in scope["assets"]:
        p = positions[loc][light]
        cx, cy = Inches(p["cx"]), Inches(p["cy"])
        ring(find_oval(slide, cx, cy, RING_D), RED)          # core out + new base
        disc(find_oval(slide, cx, cy, INNER_D), ORANGE)      # via saw cut
        done.append(light)
    for light in scope["superseded"]:
        p = positions[loc][light]
        cx, cy = Inches(p["cx"]), Inches(p["cy"])
        ring(find_oval(slide, cx, cy, RING_D), GREY)
        disc(find_oval(slide, cx, cy, INNER_D), GREY)
        done.append(light)
    for light in scope["not_affected"]:                      # already grey, verify only
        p = positions[loc][light]
        r = find_oval(slide, Inches(p["cx"]), Inches(p["cy"]), RING_D)
        if str(r.line.color.rgb) != GREY:
            raise AssertionError(f"{loc} {light}: expected a grey ring, found {r.line.color.rgb}")
        done.append(light)
    return done


def audit(slide, loc, scope, positions):
    """Every named asset carries exactly the marker the final scope implies."""
    problems = []
    want = ([(a, RED, ORANGE) for a in scope["assets"]]
            + [(a, GREY, GREY) for a in scope["superseded"]]
            + [(a, GREY, None) for a in scope["not_affected"]])
    for light, ring_c, inner_c in want:
        p = positions[loc][light]
        cx, cy = Inches(p["cx"]), Inches(p["cy"])
        try:
            r = find_oval(slide, cx, cy, RING_D)
        except LookupError as e:
            problems.append(f"{loc} {light}: {e}")
            continue
        if str(r.line.color.rgb) != ring_c:
            problems.append(f"{loc} {light}: ring {r.line.color.rgb} != {ring_c}")
        if inner_c is None:
            continue
        try:
            inner = find_oval(slide, cx, cy, INNER_D)
        except LookupError as e:
            problems.append(f"{loc} {light}: inner {e}")
            continue
        if str(inner.fill.fore_color.rgb) != inner_c:
            problems.append(f"{loc} {light}: inner {inner.fill.fore_color.rgb} != {inner_c}")
    return problems


# ------------------------------------------------------------------ table rows

def table_rows(scope):
    rows = []
    for loc in ("LOC-01", "LOC-02", "LOC-03"):
        s = scope[loc]
        for a in s["assets"]:
            rows.append((loc, a, s["action"], "Sawcut", s["remark"]))
        for a in s["superseded"]:
            rows.append((loc, a, "NOT IN REV P08 SCOPE", "—",
                         "In the Rev P06/P07 scope; not named in the final scope of work "
                         "30.07.2026. No works. Not a field verification of 'not affected'"))
        for a in s["not_affected"]:
            rows.append((loc, a, "NOT AFFECTED", "—", "Field verified — no works"))
        for a in s["rrm"]:
            rows.append((loc, a, "RRM REMOVE / PROTECT", "—",
                         "Remove / protect before saw cut, re-fix after paving"))
    return rows


# ------------------------------------------------------------------------ main

def main():
    scope = json.loads((HERE / "scope_final.json").read_text())
    positions = json.loads((ROOT / "fieldsheet" / "marker_positions.json").read_text())
    prs = Presentation(SRC)

    problems = []

    # ---------------------------------------------------------- location sheets
    for ix, meta in SHEET.items():
        slide = prs.slides[ix]
        loc = meta["loc"]
        s = scope[loc]

        remark(slide, loc, s, positions)
        problems += audit(slide, loc, s, positions)

        sub(by_id(slide, 2), "REV P06", "REV P08")
        hdr = by_id(slide, 2)
        for placeholder in (f"[XXX-ELE-SHD-{meta['dwg'][-4:]}]",):
            if placeholder in hdr.text_frame.text:
                sub(hdr, placeholder, meta["dwg"])

        sub(by_id(slide, meta["rev"]), "REV P06 (EDITABLE)", "REV P08 (FINAL SCOPE)")
        set_lines(by_id(slide, meta["source"]), [SOURCE_LINE[loc]])
        sub(by_id(slide, meta["meta"]), "Rev P06", "Rev P08 · issued 30.07.2026")
        scope_panel = by_id(slide, meta["scope"])
        set_lines(scope_panel, SCOPE_PANEL[loc] + SCOPE_TAIL)
        fit(scope_panel, meta["scope_pt"])

        notes = by_id(slide, meta["notes"])
        existing = lines_of(notes)
        if NOTE_2_OLD not in existing:
            problems.append(f"{loc}: note 2 is not the text Rev P06 carried — not replaced")
        set_lines(notes, [NOTE_2_NEW if t == NOTE_2_OLD else t for t in existing])
        fit(notes, 7.0)

        marker_rows = [(f"CORE OUT {s['core_dia']} + SIDE-ENTRY BASE (SAWCUT)", RED, ORANGE),
                       (SUPERSEDED_LABEL, GREY, GREY)]
        if s["not_affected"]:
            marker_rows.append(("FIELD VERIFIED — NOT AFFECTED", GREY, None))
        if s["rrm"]:
            marker_rows.append(("RRM — REMOVE / PROTECT / RE-FIX", MAGENTA, "FFFFFF"))
        legend_p08.rebuild(slide, marker_rows, meta["legend_top"])

    # -------------------------------------------------------- consolidated table
    s5 = prs.slides[4]
    sub(by_id(s5, 2), "REV P06", "REV P08 (FINAL SCOPE)")
    table = next(sh for sh in s5.shapes if sh.has_table).table
    rows = table_rows(scope)
    if len(rows) != len(table.rows) - 1:
        problems.append(f"table has {len(table.rows) - 1} data rows, scope needs {len(rows)}")
    for i, row in enumerate(rows, start=1):
        if i >= len(table.rows):
            break
        for j, text in enumerate(row):
            cell_text(table.rows[i].cells[j], text)

    n_core = sum(len(scope[l]["assets"]) for l in ("LOC-01", "LOC-02", "LOC-03"))
    n_sup = sum(len(scope[l]["superseded"]) for l in ("LOC-01", "LOC-02", "LOC-03"))
    n_rrm = sum(len(scope[l]["rrm"]) for l in ("LOC-01", "LOC-02", "LOC-03"))
    set_lines(by_id(s5, 4), [
        f"Totals (Rev P08 final scope): {n_core} fittings — 11 cored 8\" at LOC-01, "
        f"5 cored 12\" at LOC-02/03 — all with a new side-entry shallow base and a new "
        f"secondary cable, all via SAW CUT.  ·  {n_sup} assets removed from scope at P08 "
        f"(in the Rev P06/P07 scope, not in the final scope — no works).  ·  {n_rrm} RRMs "
        f"remove / protect / re-fix.  ·  1 field verified not affected.  ·  No duct route "
        f"remains in the AGL scope; the \"Route\" column reads Sawcut throughout by "
        f"instruction, not by measurement off the as-built duct geometry."
    ])

    # ------------------------------------------------------------- cover sheet
    c = prs.slides[0]
    sub(by_id(c, 3), "(ADIA)", "(ZIA)")
    set_lines(by_id(c, 5), ["AUH-SK-AGL-TWYE-001 (sheets 1001 / 1002 / 1003)  ·  "
                            "Revision P08 — FINAL SCOPE OF WORK"])
    set_lines(by_id(c, 7), ["Final scope of work issued 30.07.2026 — supersedes the Rev P06 / "
                            "P07 field-sheet quantities. Field verification sheets: Document_3 "
                            "(LOC-01) · Second_milling_area rev _1 (LOC-02 / LOC-03)."])
    sub(by_id(c, 11), "Issued 28.07.2026", "Issued 28.07.2026 · Rev P08 30.07.2026")
    set_lines(by_id(c, 12), [
        "REV P08 — the final scope of work. All secondary routes at all three locations are "
        "SAW CUT; no duct route remains in the AGL scope. Every affected fitting takes a new "
        "SIDE-ENTRY SHALLOW BASE — 8\" core at Location 1, 12\" at Locations 2 and 3. The "
        "affected-asset list reduces from 30 to 16. A saw cut and side-entry base detail sheet "
        "is added and is held for the detail drawing. The as-built duct geometry, registration "
        "and stated limits of Rev P06 are carried forward unchanged."
    ])

    # ------------------------------------------------------------- new sheets
    sheet_scope.build(prs)            # appended, then moved to position 2
    move_slide(prs, len(prs.slides) - 1, 1)
    sheet_sawcut.build(prs)           # appended, then moved in after the table
    move_slide(prs, len(prs.slides) - 1, 6)

    # Rev P06's legend rebuild sized the legend background panels off a fixed constant and
    # left them hanging ~0.03" past the bottom of the sheet on all three location sheets.
    # Harmless on screen, but a shape off the sheet edge is a defect on an issued drawing.
    clamped = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.left is None or sh.width is None:
                continue
            if sh.top + sh.height > prs.slide_height:
                sh.height = prs.slide_height - sh.top
                clamped += 1
            if sh.left + sh.width > prs.slide_width:
                sh.width = prs.slide_width - sh.left
                clamped += 1

    prs.save(OUT)

    if problems:
        print("VERIFICATION FAILED")
        for p in problems:
            print("  -", p)
        raise SystemExit(1)

    print(f"{OUT.name} — {len(prs.slides)} sheets"
          + (f" · {clamped} shape edge(s) clamped onto the sheet" if clamped else ""))
    print(f"  in scope {n_core} · superseded {n_sup} · RRM {n_rrm}")
    for loc in ("LOC-01", "LOC-02", "LOC-03"):
        s = scope[loc]
        print(f"  {loc}: {len(s['assets'])} @ {s['core_dia']} side-entry, sawcut · "
              f"{len(s['superseded'])} superseded")


if __name__ == "__main__":
    main()
