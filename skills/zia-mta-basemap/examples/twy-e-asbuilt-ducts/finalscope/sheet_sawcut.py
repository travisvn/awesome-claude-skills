"""Sheet 7 — SAW CUT & SIDE-ENTRY SHALLOW BASE DETAIL.

New at Rev P08, and deliberately a HOLD sheet. The final scope moves every
secondary route to saw cut and every new base to side entry, and says the saw cut
detail drawing will be provided. It has not been provided yet, so this sheet
carries the empty frame for it, what the detail has to answer before it can be
built to, and the quantities that will hang off it. It does not invent a section.
"""

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from deckkit import BODY, NAVY
from sheetkit import footer, new_sheet, stack, table

RED = "B3261E"
AMBER = "8A6900"

MUST_SHOW = [
    ("Cut geometry", True, NAVY),
    ("1.  Saw cut width and depth, and the tolerance on each.", False, BODY),
    ("2.  Minimum cover from finished pavement surface to the top of the new secondary cable.",
     False, BODY),
    ("3.  Cut termination at the light base — how the cut meets the side entry, and the "
     "radius or chamfer at the entry.", False, BODY),
    ("4.  Minimum offset from slab joints, existing cuts and the milling edge, and what to do "
     "where the cut would run within that offset.", False, BODY),
    ("Cable and backfill", True, NAVY),
    ("5.  Cable protection within the cut — sleeve, tape or direct — and its fixing.", False, BODY),
    ("6.  Backfill and sealant specification, layer by layer, with cure time before the area "
     "returns to operational service.", False, BODY),
    ("7.  Reinstatement of the pavement surface and the finished surface tolerance.", False, BODY),
    ("Transition", True, NAVY),
    ("8.  How the saw cut route transitions into the existing duct in the NON-CONSTRUCTION "
     "area — this is open technical query Q3 and it governs the ends of every run.", False, RED),
    ("9.  Treatment where the saw cut crosses the existing 4 x 110 mm secondary duct bank "
     "exposed at 50 mm milling depth.", False, RED),
]

BASE_DETAIL = [
    ("10.  Side-entry shallow base — depth, entry height above the base floor, entry bore and "
     "the seal at the entry.", False, BODY),
    ("11.  Core diameter and core depth for 8\" and for 12\", and the grout / bedding "
     "specification for each.", False, BODY),
    ("12.  At Location 2, the method for coring out the EXISTING shallow base and making good "
     "before the new 12\" base is set.", False, BODY),
    ("13.  Setting-out tolerance of the base relative to the as-built fitting position, and "
     "the level tolerance to finished pavement.", False, BODY),
]

CLOSING = [
    ("Items 8 and 9 are the two that are not purely detailing. Item 8 is open technical query "
     "Q3 — the transition into the non-construction area was never answered, and every run has "
     "two ends. Item 9 is a physical clash with a duct bank the field already found at 50 mm. "
     "Neither can be closed out on site.", False, RED),
]

SCHEDULE = [
    ["Location", "Fittings", "Core", "New base", "Route", "Sheet"],
    ["LOC-01", "11  (6 SBC · 5 TCC)", "8\"", "Side-entry shallow base", "Saw cut", "1001"],
    ["LOC-02", "1  (TCCECH-03/008)", "12\"", "Side-entry shallow base — existing base cored out",
     "Saw cut", "1002"],
    ["LOC-03", "4  (TCC only)", "12\"", "Side-entry shallow base", "Saw cut", "1003"],
    ["TOTAL", "16", "11 @ 8\" · 5 @ 12\"", "16 side-entry shallow bases", "Saw cut throughout", "—"],
]


def _hold_frame(slide, x, y, w, h):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor.from_string("FBFCFD")
    box.line.color.rgb = RGBColor.from_string(RED)
    box.line.width = Pt(1.5)
    ln = box._element.find(".//" + qn("a:ln"))
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    box.shadow.inherit = False
    box.name = "P08 HOLD · SAW CUT DETAIL FRAME"

    tb = slide.shapes.add_textbox(x + Inches(0.4), y + h / 2 - Inches(0.85),
                                  w - Inches(0.8), Inches(1.7))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, (text, size, bold, colour) in enumerate([
        ("SAW CUT DETAIL DRAWING TO BE ISSUED", 15, True, RED),
        ("The final scope of work records that the saw cut detail drawing will be provided. "
         "It has not been issued at the date of this revision.", 9.5, False, BODY),
        ("Paste the issued detail into this frame and re-issue the sheet. Until it is issued "
         "and accepted, NO SAW CUTTING IS TO COMMENCE — hold point H6.", 9.5, True, RED),
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        try:
            from pptx.enum.text import PP_ALIGN
            p.alignment = PP_ALIGN.CENTER
        except ImportError:
            pass
        r = p.add_run()
        r.text = text
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = RGBColor.from_string(colour)


def build(prs):
    slide = new_sheet(
        prs, "SAW CUT & SIDE-ENTRY SHALLOW BASE DETAIL — TWY E (E4–E6)  ·  REV P08  ·  HOLD")

    lx, rx, lw, rw = Inches(0.5), Inches(8.42), Inches(7.61), Inches(7.61)

    # sized so the three left-column panels close above the footer at 11.10"
    frame_h = Inches(4.35)
    table_h = Inches(0.30) * len(SCHEDULE)
    (fx, fy, fw, _), (tx, ty, tw, _), _ = stack(slide, lx, Inches(1.00), lw, [
        ("DETAIL — HELD FOR ISSUE", [], 8.5, 0, frame_h),
        ("SAW CUT & CORING SCHEDULE — REV P08 FINAL SCOPE", [], 8.0, 0, table_h),
        ("SIDE-ENTRY SHALLOW BASE — WHAT THE DETAIL MUST ALSO SHOW", BASE_DETAIL, 10.5, 9, Emu(0)),
    ])
    _hold_frame(slide, fx, fy, fw, frame_h)
    table(slide, tx, ty, tw, SCHEDULE,
          [Inches(1.05), Inches(1.55), Inches(1.05), Inches(2.35), Inches(1.05), Inches(0.56)],
          size=9.0, head_size=9.5)

    stack(slide, rx, Inches(1.00), rw, [
        ("WHAT THE SAW CUT DETAIL MUST ANSWER BEFORE IT CAN BE BUILT TO",
         MUST_SHOW + CLOSING, 10.5, 9, Emu(0)),
    ])

    footer(slide,
           "This sheet is a HOLD. It states what the awaited detail has to resolve and the "
           "quantities that depend on it; it does not show a section, because none has been "
           "issued. Quantities are per the Rev P08 final scope of work (30.07.2026) and are "
           "FOR CHECK by the Engineer.     ·     AUH-SK-AGL-TWYE-001 REV P08  ·  Prepared: "
           "Mohammed, AGL Team Leader — ADB SAFEGATE  ·  Approved: Ragesh Menon, AGL Manager")
    return slide
