"""Rewrite column 1 of a location sheet's legend (the works-action rows).

The as-built asset column added at Rev P06 is left untouched — it describes the
base layer, not the scope, and nothing in the final scope changes it.

Rows are harvested from the sheet rather than redrawn, so swatch geometry, text
size and colour stay exactly as authored; only the count, the labels and the
swatch colours change. Rows are then re-packed at the sheet's own row pitch so
deleting or adding one leaves no gap.
"""

from pptx.util import Emu, Inches

from deckkit import RING_D, INNER_D, PANEL_X, clone, delete, disc, is_oval, ring, set_lines

COL2_X = Inches(13.0)      # anything at or right of this is the as-built asset column
ROW_GAP = Emu(160000)      # shapes closer than this vertically belong to one row


def _col1_rows(slide, top_min):
    items = [sh for sh in slide.shapes
             if sh.left is not None and PANEL_X <= sh.left < COL2_X
             and sh.top is not None and sh.top >= top_min
             and sh.width is not None and sh.width < Inches(3)]
    items.sort(key=lambda s: (s.top, s.left))
    rows = []
    for sh in items:
        if rows and sh.top - rows[-1][0].top < ROW_GAP:
            rows[-1].append(sh)
        else:
            rows.append([sh])
    return rows


def _parts(row):
    outer = next((s for s in row if is_oval(s) and abs(s.width - RING_D) < Inches(0.03)), None)
    inner = next((s for s in row if is_oval(s) and abs(s.width - INNER_D) < Inches(0.03)), None)
    label = next((s for s in row
                  if s.has_text_frame and s.text_frame.text.strip()), None)
    return outer, inner, label


def rebuild(slide, marker_rows, legend_body_top):
    """marker_rows: ordered [(label, ring_colour, inner_colour_or_None), ...]

    Linework rows (duct / works area / centreline / edge) keep their order and
    follow the marker rows, exactly as the sheet already has them.
    """
    rows = _col1_rows(slide, legend_body_top)
    classified = [(r, _parts(r)) for r in rows]
    marker_existing = [(r, p) for r, p in classified if p[0] is not None]
    line_rows = [r for r, p in classified if p[0] is None]

    if not marker_existing:
        raise LookupError("no works-action legend rows found on this sheet")

    pitch = None
    tops = sorted({min(s.top for s in r) for r in rows})
    if len(tops) > 1:
        pitch = min(b - a for a, b in zip(tops, tops[1:]) if b - a > ROW_GAP)
    if not pitch:
        pitch = Emu(272000)
    first_top = tops[0]

    proto_row, (proto_outer, proto_inner, proto_label) = marker_existing[0]
    if proto_inner is None:
        proto_inner = next(p[1] for _, p in marker_existing if p[1] is not None)

    kept = []
    for i, (label, ring_c, inner_c) in enumerate(marker_rows):
        if i < len(marker_existing):
            row, (outer, inner, lab) = marker_existing[i]
            row = list(row)
        else:
            outer = clone(slide, proto_outer)
            inner = clone(slide, proto_inner)
            lab = clone(slide, proto_label)
            row = [outer, inner, lab]
            base = min(s.top for s in marker_existing[-1][0])
            for s, src in ((outer, proto_outer), (inner, proto_inner), (lab, proto_label)):
                s.left = src.left
                s.top = base + (src.top - min(t.top for t in proto_row))
            inner = inner
        ring(outer, ring_c)
        if inner_c is None:
            if inner is not None:
                delete(inner)
                row = [s for s in row if s is not inner]
        else:
            if inner is None:
                inner = clone(slide, proto_inner)
                inner.left = proto_inner.left
                inner.top = outer.top + (proto_inner.top - proto_outer.top)
                row.append(inner)
            disc(inner, inner_c)
        set_lines(lab, [label])
        kept.append(row)

    for row, _ in marker_existing[len(marker_rows):]:
        for s in row:
            delete(s)

    for i, row in enumerate(kept + line_rows):
        top = min(s.top for s in row)
        dy = (first_top + i * pitch) - top
        for s in row:
            s.top = s.top + dy

    return len(kept), len(line_rows)
