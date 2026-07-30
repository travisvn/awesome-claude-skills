#!/usr/bin/env python3
"""Check the built Rev P08 deck against scope_final.json.

    python finalscope/check_p08.py

Reads only the saved .pptx — nothing from the build's memory — so it catches a
build that silently stopped short as well as one that marked the wrong asset.

Checks
  1. every asset named in the final scope carries the marker the scope implies,
     at the position verified against the source DXF in marker_positions.json
  2. no asset outside the final scope still carries an in-scope marker
  3. the consolidated table lists exactly the scope, once each
  4. no sheet still says "Coring at Location 1 only" or "REV P06" in its header
  5. no text box overflows the sheet
"""

import json
import pathlib
import sys

HERE = pathlib.Path(__file__).parent
ROOT = HERE.parent
sys.path.insert(0, str(HERE))

from pptx import Presentation                                    # noqa: E402
from pptx.util import Emu, Inches                                # noqa: E402

from deckkit import GREY, INNER_D, ORANGE, RED, RING_D, find_oval  # noqa: E402

DECK = ROOT / "out" / "TWY-E-AGL-SHOPDWG_RevP08_FINAL-SCOPE.pptx"
SHEET_IX = {"LOC-01": 2, "LOC-02": 3, "LOC-03": 4}   # after the scope sheet is inserted


def main():
    scope = json.loads((HERE / "scope_final.json").read_text())
    positions = json.loads((ROOT / "fieldsheet" / "marker_positions.json").read_text())
    prs = Presentation(DECK)
    bad = []

    # 1 + 2 — markers
    checked = 0
    for loc, ix in SHEET_IX.items():
        slide = prs.slides[ix]
        s = scope[loc]
        want = ([(a, RED, ORANGE) for a in s["assets"]]
                + [(a, GREY, GREY) for a in s["superseded"]]
                + [(a, GREY, None) for a in s["not_affected"]])
        named = {a for a, _, _ in want}
        for light, ring_c, inner_c in want:
            p = positions[loc][light]
            cx, cy = Inches(p["cx"]), Inches(p["cy"])
            try:
                r = find_oval(slide, cx, cy, RING_D)
            except LookupError as e:
                bad.append(f"{loc} {light}: {e}")
                continue
            if str(r.line.color.rgb) != ring_c:
                bad.append(f"{loc} {light}: ring {r.line.color.rgb}, expected {ring_c}")
            if inner_c is not None:
                try:
                    inner = find_oval(slide, cx, cy, INNER_D)
                except LookupError as e:
                    bad.append(f"{loc} {light}: inner {e}")
                    continue
                if str(inner.fill.fore_color.rgb) != inner_c:
                    bad.append(f"{loc} {light}: inner {inner.fill.fore_color.rgb}, "
                               f"expected {inner_c}")
            checked += 1
        # any other field-sheet light on this sheet must not read as in scope
        for light, p in positions[loc].items():
            if light in named:
                continue
            try:
                r = find_oval(slide, Inches(p["cx"]), Inches(p["cy"]), RING_D)
            except LookupError:
                continue
            if str(r.line.color.rgb) == RED:
                bad.append(f"{loc} {light}: carries an in-scope marker but is not in the "
                           f"final scope")

    # 3 — consolidated table
    table = next(sh for sh in prs.slides[5].shapes if sh.has_table).table
    rows = [(c.text for c in r.cells) for r in list(table.rows)[1:]]
    listed = {}
    for row in rows:
        cells = list(row)
        listed.setdefault(cells[1], []).append((cells[0], cells[2], cells[3]))
    for loc in SHEET_IX:
        s = scope[loc]
        for a in s["assets"]:
            got = listed.get(a)
            if not got:
                bad.append(f"table: {a} missing")
            elif len(got) > 1:
                bad.append(f"table: {a} listed {len(got)} times")
            elif got[0][2] != "Sawcut":
                bad.append(f"table: {a} route {got[0][2]!r}, expected 'Sawcut'")
        for a in s["superseded"]:
            got = listed.get(a)
            if not got:
                bad.append(f"table: superseded {a} missing")
            elif got[0][1] != "NOT IN REV P08 SCOPE":
                bad.append(f"table: superseded {a} action {got[0][1]!r}")
    n_scope = sum(len(scope[l][k]) for l in SHEET_IX
                  for k in ("assets", "superseded", "not_affected", "rrm"))
    if len(rows) != n_scope:
        bad.append(f"table has {len(rows)} rows, the scope has {n_scope} entries")

    # 4 — retired Rev P06 statements
    for ix, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text
            if "Coring at Location 1 only" in t:
                bad.append(f"sheet {ix}: still says 'Coring at Location 1 only'")
            if ix in (3, 4, 5) and "REV P06" in t and "Rev P06/P07" not in t:
                bad.append(f"sheet {ix}: still carries 'REV P06' — {t[:60]!r}")

    # 5 — overflow off the sheet
    for ix, slide in enumerate(prs.slides, start=1):
        for sh in slide.shapes:
            if sh.left is None or sh.width is None:
                continue
            if sh.left + sh.width > prs.slide_width or sh.top + sh.height > prs.slide_height:
                bad.append(f"sheet {ix}: {sh.name!r} runs off the sheet "
                           f"(right {Emu(sh.left + sh.width).inches:.2f}\", "
                           f"bottom {Emu(sh.top + sh.height).inches:.2f}\")")

    if bad:
        print(f"FAILED — {len(bad)} problem(s)")
        for b in bad:
            print("  -", b)
        raise SystemExit(1)
    print(f"OK — {checked} assets marked per the final scope, {len(rows)} table rows, "
          f"{len(prs.slides)} sheets")


if __name__ == "__main__":
    main()
