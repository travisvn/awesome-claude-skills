"""Sheet furniture for the two sheets added at Rev P08.

Both are drawn to the idiom the Rev P06 deck already uses: Arial throughout, a
20 pt navy title across the top, and content in white panels with a hairline
border and a small navy caps heading.
"""

import copy

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from deckkit import BODY, MUTED, NAVY, WHITE

SLIDE_W = Inches(16.53)
MARGIN = Inches(0.5)


def new_sheet(prs, title, background_from=4):
    """A blank sheet with the deck's background and its 20 pt navy title."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)

    src = prs.slides[background_from]._element.find(qn("p:cSld") + "/" + qn("p:bg"))
    if src is not None:
        cSld = slide._element.find(qn("p:cSld"))
        cSld.insert(0, copy.deepcopy(src))

    box = slide.shapes.add_textbox(MARGIN, Inches(0.30), Inches(15.5), Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Emu(0)
    r = tf.paragraphs[0].add_run()
    r.text = title
    r.font.name = "Arial"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = RGBColor.from_string(NAVY)
    return slide


def section(slide, x, y, w, h, heading):
    """A bordered white panel with a navy caps heading; returns the body origin."""
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor.from_string(WHITE)
    box.line.color.rgb = RGBColor.from_string("D3D8DE")
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    box.name = f"P08 PANEL · {heading[:28]}"

    hd = slide.shapes.add_textbox(x + Inches(0.16), y + Inches(0.12),
                                  w - Inches(0.32), Inches(0.26))
    tf = hd.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    r = tf.paragraphs[0].add_run()
    r.text = heading
    r.font.name = "Arial"
    r.font.size = Pt(9.5)
    r.font.bold = True
    r.font.color.rgb = RGBColor.from_string(NAVY)
    return x + Inches(0.16), y + Inches(0.44), w - Inches(0.32)


def body(slide, x, y, w, lines, size=8.0, colour=BODY, gap=Pt(3)):
    """lines: str, or (str, bold), or (str, bold, colour)."""
    h = Inches(0.22) * max(len(lines), 1)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        if isinstance(line, str):
            text, bold, col = line, False, colour
        elif len(line) == 2:
            text, bold = line
            col = colour
        else:
            text, bold, col = line
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = gap
        r = p.add_run()
        r.text = text
        r.font.name = "Arial"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = RGBColor.from_string(col)
    return tb


HEAD_H = Inches(0.44)      # panel top to first body line
PAD_B = Inches(0.18)       # last body line to panel bottom


def _wrapped(lines, w, size):
    """Line count once wrapped, deliberately pessimistic so a panel never clips.

    Arial's average glyph is a little under half an em; 0.0062 * size" per character
    under-counts the characters that fit, which over-counts the lines needed.
    """
    per_line = max(1, int(Emu(w).inches / (0.0062 * size)))
    n = 0
    for line in lines:
        text = line if isinstance(line, str) else line[0]
        n += max(1, -(-len(text) // per_line))
    return n


def panel_height(lines, w, size=8.0, gap_pt=6, extra=Emu(0)):
    """The panel height that fits `lines` at `size` in a body of width `w`."""
    n = _wrapped(lines, w - Inches(0.32), size)
    text_h = Inches(n * size * 1.22 / 72 + len(lines) * gap_pt / 72)
    return HEAD_H + text_h + PAD_B + extra


def stack(slide, x, y, w, blocks, gap=Inches(0.16)):
    """Lay panels down a column, each sized to its own content.

    blocks: (heading, lines, size, gap_pt, extra_height). `lines` may be empty when
    the panel's content is drawn by the caller into the returned body origin.
    """
    out, cur = [], y
    for heading, lines, size, gap_pt, extra in blocks:
        h = panel_height(lines, w, size, gap_pt, extra)
        bx, by, bw = section(slide, x, cur, w, h, heading)
        if lines:
            body(slide, bx, by, bw, lines, size=size, gap=Pt(gap_pt))
        out.append((bx, by, bw, h))
        cur = cur + h + gap
    return out


def footer(slide, text, y=Inches(11.10)):
    tb = slide.shapes.add_textbox(MARGIN, y, Inches(15.5), Inches(0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(7.5)
    r.font.color.rgb = RGBColor.from_string(MUTED)
    return tb


def table(slide, x, y, w, rows, col_w, size=8.0, head_size=8.5):
    """A table styled like the consolidated scope table: navy head, plain body."""
    n_rows, n_cols = len(rows), len(rows[0])
    gf = slide.shapes.add_table(n_rows, n_cols, x, y, w, Inches(0.26) * n_rows)
    tbl = gf.table
    tbl.first_row = True
    for j, cw in enumerate(col_w):
        tbl.columns[j].width = cw
    for i, row in enumerate(rows):
        tbl.rows[i].height = Inches(0.26)
        for j, text in enumerate(row):
            cell = tbl.cell(i, j)
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Emu(18000)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor.from_string(NAVY)
            else:
                cell.fill.background()
            tf = cell.text_frame
            tf.word_wrap = True
            r = tf.paragraphs[0].add_run()
            r.text = text
            r.font.name = "Arial"
            r.font.size = Pt(head_size if i == 0 else size)
            r.font.bold = i == 0
            r.font.color.rgb = RGBColor.from_string(WHITE if i == 0 else BODY)
    return gf
