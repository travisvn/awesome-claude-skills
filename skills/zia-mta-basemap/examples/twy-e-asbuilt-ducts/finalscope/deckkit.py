"""Shared helpers for editing the TWY E deck in place.

Everything here preserves the source deck's own formatting rather than recreating
it: text is rewritten run-by-run so each paragraph keeps the font it was authored
with, and new shapes are cloned from an existing shape of the same kind.
"""

import copy

from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

# deck palette, harvested from the Rev P06 sheets
NAVY = "1E2761"
GOLD = "C9A227"
BODY = "1F2937"
MUTED = "5F6368"
PANEL = "CADCFC"
WHITE = "FFFFFF"

RED = "CC0000"        # core out + new base
ORANGE = "E8710A"     # saw cut route
BLUE = "0055CC"       # duct route (retired from the works actions at P08)
GREEN = "1E8E3E"      # new cable only (retired at P08)
GREY = "9AA0A6"       # not affected / superseded
MAGENTA = "BB00BB"    # RRM remove-protect-refix

RING_D = Inches(0.23)
INNER_D = Inches(0.12)
PANEL_X = Inches(11.20)


# ------------------------------------------------------------------ shape kinds

def is_oval(sh):
    try:
        return "OVAL" in str(sh.auto_shape_type)
    except (ValueError, AttributeError):
        return False


def centre(sh):
    return (sh.left + sh.width // 2, sh.top + sh.height // 2)


def by_id(slide, shape_id):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise LookupError(f"no shape id={shape_id} on this slide")


def find_oval(slide, cx, cy, diameter, tol=Inches(0.02)):
    """The single oval of `diameter` centred within `tol` of (cx, cy)."""
    hits = [
        sh for sh in slide.shapes
        if is_oval(sh)
        and abs(sh.width - diameter) < Inches(0.01)
        and abs(centre(sh)[0] - cx) < tol
        and abs(centre(sh)[1] - cy) < tol
    ]
    if len(hits) != 1:
        raise LookupError(
            f"expected 1 oval d={Emu(diameter).inches:.3f}\" near "
            f"({Emu(cx).inches:.3f},{Emu(cy).inches:.3f}), found {len(hits)}"
        )
    return hits[0]


def delete(sh):
    sh._element.getparent().remove(sh._element)


def clone(slide, template):
    """Deep-copy a shape's XML so the copy inherits its exact styling."""
    new = copy.deepcopy(template._element)
    template._element.addnext(new)
    sh = next(s for s in slide.shapes if s._element is new)
    max_id = max(s.shape_id for s in slide.shapes if s._element is not new)
    sh._element.nvSpPr.cNvPr.set("id", str(max_id + 1))
    sh._element.nvSpPr.cNvPr.set("name", f"P08 {max_id + 1}")
    return sh


# ------------------------------------------------------------------- appearance

def ring(shape, colour, width_pt=2.25):
    shape.fill.background()
    shape.line.color.rgb = RGBColor.from_string(colour)
    shape.line.width = Pt(width_pt)


def disc(shape, colour, width_pt=0.5):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(colour)
    shape.line.color.rgb = RGBColor.from_string(colour)
    shape.line.width = Pt(width_pt)


# ------------------------------------------------------------------------- text

def _font_of(run):
    return dict(
        size=run.font.size,
        name=run.font.name,
        bold=run.font.bold,
        rgb=run.font.color.rgb if run.font.color and run.font.color.type is not None else None,
    )


def set_lines(shape, lines):
    """Replace a text frame's paragraphs, reusing each paragraph's own formatting.

    Paragraph i keeps the font of original paragraph i; extra paragraphs inherit
    the last original, so a longer scope panel still reads as one block.
    """
    tf = shape.text_frame
    fonts = [_font_of(p.runs[0]) for p in tf.paragraphs if p.runs]
    if not fonts:
        raise ValueError(f"shape id={shape.shape_id} has no runs to take a font from")
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        f = fonts[min(i, len(fonts) - 1)]
        r.font.size = f["size"]
        r.font.name = f["name"]
        r.font.bold = f["bold"]
        if f["rgb"] is not None:
            r.font.color.rgb = f["rgb"]


def fit(shape, size_pt):
    """Top-anchor a panel and set every run to `size_pt`.

    The Rev P06 scope and notes panels were filled to their lower edge at the text
    they carried. The final scope is longer on every sheet, so they are re-set a
    little smaller and pinned to the top of the panel — the boxes cannot grow
    without moving the legend block off the bottom of the sheet.
    """
    from pptx.enum.text import MSO_ANCHOR
    shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    shape.text_frame.word_wrap = True
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size_pt)


def sub(shape, old, new):
    """Replace `old` inside a text frame, keeping every run's own formatting."""
    hit = False
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if old in r.text:
                r.text = r.text.replace(old, new)
                hit = True
    if not hit:
        raise LookupError(f"{old!r} not found in shape id={shape.shape_id}")


def lines_of(shape):
    return [p.text for p in shape.text_frame.paragraphs if p.text.strip()]


def cell_text(cell, text, bold=None):
    tf = cell.text_frame
    proto = tf.paragraphs[0].runs[0]
    f = dict(size=proto.font.size, name=proto.font.name, bold=proto.font.bold,
             rgb=proto.font.color.rgb
             if proto.font.color and proto.font.color.type is not None else None)
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = f["size"]
    r.font.name = f["name"]
    r.font.bold = f["bold"] if bold is None else bold
    if f["rgb"] is not None:
        r.font.color.rgb = f["rgb"]


# ----------------------------------------------------------------- new textboxes

def textbox(slide, x, y, w, h, lines, size=9, bold=False, colour=BODY,
            spacing=None, name=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if spacing is not None:
            p.space_after = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = "Arial"
        r.font.color.rgb = RGBColor.from_string(colour)
    if name:
        tb.name = name
    return tb


def panel(slide, x, y, w, h, fill=WHITE, line=None, name=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = RGBColor.from_string(line)
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    if name:
        sh.name = name
    return sh


# -------------------------------------------------------------- slide ordering

def move_slide(prs, old_index, new_index):
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    sldIdLst.remove(ids[old_index])
    sldIdLst.insert(new_index, ids[old_index])
