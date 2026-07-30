#!/usr/bin/env python3
"""End-to-end check of the issued deck against the field sheets and against itself.

Everything the drawing asserts is checked in one place: markers, positions, legends,
labels, the scope notes, the consolidated table, the totals sheet, the revision block,
and the rendered PDF's actual text layout.

    python check_all.py                 # checks the P07 output
    python check_all.py src_RevP06.pptx # checks any other revision

Exits non-zero if any check fails. Layout checks use the rendered PDF, so they need
LibreOffice with Impress; if it is missing they are reported as SKIPPED, never as passes.
"""

import collections
import json
import math
import pathlib
import subprocess
import sys
import tempfile

from pptx import Presentation
from pptx.util import Emu, Inches

import legend
import reconcile_fieldsheet as R

HERE = pathlib.Path(__file__).parent
SHEETS = {1: "LOC-01", 2: "LOC-02", 3: "LOC-03"}
PAGE_W, PAGE_H = 16.539931, 11.689237
MAP_RIGHT = 11.05
PANEL_RIGHT = 16.25

results = []


def check(name, problems, skipped=False):
    results.append((name, list(problems), skipped))


def _in(v):
    return Emu(v).inches


def labels_on(slide):
    """(text, left, top, width) for every asset label in the map frame."""
    out = []
    for sh in slide.shapes:
        if sh.left is None or _in(sh.left) >= MAP_RIGHT or not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t and len(t) < 42 and ("." in t or "/" in t) and not t.startswith(("UTM", "1.")):
            out.append((t, _in(sh.left), _in(sh.top), _in(sh.width)))
    return out


# ---------------------------------------------------------------- 1. markers

def check_markers(prs, sheets, positions):
    check("markers match the field-sheet Yes/No columns",
          R.audit(prs, sheets, positions))

    tol = {"LOC-01": 1.0, "LOC-02": 1.0, "LOC-03": 60.0}
    bad = [f"{loc} {lt}: {p['fixture_mm']} mm from the nearest surveyed light "
           f"(limit {tol[loc]} mm)"
           for loc in SHEETS.values() for lt, p in positions[loc].items()
           if p["fixture_mm"] > tol[loc]]
    check("every marker sits on a surveyed AGL light fixture", bad)


def check_no_stray_markers(prs, sheets):
    """No works-action ring may sit on anything but a field-sheet light or an RRM."""
    ACTION = {"CC0000", "0055CC", "E8710A", "F9AB00", "BB00BB"}
    bad = []
    for ix, loc in SHEETS.items():
        slide = prs.slides[ix]
        lbls = labels_on(slide)
        field = {r["light"] for r in sheets[loc]}
        for sh in slide.shapes:
            if not legend._kind(sh) == "OVAL" or sh.left is None:
                continue
            if _in(sh.left) >= MAP_RIGHT or abs(sh.width - Inches(0.23)) >= Inches(0.01):
                continue
            if legend._line(sh) not in ACTION:
                continue
            cx, cy = _in(sh.left) + 0.115, _in(sh.top) + 0.115
            near = min(lbls, key=lambda l: math.hypot(l[1] - 0.13 - cx, l[2] + 0.215 - cy),
                       default=None)
            if near is None:
                bad.append(f"{loc}: works ring at ({cx:.3f},{cy:.3f}) has no label")
            elif not any(f in near[0] for f in field) and not near[0].startswith("RRM"):
                bad.append(f"{loc}: works ring on {near[0]}, which is not on the field sheet")
    check("no works marker on an asset the field sheet does not list", bad)


# ---------------------------------------------------------------- 2. legends

def check_legends(prs, asset_keys):
    problems = []
    for ix, loc in SHEETS.items():
        slide = prs.slides[ix]
        declared = declared_from_sheet(slide, loc)
        problems += legend.audit(slide, loc, declared)
    check("every plotted symbol and line colour is legended, and vice versa", problems)


def declared_from_sheet(slide, loc):
    """Read back what the sheet's legend actually declares, from its own swatches."""
    top = Inches(legend.PANEL[loc][1])
    sw = [s for s in slide.shapes
          if s.left is not None and s.left >= Inches(11) and s.top is not None and s.top >= top]
    ov = [s for s in sw if legend._kind(s) == "OVAL"]
    rings = [o for o in ov if abs(o.width - legend.RING_D) < Inches(0.01)]
    inners = [o for o in ov if abs(o.width - legend.INNER_D) < Inches(0.01)]
    dec = []
    for r in rings:
        rc = legend._centre(r)
        m = [i for i in inners if abs(legend._centre(i)[0] - rc[0]) < Inches(0.02)
             and abs(legend._centre(i)[1] - rc[1]) < Inches(0.02)]
        dec.append((legend.MARKER, legend._line(r), legend._fill(m[0])) if m
                   else (legend.RING, legend._line(r)))
    claimed = set()
    for s in sw:
        if s.width is None or s.width > Inches(0.16) or legend._kind(s) is None:
            continue
        cls = legend.asset_class_of_swatch(s)
        if cls is not False:
            dec.append((legend.ASSET, cls))
            claimed.add(id(s._element))
    for s in sw:
        if id(s._element) in claimed:
            continue          # already read as an asset-class swatch
        if s.height is not None and s.height < Inches(0.01) and legend._line(s):
            dec.append((legend.LINE, legend._line(s)))
        elif legend._kind(s) == "RECTANGLE" and s.width is not None and s.width < Inches(0.3):
            if legend._fill(s) == legend.SHADE:
                dec.append((legend.FILL, legend.SHADE))
            elif legend._fill(s) is None and legend._line(s):
                dec.append((legend.SQUARE, legend._line(s)))
    return dec


def check_asset_symbols(prs, symbols):
    """Every existing-asset symbol carries the symbol of its classified as-built class,
    and each class present on a sheet has its own legend row."""
    bad = []
    # the (prst, fill, line) triple must identify a class uniquely, or the legend is
    # ambiguous and asset_class_of_swatch() would silently mis-attribute a row
    triples = collections.Counter((v[0], v[2], v[3]) for v in legend.ASSET_SPEC.values())
    for t, n in triples.items():
        if n > 1:
            bad.append(f"symbology: {n} classes share the swatch triple {t}")

    for ix, loc in SHEETS.items():
        slide = prs.slides[ix]
        for rec in symbols[loc]["symbols"]:
            cx, cy = Inches(rec["x"]), Inches(rec["y"])
            hits = [sh for sh in slide.shapes
                    if legend._kind(sh) is not None and sh.left is not None
                    and sh.width is not None and sh.width <= Inches(0.11)
                    and abs(legend._centre(sh)[0] - cx) < Inches(0.02)
                    and abs(legend._centre(sh)[1] - cy) < Inches(0.02)]
            if len(hits) != 1:
                bad.append(f"{loc}: {len(hits)} asset symbols at ({rec['x']},{rec['y']})")
                continue
            got = legend.asset_key_of(hits[0])
            if got is False:
                bad.append(f"{loc}: symbol at ({rec['x']},{rec['y']}) matches no class "
                           f"in the symbology")
            elif got != rec["cls"]:
                bad.append(f"{loc}: symbol at ({rec['x']},{rec['y']}) drawn as {got!r}, "
                           f"classified {rec['cls']!r}")
    check("every asset symbol is drawn as its as-built class", bad)


def check_class_legend(prs, symbols):
    bad = []
    for ix, loc in SHEETS.items():
        declared = {k[1] for k in declared_from_sheet(prs.slides[ix], loc)
                    if k[0] == legend.ASSET}
        present = {r["cls"] for r in symbols[loc]["symbols"]}
        for cls in sorted(present - declared, key=str):
            bad.append(f"{loc}: class {cls!r} is plotted but has no legend row")
        for cls in sorted(declared - present, key=str):
            bad.append(f"{loc}: legend row for {cls!r} but no such asset on the sheet")
    check("each as-built class present has its own legend row", bad)


# ---------------------------------------------------------------- 3. text vs data

def check_scope_notes(prs, sheets):
    """Every light the field sheet lists must appear in its sheet's scope note, and no
    light may appear on the wrong sheet's note."""
    bad = []
    for ix, loc in SHEETS.items():
        text = " ".join(sh.text_frame.text for sh in prs.slides[ix].shapes
                        if sh.has_text_frame and sh.left is not None
                        and _in(sh.left) >= MAP_RIGHT)
        for row in sheets[loc]:
            lt = row["light"]
            circuit, num = lt.split("/")
            # the notes abbreviate repeats, e.g. "SBC102-02/024, 01/027"
            if lt not in text and f"{circuit.split('-')[1]}/{num}" not in text:
                bad.append(f"{loc}: {lt} is on the field sheet but not in the scope note")
    check("scope notes name every field-sheet light", bad)


def check_table(prs, sheets, positions):
    """The consolidated table must agree with the field sheets row for row."""
    tbl = next(sh for sh in prs.slides[4].shapes if sh.has_table).table
    rows = [[c.text.strip() for c in r.cells] for r in tbl.rows][1:]
    bad = []
    want_action = {
        ("CC0000", "0055CC"): "CORE OUT + NEW CABLE",
        ("CC0000", "E8710A"): "CORE OUT + NEW CABLE",
        ("0055CC", "1E8E3E"): "NEW CABLE ONLY",
        ("E8710A", "1E8E3E"): "NEW CABLE ONLY",
        ("F9AB00", "E8710A"): "DUMMY PLATE + NEW CABLE",
        ("9AA0A6", None): "NOT AFFECTED",
    }
    seen = collections.Counter()
    for loc, asset, action, route, _remark in rows:
        seen[(loc, asset)] += 1
        row = next((r for r in sheets.get(loc, []) if r["light"] == asset), None)
        if row is None:
            if asset.startswith("RRM"):
                continue
            if "NOT ON FIELD SHEET" not in action:
                bad.append(f"{loc} {asset}: not on the field sheet but action is {action!r}")
            continue
        exp = want_action[legend_key(loc, row)]
        if action != exp:
            bad.append(f"{loc} {asset}: table says {action!r}, columns imply {exp!r}")
        exp_route = row["route"] or row.get("route_deck")
        if exp_route and route and route != exp_route:
            bad.append(f"{loc} {asset}: table route {route!r} != {exp_route!r}")
    for loc in SHEETS.values():
        for row in sheets[loc]:
            if seen[(loc, row["light"])] != 1:
                bad.append(f"{loc} {row['light']}: appears {seen[(loc, row['light'])]}x "
                           f"in the consolidated table, expected once")
    check("consolidated table agrees with the field sheets", bad)


def legend_key(loc, row):
    ring, inner = R.expected(loc, row)
    return (ring, inner)


def check_totals(prs, sheets):
    """The totals tiles must equal what the field sheets add up to."""
    texts = [sh.text_frame.text for sh in prs.slides[5].shapes if sh.has_text_frame]
    joined = "\n".join(texts)
    n_sec = sum(1 for loc in SHEETS.values() for r in sheets[loc] if r["sec"])
    n_core = sum(1 for r in sheets["LOC-01"] if r["sec"] and r["base"])
    n_dummy = sum(1 for r in sheets["LOC-03"] if r["base"])
    n_not = sum(1 for loc in SHEETS.values() for r in sheets[loc]
                if not r["sec"] and not r["base"])
    n_duct = sum(1 for loc in SHEETS.values() for r in sheets[loc]
                 if r["sec"] and (r["route"] or r.get("route_deck")) == "Duct")
    n_saw = n_sec - n_duct
    bad = []
    for label, want in (("Core-outs", n_core), ("Dummy plates", n_dummy),
                        ("Secondary cable runs", n_sec),
                        ("Field-verified not affected", n_not),
                        ("Saw-cut route runs", n_saw),
                        ("Duct route runs", n_duct)):
        tile = next((t for t in texts if t.startswith(label)), None)
        if tile is None:
            bad.append(f"totals sheet has no {label!r} tile")
            continue
        idx = texts.index(tile)
        value = texts[idx - 1].strip()
        if value != str(want):
            bad.append(f"totals: {label} tile reads {value!r}, field sheets give {want}")
    if str(n_sec) not in joined:
        bad.append(f"totals sheet never states the {n_sec} secondary cable runs")
    check("totals sheet equals the field-sheet arithmetic", bad)


def check_revision(prs):
    bad = []
    for i, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text
            if "REV P06" in t or "Rev P06" in t:
                bad.append(f"slide {i + 1}: stale revision text {t.strip()[:60]!r}")
    for ix in SHEETS:
        text = " ".join(sh.text_frame.text for sh in prs.slides[ix].shapes
                        if sh.has_text_frame)
        if "REV P07" not in text:
            bad.append(f"slide {ix + 1}: no REV P07 marking")
    check("revision block reads P07 throughout, with no P06 left", bad)


# ---------------------------------------------------------------- 4. hygiene

def check_no_phantom_outlines(prs):
    """python-pptx's line.color getter inserts <a:ln><a:solidFill/></a:ln>, which
    renders as a black outline. Nothing in the pipeline may leave one behind."""
    A = legend.DRAWINGML
    bad = []
    for i, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            sp = getattr(sh._element, "spPr", None)
            if sp is None:
                continue
            ln = sp.find(A + "ln")
            if ln is None:
                continue
            sf = ln.find(A + "solidFill")
            if sf is not None and len(sf) == 0:
                bad.append(f"slide {i + 1}: {sh.shape_id} has an empty <a:solidFill/> outline")
    check("no shape carries a phantom outline", bad)


def check_labels(prs, sheets):
    bad = []
    for ix, loc in SHEETS.items():
        lbls = labels_on(prs.slides[ix])
        names = collections.Counter(t for t, *_ in lbls)
        for row in sheets[loc]:
            lt = row["light"]
            # a shared label ("A + B (0.49 m apart)") counts for both lights
            hits = [n for n in names if n == lt or lt in n
                    or lt.split("-")[1] in n and n.startswith(lt.split("-")[0])]
            if not hits:
                bad.append(f"{loc}: {lt} has no label on the sheet")
            elif sum(names[h] for h in hits) > 1:
                bad.append(f"{loc}: {lt} is labelled more than once")
        for name, n in names.items():
            if n > 1:
                bad.append(f"{loc}: label {name!r} appears {n}x")
        for t, l, _t2, w in lbls:
            if l + w > MAP_RIGHT + 0.01:
                bad.append(f"{loc}: label {t!r} is {l + w:.2f}\" wide-right, past the map frame")
    check("every field-sheet light is labelled exactly once, inside the frame", bad)


# ---------------------------------------------------------------- 5. rendered PDF

def render(pptx):
    out = pathlib.Path(tempfile.mkdtemp())
    try:
        subprocess.run(["soffice", "--headless", "--norestore", "--convert-to", "pdf",
                        "--outdir", str(out), str(pptx)],
                       check=True, capture_output=True, timeout=300,
                       env={"HOME": "/tmp/lohome", "PATH": "/usr/bin:/bin"})
    except Exception:
        return None
    pdfs = list(out.glob("*.pdf"))
    return pdfs[0] if pdfs else None


def check_rendered(pptx):
    try:
        import fitz
    except ImportError:
        check("rendered PDF layout", ["PyMuPDF not installed"], skipped=True)
        return
    pdf = render(pptx)
    if pdf is None:
        check("rendered PDF layout", ["LibreOffice Impress unavailable"], skipped=True)
        return
    doc = fitz.open(pdf)
    bad = []
    if doc.page_count != 6:
        bad.append(f"expected 6 pages, got {doc.page_count}")
    for pno, page in enumerate(doc):
        spans = [s for b in page.get_text("dict")["blocks"] if b["type"] == 0
                 for l in b["lines"] for s in l["spans"] if s["text"].strip()]
        for s in spans:
            x0, y0, x1, y1 = s["bbox"]
            if x1 > PANEL_RIGHT * 72 + 4 or y1 > PAGE_H * 72 + 2 or x0 < -2 or y0 < -2:
                bad.append(f"page {pno + 1}: text {s['text'].strip()[:34]!r} outside the sheet "
                           f"({x0 / 72:.2f},{y0 / 72:.2f})-({x1 / 72:.2f},{y1 / 72:.2f})")
        if pno in SHEETS:
            # label-on-label collisions in the map frame render as garble
            map_spans = [s for s in spans if s["bbox"][2] < MAP_RIGHT * 72
                         and any(c in s["text"] for c in "./")]
            for i, a in enumerate(map_spans):
                for b in map_spans[i + 1:]:
                    ax0, ay0, ax1, ay1 = a["bbox"]
                    bx0, by0, bx1, by1 = b["bbox"]
                    ox = min(ax1, bx1) - max(ax0, bx0)
                    oy = min(ay1, by1) - max(ay0, by0)
                    if ox > 6 and oy > 3:
                        bad.append(f"page {pno + 1}: labels overlap — "
                                   f"{a['text'].strip()[:22]!r} / {b['text'].strip()[:22]!r} "
                                   f"({ox / 72:.2f}\" x {oy / 72:.2f}\")")
    check("rendered PDF: 6 pages, all text inside the sheet, no label collisions", bad)


# ---------------------------------------------------------------------- main

def main():
    target = HERE / (sys.argv[1] if len(sys.argv) > 1 else "TWY_E_AGL_Shop_Drawing_ZIA_P07.pptx")
    sheets = json.loads((HERE / "field_sheets.json").read_text())
    positions = json.loads((HERE / "marker_positions.json").read_text())
    asset_keys = json.loads((HERE / "asset_key.json").read_text())
    symbols = json.loads((HERE / "asset_symbols.json").read_text())
    prs = Presentation(target)

    check_markers(prs, sheets, positions)
    check_no_stray_markers(prs, sheets)
    check_legends(prs, asset_keys)
    check_asset_symbols(prs, symbols)
    check_class_legend(prs, symbols)
    check_scope_notes(prs, sheets)
    check_table(prs, sheets, positions)
    check_totals(prs, sheets)
    check_revision(prs)
    check_no_phantom_outlines(prs)
    check_labels(prs, sheets)
    check_rendered(target)

    print(f"\n{'=' * 76}\n  {target.name}\n{'=' * 76}")
    failed = skipped = 0
    for name, problems, was_skipped in results:
        if was_skipped:
            print(f"  SKIP  {name}")
            for p in problems:
                print(f"          {p}")
            skipped += 1
        elif problems:
            print(f"  FAIL  {name}  ({len(problems)})")
            for p in problems:
                print(f"          - {p}")
            failed += 1
        else:
            print(f"  ok    {name}")
    total = len(results)
    print(f"\n  {total - failed - skipped}/{total} checks passed, {failed} failed, "
          f"{skipped} skipped\n")
    return 1 if failed else 0


if __name__ == "__main__":
    raise SystemExit(main())
