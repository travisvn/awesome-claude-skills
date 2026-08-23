#!/usr/bin/env python3
"""Rebuild each location sheet's legend so it accounts for every symbol plotted on
that sheet, and add an as-built asset key for the label prefixes.

Rev P06's legends were incomplete. The AGL feed manhole / handhole square is drawn on
all three sheets but P05's legend row for it was dropped, and P06 reused the freed row
to relabel the grey ring as "Dummy Plate" — see reconcile_fieldsheet.py. The milling /
cut area shading was never legended on any sheet. And every existing asset — centreline
light, stop bar light, edge light, sign foundation, handhole, manhole, civil pit, RRM —
plots as one undifferentiated teal dot, identifiable only from its label prefix, which
nothing on the sheet explained.

Rather than invent new symbology on an issued drawing, the existing-asset dot keeps one
legend row and the sheet gains an ASSET KEY (AS-BUILT) block naming each label prefix
present. Prefix -> asset class is not guessed: it is read from the classifier in the
zia-mta-basemap skill, by matching each plotted symbol to its nearest fixture in
Z1-Z2-Z3-MTA_SEGMENTATION.dxf. `python asset_key.py` regenerates asset_key.json and
records the match quality; AGL light prefixes resolve unanimously at 0.2-1.2 m, civil
prefixes at ~2.0-2.2 m, which is the documented civil-symbol plot offset.

`build()` lays the legend out in two columns and then audits itself both ways: every
point symbol plotted in the map frame must have a legend row, and every legend row's
swatch must correspond to something plotted.
"""

import copy
import json
import pathlib

from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

HERE = pathlib.Path(__file__).parent

GREY = "9AA0A6"
RED = "CC0000"
BLUE = "0055CC"
ORANGE = "E8710A"
AMBER = "F9AB00"
GREEN = "1E8E3E"
TEAL = "12A5B8"
MAGENTA = "BB00BB"
WHITE = "FFFFFF"
BLACK = "000000"
SHADE = "BFBFBF"
DEMARC = "CC0000"
CENTRELINE = "C7A400"
STOPBAR = "8B0000"
EDGE = "7A8288"
INDIC_DUCT = "C9CDD2"


# --- existing-asset symbology, one per as-built class -----------------------
# Shape carries the family, colour carries the AGL optical colour where the asset has
# one (green centreline, red stop bar, blue edge, yellow guard light). Works action is
# conveyed only by the ring drawn over the symbol, so a green asset dot inside a red
# ring stays unambiguous: green = centreline light, red ring = core out.
#   class -> (prst, size_in, fill or None, line, line_pt)
ASSET_SPEC = {
    "Taxiway centreline light":      ("ellipse",  0.075, "00A650", "00A650", 0.4),
    "Stop bar light":                ("ellipse",  0.075, "D32F2F", "D32F2F", 0.4),
    "Taxiway edge light":            ("ellipse",  0.075, "1565C0", "1565C0", 0.4),
    "Runway guard light / RRM":      ("ellipse",  0.075, "F2C200", "9A7B00", 0.5),
    "Existing light base":           ("ellipse",  0.080, None,     "7A8288", 0.75),
    "Sign foundation":               ("rect",     0.090, "5F6368", "5F6368", 0.4),
    "Handhole":                      ("rect",     0.075, None,     "12A5B8", 0.75),
    "Manhole":                       ("rect",     0.085, "455A64", "455A64", 0.4),
    "Existing manhole":              ("rect",     0.085, None,     "455A64", 0.75),
    "Existing transformer handhole": ("diamond",  0.095, "12A5B8", "0B7C8A", 0.4),
    "Earthing pit":                  ("diamond",  0.095, None,     "827717", 0.75),
    "Earthing point":                ("triangle", 0.090, "827717", "827717", 0.4),
}
GENERIC_ASSET = ("ellipse", 0.064, "12A5B8", "12A5B8", 0.5)   # unclassified fallback
ASSET = "asset"

# --- layout ----------------------------------------------------------------
COL_X = (Inches(11.27), Inches(13.62))     # swatch left, per column
TEXT_X = (Inches(11.52), Inches(13.87))    # text left, per column
TEXT_W = (Inches(2.05), Inches(2.28))      # text width, per column
FONT_PT = (6.4, 6.0)                       # row font size, per column
# average glyph width as a fraction of the point size: column A is all caps
# and measures wider than column B's mixed-case asset classes.
CHAR_W = (0.58, 0.48)
ROW_H = Inches(0.195)
RING_D = Inches(0.21)
INNER_D = Inches(0.11)
DOT_D = Inches(0.10)
SQ_D = Inches(0.17)
LINE_W = Inches(0.26)

# --- what the map actually plots -------------------------------------------
# key -> (legend kind, args). Detected styles are matched against these keys.
MARKER = "marker"; RING = "ring"; DOT = "dot"; SQUARE = "square"; FILL = "fill"
LINE = "line"; HEAD = "head"; TEXT = "text"

ROWS = {
    "LOC-01": [
        (HEAD, "WORKS ACTIONS"),
        (MARKER, RED, BLUE, "CORE OUT + NEW CABLE (DUCT)"),
        (MARKER, RED, ORANGE, "CORE OUT + NEW CABLE (SAWCUT)"),
        (MARKER, BLUE, GREEN, "NEW SEC. CABLE ONLY (DUCT)"),
        (MARKER, ORANGE, GREEN, "NEW SEC. CABLE ONLY (SAWCUT)"),
        (RING, GREY, "NOT AFFECTED / NOT ON FIELD SHEET"),
        (HEAD, "AREAS & LINEWORK"),
        (FILL, SHADE, "MILLING / CUT AREA"),
        (LINE, DEMARC, "AGL WORKS DEMARCATION (GOVERNING)"),
        (LINE, MAGENTA, "SECONDARY DUCT CROSSING CUT (INDIC.)"),
        (LINE, INDIC_DUCT, "INDICATIVE DUCT — LIGHT TO MH / HH"),
        (LINE, CENTRELINE, "TAXIWAY CENTRELINE (THROUGH TCC)"),
        (LINE, STOPBAR, "STOP BAR (THROUGH SBC)"),
        (LINE, EDGE, "TWY EDGE / PAVEMENT (INDIC. 23 m)"),
    ],
    "LOC-02": [
        (HEAD, "WORKS ACTIONS"),
        (MARKER, BLUE, GREEN, "NEW SEC. CABLE ONLY (DUCT)"),
        (MARKER, MAGENTA, WHITE, "RRM — REMOVE / PROTECT / RE-FIX"),
        (HEAD, "AREAS & LINEWORK"),
        (FILL, SHADE, "MILLING / CUT AREA"),
        (LINE, DEMARC, "AGL WORKS DEMARCATION (GOVERNING)"),
        (LINE, MAGENTA, "SECONDARY DUCT CROSSING CUT (INDIC.)"),
        (LINE, INDIC_DUCT, "INDICATIVE DUCT — LIGHT TO MH / HH"),
        (LINE, CENTRELINE, "TAXIWAY CENTRELINE (THROUGH TCC)"),
        (LINE, EDGE, "TWY EDGE / PAVEMENT (INDIC. 23 m)"),
    ],
    "LOC-03": [
        (HEAD, "WORKS ACTIONS"),
        (MARKER, AMBER, ORANGE, "BASE RETAINED + DUMMY PLATE (SAWCUT)"),
        (MARKER, BLUE, GREEN, "NEW SEC. CABLE ONLY (DUCT)"),
        (MARKER, MAGENTA, WHITE, "RRM — REMOVE / PROTECT / RE-FIX"),
        (SQUARE, BLACK, "FITTING LIFTED + OPENING PROTECTED"),
        (HEAD, "AREAS & LINEWORK"),
        (FILL, SHADE, "MILLING / CUT AREA"),
        (LINE, DEMARC, "AGL WORKS DEMARCATION (GOVERNING)"),
        (LINE, MAGENTA, "SECONDARY DUCT CROSSING CUT (INDIC.)"),
        (LINE, INDIC_DUCT, "INDICATIVE DUCT — LIGHT TO MH / HH"),
        (LINE, CENTRELINE, "TAXIWAY CENTRELINE (THROUGH TCC)"),
        (LINE, STOPBAR, "STOP BAR (THROUGH SBC)"),
        (LINE, EDGE, "TWY EDGE / PAVEMENT (INDIC. 23 m)"),
    ],
}

PANEL = {                      # legend panel rect: (left, top, width, height)
    "LOC-01": (11.05, 7.707, 5.15, 3.74),
    "LOC-02": (11.05, 6.695, 5.15, 2.84),
    "LOC-03": (11.05, 6.948, 5.15, 3.44),
}
HEADER_TEXT = "LEGEND — THIS SHEET"
SHEET_BOTTOM = Inches(11.60)   # keep clear of the sheet edge at 11.689
ROW_TOP = Inches(0.32)         # first row, below the panel header
PAD = Inches(0.10)


# --- shape helpers ----------------------------------------------------------

DRAWINGML = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

def _kind(sh):
    try:
        return str(sh.auto_shape_type).split(" ")[0]
    except (ValueError, AttributeError):
        return None            # freeform / polyline segment, not a point symbol


# python-pptx's shape.fill / shape.line accessors are NOT read-only: asking a shape
# with no explicit <a:ln> for line.color inserts <a:ln><a:solidFill/></a:ln>, which
# renders as a black outline. An audit must never alter the drawing it inspects, so
# these read the XML directly.

def _srgb(parent, tag):
    if parent is None:
        return None
    el = parent.find(DRAWINGML + tag)
    if el is None:
        return None
    clr = el.find(DRAWINGML + "srgbClr")
    return None if clr is None else clr.get("val")


def _fill(sh):
    return _srgb(sh._element.spPr, "solidFill")


def _line(sh):
    ln = sh._element.spPr.find(DRAWINGML + "ln")
    return _srgb(ln, "solidFill")


def _dash(sh):
    """The stroke's prstDash value, or None for a solid line."""
    ln = sh._element.spPr.find(DRAWINGML + "ln")
    if ln is None:
        return None
    d = ln.find(DRAWINGML + "prstDash")
    return None if d is None else d.get("val")


def _set_dash(sh, val):
    ln = sh._element.spPr.find(DRAWINGML + "ln")
    d = ln.find(DRAWINGML + "prstDash")
    if val is None:
        if d is not None:
            ln.remove(d)
        return
    if d is None:
        d = ln.makeelement(DRAWINGML + "prstDash", {})
        ln.append(d)
    d.set("val", val)


def _prst(sh):
    g = sh._element.spPr.find(DRAWINGML + "prstGeom")
    return None if g is None else g.get("prst")


def _set_prst(sh, prst):
    g = sh._element.spPr.find(DRAWINGML + "prstGeom")
    if g is not None:
        g.set("prst", prst)


def _centre(sh):
    return (sh.left + sh.width // 2, sh.top + sh.height // 2)


def asset_key_of(sh):
    """The ASSET_SPEC key a plotted existing-asset symbol matches, or None."""
    got = (_prst(sh), round(_in(sh.width), 3), _fill(sh), _line(sh))
    for cls, (prst, size, fill, line, _w) in ASSET_SPEC.items():
        if got == (prst, round(size, 3), fill, line):
            return cls
    if got == (GENERIC_ASSET[0], round(GENERIC_ASSET[1], 3), GENERIC_ASSET[2],
               GENERIC_ASSET[3]):
        return None          # generic fallback dot, legended separately
    return False             # not an asset symbol at all


def asset_class_of_swatch(sh):
    """Match a legend swatch to its class on (prst, fill, line) — swatches are scaled
    up from the map symbol, so size is deliberately ignored. The three-tuple is unique
    across ASSET_SPEC; asserted by check_all.py."""
    got = (_prst(sh), _fill(sh), _line(sh))
    for cls, (prst, _size, fill, line, _w) in ASSET_SPEC.items():
        if got == (prst, fill, line):
            return cls
    if got == (GENERIC_ASSET[0], GENERIC_ASSET[2], GENERIC_ASSET[3]):
        return None
    return False


def draw_asset(sh, cls):
    """Restyle a plotted asset symbol as its as-built class, about its own centre."""
    prst, size, fill, line, line_pt = ASSET_SPEC.get(cls, GENERIC_ASSET)
    cx, cy = _centre(sh)
    d = Inches(size)
    sh.width = sh.height = d
    sh.left, sh.top = cx - d // 2, cy - d // 2
    _set_prst(sh, prst)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor.from_string(fill)
    sh.line.color.rgb = RGBColor.from_string(line)
    sh.line.width = Pt(line_pt)


def _in(v):
    return Emu(v).inches


def _chars_per_line(col):
    return max(1, int(_in(TEXT_W[col]) * 72 / (FONT_PT[col] * CHAR_W[col])))


def plotted_symbols(slide, map_right=Inches(11.0)):
    """Every point symbol in the map frame, as {style_key: count}.

    A point symbol is a real autoshape (OVAL / RECTANGLE); the sheets' linework is
    freeform, for which auto_shape_type raises and _kind returns None. Label boxes
    carry text and are skipped.
    """
    out = {}
    ovals, others = [], []
    for sh in slide.shapes:
        if sh.left is None or sh.left >= map_right:
            continue
        k = _kind(sh)
        if k is None:
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            continue
        (ovals if k == "OVAL" else others).append(sh)

    rings = [o for o in ovals if abs(o.width - Inches(0.23)) < Inches(0.01)]
    inners = [o for o in ovals if abs(o.width - Inches(0.12)) < Inches(0.01)]

    # existing-asset symbols: everything small, keyed by its ASSET_SPEC tuple
    for sh in ovals + others:
        if sh.width is None or sh.width > Inches(0.11):
            continue
        k = asset_key_of(sh)
        if k is False:
            out[("unknown", _prst(sh), round(_in(sh.width), 3), _fill(sh), _line(sh))] = \
                out.get(("unknown", _prst(sh), round(_in(sh.width), 3), _fill(sh),
                         _line(sh)), 0) + 1
        else:
            key = (ASSET, k)
            out[key] = out.get(key, 0) + 1

    for r in rings:
        rc = _centre(r)
        match = [i for i in inners
                 if abs(_centre(i)[0] - rc[0]) < Inches(0.02)
                 and abs(_centre(i)[1] - rc[1]) < Inches(0.02)]
        if match:
            key = (MARKER, _line(r), _fill(match[0]))
        else:
            key = (RING, _line(r))
        out[key] = out.get(key, 0) + 1
    for sh in others:
        if sh.width is not None and sh.width <= Inches(0.11):
            continue          # already counted as an asset symbol
        if _fill(sh) == SHADE:
            out[(FILL, SHADE)] = out.get((FILL, SHADE), 0) + 1
        elif _kind(sh) == "RECTANGLE" and _fill(sh) is None and _line(sh):
            key = (SQUARE, _line(sh))
            out[key] = out.get(key, 0) + 1
    return out


def plotted_linework(slide, map_right=Inches(11.0)):
    """Line colours used by the sheet's linework, as {colour: count}."""
    seen = {}
    for sh in slide.shapes:
        if sh.left is None or sh.left >= map_right or _kind(sh) is not None:
            continue
        c = _line(sh)
        if c:
            seen[c] = seen.get(c, 0) + 1
    return seen


def linework_dashes(slide, map_right=Inches(11.0)):
    """The dash style each line colour is actually drawn with, so the legend swatch
    matches the line rather than whichever template it was cloned from."""
    votes = {}
    for sh in slide.shapes:
        if sh.left is None or sh.left >= map_right or _kind(sh) is not None:
            continue
        c = _line(sh)
        if c:
            votes.setdefault(c, {}).setdefault(_dash(sh), 0)
            votes[c][_dash(sh)] += 1
    return {c: max(v, key=v.get) for c, v in votes.items()}


# --- builders ---------------------------------------------------------------

class Templates:
    """Detached XML snapshots of the shapes the deck already uses, taken before the
    panel is cleared, so new rows inherit the deck's styling even though the originals
    are deleted."""

    def __init__(self, slide):
        ovals = [s for s in slide.shapes
                 if _kind(s) == "OVAL" and s.left is not None and s.left >= Inches(11)]
        ring = next(o for o in ovals if abs(o.width - RING_D) < Inches(0.01))
        inner = next(o for o in ovals if abs(o.width - INNER_D) < Inches(0.01))
        dot = next(o for o in ovals if abs(o.width - DOT_D) < Inches(0.005))
        line = next(s for s in slide.shapes
                    if s.left is not None and s.left >= Inches(11)
                    and s.height is not None and s.height < Inches(0.01)
                    and s.width is not None and abs(s.width - Inches(0.3)) < Inches(0.02))
        text = next(s for s in slide.shapes
                    if s.has_text_frame and s.left is not None
                    and s.left >= Inches(11.5)
                    and s.text_frame.text.strip().startswith(("CORE OUT", "NEW SEC.",
                                                              "BASE RETAINED")))
        self.ring = copy.deepcopy(ring._element)
        self.inner = copy.deepcopy(inner._element)
        self.dot = copy.deepcopy(dot._element)
        self.line = copy.deepcopy(line._element)
        self.text = copy.deepcopy(text._element)
        self.head = next(s for s in slide.shapes
                         if s.has_text_frame and "LEGEND" in s.text_frame.text)
        self.square = None      # set by build(), may be snapshotted from another slide


def _clone(slide, template_el, left, top, width=None, height=None):
    el = copy.deepcopy(template_el)
    slide.shapes._spTree.append(el)
    sh = next(s for s in slide.shapes if s._element is el)
    ids = [s.shape_id for s in slide.shapes if s._element is not el]
    sh._element.nvSpPr.cNvPr.set("id", str(max(ids) + 1))
    sh._element.nvSpPr.cNvPr.set("name", f"Legend {max(ids) + 1}")
    sh.left, sh.top = left, top
    if width is not None:
        sh.width = width
    if height is not None:
        sh.height = height
    return sh


def _retext(shape, text, size=None, bold=None, rgb=None):
    tf = shape.text_frame
    proto = tf.paragraphs[0].runs[0]
    f = dict(size=proto.font.size, name=proto.font.name, bold=proto.font.bold,
             rgb=proto.font.color.rgb if proto.font.color and proto.font.color.type is not None
             else None)
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(size) if size else f["size"]
    r.font.name = f["name"]
    r.font.bold = f["bold"] if bold is None else bold
    r.font.color.rgb = RGBColor.from_string(rgb) if rgb else (f["rgb"] or RGBColor(0x1F, 0x29, 0x37))
    return shape


def _ring(sh, colour, width_pt=1.5):
    sh.fill.background()
    sh.line.color.rgb = RGBColor.from_string(colour)
    sh.line.width = Pt(width_pt)


def _disc(sh, colour):
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string(colour)
    sh.line.color.rgb = RGBColor.from_string(colour)
    sh.line.width = Pt(0.5)


def build(slide, loc, asset_counts, square_template=None):
    """Clear the legend panel and lay it out afresh in two columns.

    Returns the list of style keys the legend declares, for the caller to audit.
    """
    left, top, width, height = (Inches(v) for v in PANEL[loc])
    dashes = linework_dashes(slide)
    tpl = Templates(slide)
    tpl.square = square_template

    # keep the panel rect and the header; drop every other shape inside the panel
    keep = {id(tpl.head)}
    panel_rect = None
    for sh in list(slide.shapes):
        if sh.left is None or sh.width is None:
            continue
        if sh.left < Inches(11.0):
            continue
        if sh.top is None or sh.top < top - Inches(0.02):
            continue
        if abs(sh.width - width) < Inches(0.02) and abs(sh.height - height) < Inches(0.02):
            panel_rect = sh
            continue
        if id(sh) in keep:
            continue
        sh._element.getparent().remove(sh._element)
    assert panel_rect is not None, f"{loc}: legend panel rect not found"

    _retext(tpl.head, HEADER_TEXT)
    tpl.head.left, tpl.head.top = left + Inches(0.15), top + Inches(0.08)

    rows = list(ROWS[loc])
    # one row per as-built asset class actually plotted on this sheet, as the as-built
    # does — ordered AGL lights first, then civil, each by descending count
    counts = asset_counts
    order = [c for c in ASSET_SPEC if c in counts]
    rows_b = [(HEAD, "EXISTING ASSETS (AS-BUILT)")]
    rows_b += [(ASSET, c, f"{c}  ×{counts[c]}") for c in order]
    if counts.get(None):
        rows_b.append((ASSET, None, f"Other existing asset  ×{counts[None]}"))
    rows_b.append((SQUARE, MAGENTA, "AGL feed manhole / handhole"))

    def units(items, col):
        """Row slots the column needs, counting wrapped text as more than one."""
        n = 0
        for r in items:
            txt = r[-1]
            cpl = _chars_per_line(col)
            n += max(1, -(-len(txt) // cpl))
        return n

    # grow the panel if the rows need more room than Rev P06's box had
    needed = ROW_TOP + max(units(rows, 0), units(rows_b, 1)) * ROW_H + PAD
    room = SHEET_BOTTOM - top
    assert needed <= room, (f"{loc}: legend needs {_in(needed):.2f}\" but only "
                            f"{_in(room):.2f}\" is available above the sheet edge")
    if needed > height:
        height = needed
        panel_rect.height = height

    declared = []
    for col, items in ((0, rows), (1, rows_b)):
        y = top + ROW_TOP
        for row in items:
            kind = row[0]
            if kind == HEAD:
                _retext(_clone(slide, tpl.text, TEXT_X[col] - Inches(0.25), y,
                               width=TEXT_W[col] + Inches(0.25)), row[1], size=6.6,
                        bold=True, rgb="1E2761")
                y += ROW_H
                continue
            sy = y + (ROW_H - RING_D) // 2
            if kind == MARKER:
                _ring(_clone(slide, tpl.ring, COL_X[col], sy), row[1])
                _disc(_clone(slide, tpl.inner,
                             COL_X[col] + (RING_D - INNER_D) // 2,
                             sy + (RING_D - INNER_D) // 2), row[2])
                declared.append((MARKER, row[1], row[2]))
                text = row[3]
            elif kind == RING:
                _ring(_clone(slide, tpl.ring, COL_X[col], sy), row[1])
                declared.append((RING, row[1]))
                text = row[2]
            elif kind == ASSET:
                sw = _clone(slide, tpl.dot, COL_X[col], sy)
                prst, size, _f, _l, _w = ASSET_SPEC.get(row[1], GENERIC_ASSET)
                sw.width = sw.height = Inches(min(size * 1.7, 0.15))  # read larger,
                #                                            but clear of the row pitch
                sw.left = COL_X[col] + (RING_D - sw.width) // 2
                sw.top = sy + (RING_D - sw.height) // 2
                _set_prst(sw, prst)
                spec = ASSET_SPEC.get(row[1], GENERIC_ASSET)
                if spec[2] is None:
                    sw.fill.background()
                else:
                    sw.fill.solid()
                    sw.fill.fore_color.rgb = RGBColor.from_string(spec[2])
                sw.line.color.rgb = RGBColor.from_string(spec[3])
                sw.line.width = Pt(spec[4] * 1.6)
                declared.append((ASSET, row[1]))
                text = row[2]
            elif kind == DOT:
                _disc(_clone(slide, tpl.dot, COL_X[col] + (RING_D - DOT_D) // 2,
                             sy + (RING_D - DOT_D) // 2), row[1])
                declared.append((DOT, row[1]))
                text = row[2]
            elif kind == SQUARE:
                assert tpl.square is not None, f"{loc}: no square template"
                sq = _clone(slide, tpl.square, COL_X[col] + (RING_D - SQ_D) // 2,
                            sy + (RING_D - SQ_D) // 2, width=SQ_D, height=SQ_D)
                _ring(sq, row[1])
                declared.append((SQUARE, row[1]))
                text = row[2]
            elif kind == FILL:
                sq = _clone(slide, tpl.square, COL_X[col], sy + Inches(0.04),
                            width=Inches(0.22), height=Inches(0.13))
                sq.fill.solid()
                sq.fill.fore_color.rgb = RGBColor.from_string(row[1])
                sq.line.color.rgb = RGBColor.from_string(EDGE)
                sq.line.width = Pt(0.5)
                declared.append((FILL, row[1]))
                text = row[2]
            elif kind == LINE:
                sw = _clone(slide, tpl.line, COL_X[col] - Inches(0.02),
                            y + ROW_H // 2, width=LINE_W)
                sw.line.color.rgb = RGBColor.from_string(row[1])
                _set_dash(sw, dashes.get(row[1]))
                declared.append((LINE, row[1]))
                text = row[2]
            elif kind == TEXT:
                text = row[1]
            else:
                raise ValueError(kind)
            cpl = _chars_per_line(col)
            n_lines = max(1, -(-len(text) // cpl))
            _retext(_clone(slide, tpl.text, TEXT_X[col], y, width=TEXT_W[col],
                           height=Inches(0.13) * n_lines), text, size=FONT_PT[col])
            y += ROW_H * n_lines
        overflow = _in(y) - (_in(top) + _in(height))
        assert overflow < 0, f"{loc} column {col}: legend overflows panel by {overflow:.3f}\""
    return declared


def audit(slide, loc, declared):
    """Every plotted point symbol legended, and every legend swatch actually plotted."""
    plotted = plotted_symbols(slide)
    lines = plotted_linework(slide)
    problems = []
    dec = set(declared)
    for key, n in sorted(plotted.items(), key=lambda kv: str(kv[0])):
        if key not in dec:
            problems.append(f"{loc}: {n} x {key} plotted but not in the legend")
    for key in sorted(dec, key=str):
        if key[0] == LINE:
            if key[1] not in lines:
                problems.append(f"{loc}: legend row {key} has no matching linework")
        elif key not in plotted:
            problems.append(f"{loc}: legend row {key} is not plotted on the sheet")
    declared_lines = {k[1] for k in dec if k[0] == LINE}
    for colour, n in sorted(lines.items()):
        if colour not in declared_lines:
            problems.append(f"{loc}: {n} x {colour} linework plotted but not in the legend")
    return problems


def load_asset_counts(loc):
    """{class: count} for the sheet, from asset_symbols.json (None = unclassified)."""
    data = json.loads((HERE / "asset_symbols.json").read_text())[loc]
    counts = dict(data["counts"])
    n_generic = sum(1 for r in data["symbols"] if r["cls"] is None)
    if n_generic:
        counts[None] = n_generic
    return counts


def load_symbols(loc):
    return json.loads((HERE / "asset_symbols.json").read_text())[loc]["symbols"]
