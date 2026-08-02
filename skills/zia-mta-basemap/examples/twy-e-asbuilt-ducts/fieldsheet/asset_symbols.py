#!/usr/bin/env python3
"""Classify every existing-asset symbol on the three sheets, so each asset class can be
drawn with its own symbol and given its own legend row — as the as-built does.

Rev P07 up to now drew every existing asset as one teal dot, with a label-prefix key
explaining what each was. This resolves the class of *every* plotted dot, including the
majority that carry no label, by registering the sheet onto
Z1-Z2-Z3-MTA_SEGMENTATION.dxf and reading `asset_type` from the skill's classifier for
the nearest fixture.

Two things make that trustworthy rather than a guess:

  - **The systematic plot offset is removed first.** Civil symbols sit ~1.4-1.9 m from
    their surveyed insertion point and AGL symbols ~0.7-1.0 m (see ../README.md). The
    offset is estimated per sheet per family from the *labelled* symbols, whose class is
    known independently from the label prefix, then subtracted before matching.
  - **The result is cross-validated on the labelled symbols.** Where a symbol carries a
    label, the prefix gives its class independently; agreement between the two is
    reported and must exceed AGREEMENT_MIN.

Ambiguous symbols — a rival class nearly as close — are recorded with
`ambiguous: true` and fall back to the generic existing-asset dot rather than being
forced into a class.

Run from this directory:  python asset_symbols.py
"""

import collections
import json
import pathlib
import sys

import numpy as np

HERE = pathlib.Path(__file__).parent
ROOT = HERE.parent
sys.path.insert(0, "/root/.claude/skills/zia-mta-basemap/scripts")
sys.path.insert(0, str(ROOT / "pipeline"))

from basemap import load_fixtures          # noqa: E402
from register import apply                 # noqa: E402
from pptx import Presentation              # noqa: E402
from pptx.util import Emu                  # noqa: E402

DECK = HERE / "TWY_E_AGL_Shop_Drawing_ZIA_P07.pptx"
OUT = HERE / "asset_symbols.json"
SHEETS = {1: "LOC-01", 2: "LOC-02", 3: "LOC-03"}
EMU_PER_INCH = 914400.0
DOT_IN = 0.064
LABEL_DX, LABEL_DY = -0.13, 0.215
AMBIGUOUS_M = 0.5      # rival class this much closer than the winner -> ambiguous
MAX_MATCH_M = 3.0      # no fixture within this -> unclassified
AGREEMENT_MIN = 0.95   # labelled-symbol agreement between prefix and nearest fixture

# Label prefix -> as-built class. Where a symbol carries a label this is authoritative:
# it is the drawing's own assertion of what the asset is, and it outranks a proximity
# match. Nearest-fixture classification is used only for the unlabelled symbols, and the
# labelled ones measure how well it performs. Pit prefixes (P2/P4/P6/X_CV_STH_PITS) map
# to more than one civil class, so they are deliberately absent and fall through.
PREFIX_CLASS = {
    "SBC102": "Stop bar light",
    "TCC102": "Taxiway centreline light",
    "TCC103": "Taxiway centreline light",
    "TCCECH": "Taxiway centreline light",
    "TEC102": "Taxiway edge light",
    "SGC102": "Sign foundation",
    "RRM": "Runway guard light / RRM",
    "HH": "Handhole",
    "MH": "Manhole",
    "EL": "Existing light base",
}
AGL_PREFIXES = ("SBC", "TCC", "TEC", "SGC", "EL")


def kind(sh):
    try:
        return str(sh.auto_shape_type).split(" ")[0]
    except (ValueError, AttributeError):
        return None


def family(prefix):
    return "AGL" if prefix.startswith(AGL_PREFIXES) or prefix == "RRM" else "CIVIL"


def main():
    reg = json.loads((ROOT / "data" / "registration.json").read_text())
    prs = Presentation(DECK)
    out = {}
    agree = disagree = 0
    problems = []

    for slide_ix, loc in SHEETS.items():
        slide = prs.slides[slide_ix]
        T = dict(reg[loc])
        T["t"] = np.array(T["t"])

        def grid(x_in, y_in):
            return apply(T, [[x_in * EMU_PER_INCH, y_in * EMU_PER_INCH]])[0]

        dots, labels = [], []
        for sh in slide.shapes:
            if sh.left is None or Emu(sh.left).inches >= 11:
                continue
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text.strip()
                if len(t) < 42 and ("." in t or "/" in t):
                    labels.append((t, Emu(sh.left).inches, Emu(sh.top).inches))
            elif kind(sh) == "OVAL" and abs(Emu(sh.width).inches - DOT_IN) < 0.005:
                w = Emu(sh.width).inches
                dots.append((round(Emu(sh.left).inches + w / 2, 4),
                             round(Emu(sh.top).inches + w / 2, 4)))

        local = np.array([grid(x, y) for x, y in dots])
        fx = load_fixtures(bbox=(local[:, 0].min() - 40, local[:, 1].min() - 40,
                                 local[:, 0].max() + 40, local[:, 1].max() + 40),
                           assets_only=True).reset_index(drop=True)

        # --- estimate the systematic plot offset per family, from labelled symbols
        offsets = collections.defaultdict(list)
        for name, lx, ly in labels:
            g = grid(lx + LABEL_DX, ly + LABEL_DY)
            prefix = name.split(".")[0].split("-")[0]
            cls = PREFIX_CLASS.get(prefix)
            pool = fx[fx.asset_type == cls] if cls is not None else fx
            if not len(pool):
                continue
            d = np.hypot(pool.x - g[0], pool.y - g[1])
            i = d.idxmin()
            offsets[family(prefix)].append((g[0] - fx.x[i], g[1] - fx.y[i]))
        corr = {fam: np.median(np.array(v), axis=0) for fam, v in offsets.items() if v}

        # --- a symbol is "labelled" if a label sits at its marker offset
        labelled = {}
        for name, lx, ly in labels:
            want = (round(lx + LABEL_DX, 4), round(ly + LABEL_DY, 4))
            near = min(dots, key=lambda d: (d[0] - want[0]) ** 2 + (d[1] - want[1]) ** 2)
            if abs(near[0] - want[0]) <= 0.15 and abs(near[1] - want[1]) <= 0.15:
                labelled.setdefault(near, name)

        # --- classify every dot, with the offset removed
        recs = []
        for (dx, dy), p in zip(dots, local):
            best = None
            for fam, shift in corr.items():
                q = p - shift
                d = np.hypot(fx.x - q[0], fx.y - q[1]).values
                order = d.argsort()
                cls = fx.asset_type[order[0]]
                if family_of_class(cls) != fam:
                    continue
                rival = next((j for j in order[1:] if fx.asset_type[j] != cls), None)
                cand = dict(cls=cls, dist=float(d[order[0]]),
                            margin=float(d[rival] - d[order[0]]) if rival is not None else 99.0)
                if best is None or cand["dist"] < best["dist"]:
                    best = cand
            name = labelled.get((dx, dy))
            from_label = PREFIX_CLASS.get(name.split(".")[0].split("-")[0]) if name else None
            if from_label is not None:
                recs.append(dict(x=dx, y=dy, cls=from_label, source="label", label=name,
                                 nearest_cls=best["cls"] if best else None,
                                 match_m=round(best["dist"], 2) if best else None))
                continue
            if best is None or best["dist"] > MAX_MATCH_M:
                recs.append(dict(x=dx, y=dy, cls=None, unclassified=True,
                                 label=name, source="none"))
                problems.append(f"{loc}: unlabelled dot at ({dx},{dy}) has no fixture "
                                f"within {MAX_MATCH_M} m — left as a generic asset dot")
                continue
            recs.append(dict(x=dx, y=dy, cls=best["cls"], source="fixture",
                             match_m=round(best["dist"], 2),
                             margin_m=round(min(best["margin"], 99.0), 2),
                             ambiguous=best["margin"] < AMBIGUOUS_M))

        # --- cross-validate against the label prefixes
        by_pos = {(r["x"], r["y"]): r for r in recs}
        for name, lx, ly in labels:
            prefix = name.split(".")[0].split("-")[0]
            cls = PREFIX_CLASS.get(prefix)
            if cls is None:
                continue
            want = (round(lx + LABEL_DX, 4), round(ly + LABEL_DY, 4))
            near = min(by_pos, key=lambda k: (k[0] - want[0]) ** 2 + (k[1] - want[1]) ** 2)
            if abs(near[0] - want[0]) > 0.15 or abs(near[1] - want[1]) > 0.15:
                continue
            got = by_pos[near].get("nearest_cls")
            if got == cls:
                agree += 1
            else:
                disagree += 1
                problems.append(f"{loc}: {name} — label says {cls!r}, nearest fixture "
                                f"would say {got!r} (label wins)")

        counts = collections.Counter(r["cls"] for r in recs if r["cls"])
        out[loc] = dict(
            symbols=recs,
            counts=dict(counts),
            offset_correction_m={f: [round(float(v[0]), 3), round(float(v[1]), 3)]
                                 for f, v in corr.items()},
            ambiguous=sum(1 for r in recs if r.get("ambiguous")),
            unclassified=sum(1 for r in recs if r.get("unclassified")),
        )
        print(f"=== {loc}: {len(recs)} symbols, {len(counts)} classes, "
              f"{out[loc]['ambiguous']} ambiguous, {out[loc]['unclassified']} unclassified")
        for k, v in counts.most_common():
            print(f"     {k:34} {v}")

    total = agree + disagree
    rate = agree / total if total else 0.0
    print(f"\ncross-validation: on the {total} labelled symbols, nearest-fixture "
          f"classification agrees with the label prefix {agree}/{total} ({rate:.1%}) — "
          f"this is the accuracy of the method used for the unlabelled symbols")
    for p in problems:
        print("   -", p)
    if rate < AGREEMENT_MIN:
        raise SystemExit(f"agreement {rate:.1%} below {AGREEMENT_MIN:.0%} — not safe to "
                         f"redraw symbols from this classification")

    OUT.write_text(json.dumps(out, indent=1) + "\n")
    print("\nwrote", OUT)


CLASS_FAMILY = {
    "Taxiway centreline light": "AGL",
    "Stop bar light": "AGL",
    "Taxiway edge light": "AGL",
    "Runway guard light / RRM": "AGL",
    "Sign foundation": "AGL",
    "Existing light base": "AGL",
    "Handhole": "CIVIL",
    "Manhole": "CIVIL",
    "Existing manhole": "CIVIL",
    "Existing transformer handhole": "CIVIL",
    "Earthing pit": "CIVIL",
    "Earthing point": "CIVIL",
}


def family_of_class(cls):
    return CLASS_FAMILY.get(cls, "CIVIL")


if __name__ == "__main__":
    main()
