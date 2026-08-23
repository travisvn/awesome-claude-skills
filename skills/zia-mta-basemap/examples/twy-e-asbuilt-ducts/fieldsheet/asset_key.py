#!/usr/bin/env python3
"""Build asset_key.json — what each label prefix on the sheets actually is, per the
as-built classification.

Every existing asset on the three sheets plots as the same teal dot, so a reader can
only tell a stop bar light from a handhole by its label prefix, and nothing on the sheet
explained the prefixes. This resolves each prefix by registering the sheet onto
Z1-Z2-Z3-MTA_SEGMENTATION.dxf and reading `asset_type` from the skill's classifier for
the fixture nearest each labelled symbol.

Match quality is recorded and matters:

- AGL light prefixes (SBC / TCC / TCCECH / TEC / SGC) resolve unanimously at 0.2-1.2 m.
- Civil prefixes (HH / MH / RRM) resolve unanimously but at ~2.0-2.2 m. That is the
  documented civil-symbol plot offset (see ../README.md), not a misidentification.
- Pit prefixes (P2 / P4 / P6 / X_CV_STH_PITS) are genuinely mixed: the same prefix sits
  near "Existing manhole", "Existing transformer handhole" and "Earthing pit" fixtures in
  different places. Those are not forced to a single class — they are labelled with the
  generic civil-pit wording in PREFIX_OVERRIDE and flagged `unanimous: false`.

Run from this directory:  python asset_key.py
"""

import collections
import json
import pathlib
import sys

import numpy as np

HERE = pathlib.Path(__file__).parent
ROOT = HERE.parent
sys.path.insert(0, "/root/.claude/skills/zia-mta-basemap/scripts")
sys.path.insert(0, str(ROOT / "pipeline"))

from basemap import load_fixtures          # noqa: E402
from register import apply                 # noqa: E402
from pptx import Presentation              # noqa: E402
from pptx.util import Emu                  # noqa: E402

DECK = HERE / "TWY_E_AGL_Shop_Drawing_ZIA_P07.pptx"
EMU_PER_INCH = 914400.0
LABEL_DX, LABEL_DY = -0.13, 0.215

# Prefixes whose nearest-fixture class is not unanimous get explicit wording instead.
# The value is what goes on the drawing; the JSON keeps the observed classes too.
PREFIX_OVERRIDE = {
    "P2": "Civil pit (transformer handhole)",
    "P4": "Civil pit (manhole / transformer HH)",
    "P6": "Civil pit (manhole / earthing point)",
    "X_CV_STH_PITS &": "Civil pit / transformer HH",
    "EL": "Existing light base (EBASE)",
}
# Shown on the drawing in place of the raw prefix, where the raw form is unwieldy.
PREFIX_LABEL = {"X_CV_STH_PITS &": "X_CV_STH_PITS"}


def _kind(sh):
    try:
        return str(sh.auto_shape_type).split(" ")[0]
    except (ValueError, AttributeError):
        return None


def main():
    reg = json.loads((ROOT / "data" / "registration.json").read_text())
    prs = Presentation(DECK)
    out = {"order": {}}

    for slide_ix, loc in ((1, "LOC-01"), (2, "LOC-02"), (3, "LOC-03")):
        slide = prs.slides[slide_ix]
        T = dict(reg[loc])
        T["t"] = np.array(T["t"])

        def grid(x_in, y_in):
            return apply(T, [[x_in * EMU_PER_INCH, y_in * EMU_PER_INCH]])[0]

        labels, dots = [], []
        for sh in slide.shapes:
            if sh.left is None or Emu(sh.left).inches >= 11:
                continue
            if sh.has_text_frame and sh.text_frame.text.strip():
                txt = sh.text_frame.text.strip()
                if len(txt) < 26 and ("." in txt or "/" in txt):
                    labels.append((txt, Emu(sh.left).inches, Emu(sh.top).inches))
            elif _kind(sh) == "OVAL" and abs(Emu(sh.width).inches - 0.064) < 0.005:
                w = Emu(sh.width).inches
                dots.append((Emu(sh.left).inches + w / 2, Emu(sh.top).inches + w / 2))

        pts = [grid(x, y) for x, y in dots]
        xs = [p[0] for p in pts]
        ys = [p[1] for p in pts]
        fx = load_fixtures(bbox=(min(xs) - 40, min(ys) - 40, max(xs) + 40, max(ys) + 40),
                           assets_only=True)

        seen = collections.defaultdict(collections.Counter)
        dists = collections.defaultdict(list)
        first = {}
        for name, lx, ly in labels:
            g = grid(lx + LABEL_DX, ly + LABEL_DY)
            d = np.hypot(fx.x.values - g[0], fx.y.values - g[1])
            i = int(d.argmin())
            prefix = name.split(".")[0].split("-")[0]
            seen[prefix][fx.iloc[i].asset_type] += 1
            dists[prefix].append(float(d[i]))
            first.setdefault(prefix, name)

        rec = {}
        for prefix, types in seen.items():
            classes = types.most_common()
            unanimous = len(classes) == 1
            rec[prefix] = dict(
                asset_type=PREFIX_OVERRIDE.get(prefix, classes[0][0]),
                observed_classes={t: c for t, c in classes},
                unanimous=unanimous,
                count=sum(types.values()),
                median_match_m=round(float(np.median(dists[prefix])), 2),
                example=first[prefix],
                label=PREFIX_LABEL.get(prefix, prefix),
            )
            if not unanimous and prefix not in PREFIX_OVERRIDE:
                raise SystemExit(
                    f"{loc} {prefix}: classes not unanimous {classes} and no "
                    f"PREFIX_OVERRIDE — refusing to pick one for the drawing"
                )
        # AGL lights first, then civil, each alphabetically — reads like the as-built
        agl = sorted(p for p in rec if p.startswith(("SBC", "TCC", "TEC", "SGC", "EL")))
        civil = sorted(p for p in rec if p not in agl)
        out["order"][loc] = agl + civil
        out[loc] = {PREFIX_LABEL.get(p, p): rec[p] for p in rec}
        out["order"][loc] = [PREFIX_LABEL.get(p, p) for p in out["order"][loc]]

        print(f"=== {loc}")
        for p in out["order"][loc]:
            a = out[loc][p]
            mark = "" if a["unanimous"] else "   (mixed classes, generic wording used)"
            print(f"   {p:15} x{a['count']:2}  {a['median_match_m']:5.2f} m  "
                  f"{a['asset_type']}{mark}")

    (HERE / "asset_key.json").write_text(json.dumps(out, indent=1) + "\n")
    print("\nwrote", HERE / "asset_key.json")


if __name__ == "__main__":
    main()
