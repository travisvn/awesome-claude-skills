#!/usr/bin/env python3
"""Build marker_positions.json — the plotted position of every field-sheet light,
taken from the Rev P05 works overlay and verified against the source DXF.

Why P05 and not P06: Rev P06 re-plotted several LOC-01 symbols and its added
rings sit 0.15-1.16 m from the nearest surveyed light, whereas every P05 works
ring lands on a fixture insertion point to <=0.7 mm. In the
TCCECH-03/035 // 03/018 cluster the two nearest fixtures are 0.49 m apart and
P06's rings are near-equidistant from both, so its positions cannot be repaired
by snapping — P05's have to be used.

LOC-03 carries a uniform ~34-52 mm residual across all 12 lights. That is the
known bias in LOC-03's registration (fit RMS 58.2 mm, see ../README.md), not a
drawing error: a systematic offset shared by every marker on the sheet.

Run from the example root (the directory holding data/ and input/):
    python fieldsheet/build_marker_positions.py
"""

import json
import pathlib
import sys

import numpy as np

HERE = pathlib.Path(__file__).parent
ROOT = HERE.parent
SKILL_SCRIPTS = pathlib.Path("/root/.claude/skills/zia-mta-basemap/scripts")

sys.path.insert(0, str(SKILL_SCRIPTS))
sys.path.insert(0, str(ROOT / "pipeline"))

from basemap import load_fixtures          # noqa: E402
from register import apply                 # noqa: E402
from pptx import Presentation              # noqa: E402
from pptx.util import Emu                  # noqa: E402

EMU_PER_INCH = 914400.0
RING_IN = 0.23
LABEL_DX, LABEL_DY = -0.13, 0.215   # label origin -> marker centre, deck convention
# LOC-03's registration carries a uniform offset; everything else must be sub-mm.
FIXTURE_TOL_MM = {"LOC-01": 1.0, "LOC-02": 1.0, "LOC-03": 60.0}


def is_oval(sh):
    try:
        return "OVAL" in str(sh.auto_shape_type)
    except (ValueError, AttributeError):
        return False


def main():
    reg = json.loads((ROOT / "data" / "registration.json").read_text())
    sheets = json.loads((HERE / "field_sheets.json").read_text())
    prs = Presentation(ROOT / "input" / "TWYEAGLSHOPDWGEDITABLE_RevP05.pptx")

    out, failures = {}, []
    for slide_ix, loc in ((1, "LOC-01"), (2, "LOC-02"), (3, "LOC-03")):
        slide = prs.slides[slide_ix]
        T = dict(reg[loc])
        T["t"] = np.array(T["t"])

        def grid(x_in, y_in):
            return apply(T, [[x_in * EMU_PER_INCH, y_in * EMU_PER_INCH]])[0]

        rings = {}
        labels = {}
        for sh in slide.shapes:
            left, top = Emu(sh.left).inches, Emu(sh.top).inches
            if left > 11:                       # right-hand panel, not the map
                continue
            if is_oval(sh) and abs(Emu(sh.width).inches - RING_IN) < 0.01:
                w = Emu(sh.width).inches
                rings[(left + w / 2, top + w / 2)] = str(sh.line.color.rgb)
            elif sh.has_text_frame:
                txt = sh.text_frame.text.strip()
                if "/" in txt and len(txt) < 20:
                    labels[txt] = (left, top)

        pts = [grid(x, y) for x, y in rings]
        xs = [p[0] for p in pts]
        ys = [p[1] for p in pts]
        fx = load_fixtures(bbox=(min(xs) - 30, min(ys) - 30, max(xs) + 30, max(ys) + 30),
                           assets_only=True)
        lights = fx[fx.asset_type.str.contains("light", case=False)].reset_index(drop=True)

        rec = {}
        for row in sheets[loc]:
            light = row["light"]
            if light not in labels:
                failures.append(f"{loc} {light}: no label on the Rev P05 sheet")
                continue
            lx, ly = labels[light]
            want = (lx + LABEL_DX, ly + LABEL_DY)
            cx, cy = min(rings, key=lambda r: (r[0] - want[0]) ** 2 + (r[1] - want[1]) ** 2)
            g = grid(cx, cy)
            d = np.hypot(lights.x.values - g[0], lights.y.values - g[1])
            order = d.argsort()
            mm = float(d[order[0]]) * 1000
            rec[light] = dict(
                cx=round(cx, 4), cy=round(cy, 4), ring=rings[(cx, cy)],
                local_x=round(float(g[0]), 3), local_y=round(float(g[1]), 3),
                fixture_mm=round(mm, 1),
                second_nearest_m=round(float(d[order[1]]), 3),
                fixture_type=lights.iloc[order[0]].asset_type,
            )
            if mm > FIXTURE_TOL_MM[loc]:
                failures.append(f"{loc} {light}: {mm:.1f} mm from nearest light fixture")
        out[loc] = rec

    if failures:
        print("VERIFICATION FAILED")
        for f in failures:
            print("  -", f)
        raise SystemExit(1)

    (HERE / "marker_positions.json").write_text(json.dumps(out, indent=1) + "\n")
    n = sum(len(v) for v in out.values())
    worst = max(r["fixture_mm"] for v in out.values() for r in v.values())
    print(f"{n} field-sheet lights located; worst distance to a surveyed light fixture "
          f"{worst:.1f} mm")
    print("wrote", HERE / "marker_positions.json")


if __name__ == "__main__":
    main()
