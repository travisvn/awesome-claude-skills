#!/usr/bin/env python3
"""Reconcile the TWY E (E4-E6) AGL shop drawing against the issued field sheets.

Input : src_RevP06.pptx  +  field_sheets.json  (transcribed from the PDFs)
Output: TWY_E_AGL_Shop_Drawing_ZIA_P07.pptx

Rev P06 had dropped the works-action marker for three LOC-01 lights that the field
sheet records as secondary-cable affected (SBC102-02/027, TCCECH-03/035, TCCECH-03/018)
and had drawn them grey — "not affected" — while its own scope note still listed them
in scope. It also carried a red/green marker on TCC103-11/021 that matches no legend
entry and no field-sheet row. This script re-marks every asset strictly from the two
Yes/No columns of the field sheet:

    sec YES + base YES  ->  CORE OUT + NEW CABLE      outer red    / inner blue (duct) or orange (sawcut)
    sec YES + base NO   ->  NEW SEC. CABLE ONLY       outer blue (duct) or orange (sawcut) / inner green
    sec NO  + base NO   ->  not affected              outer grey ring, existing-asset dot only

LOC-02's shallow-base column is blank on the field sheet. It is confirmed as "NO"
(30.07.2026) and recorded as such in field_sheets.json, so those five lights mark as
new-secondary-cable-only, which is what the sheet already showed.

LOC-03's four shallow-base lights keep the base-retained + dummy-plate (amber/orange)
treatment already agreed on Rev P06 — the field sheet gives the Yes, not the method.
"""

import copy
import json
import pathlib

import legend

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

HERE = pathlib.Path(__file__).parent
SRC = HERE / "src_RevP06.pptx"
OUT = HERE / "TWY_E_AGL_Shop_Drawing_ZIA_P07.pptx"

GREY = "9AA0A6"      # field verified not affected / not on field sheet
RED = "CC0000"       # core out
BLUE = "0055CC"      # duct
ORANGE = "E8710A"    # sawcut
GREEN = "1E8E3E"     # new cable only
TEAL = "12A5B8"      # existing AGL asset dot

RING_D = Inches(0.23)
INNER_D = Inches(0.12)
DOT_D = Inches(0.064)
TOL = Inches(0.09)   # marker identification tolerance, ~0.23 m on the ground


# ---------------------------------------------------------------- shape helpers

def is_oval(sh):
    try:
        return "OVAL" in str(sh.auto_shape_type)
    except (ValueError, AttributeError):
        return False


def centre(sh):
    return (sh.left + sh.width // 2, sh.top + sh.height // 2)


def find_oval(slide, cx, cy, diameter, tol=TOL):
    """The single oval of `diameter` centred within `tol` of (cx, cy)."""
    hits = [
        sh for sh in slide.shapes
        if is_oval(sh)
        and abs(sh.width - diameter) < Inches(0.01)
        and abs(centre(sh)[0] - cx) < tol
        and abs(centre(sh)[1] - cy) < tol
    ]
    if len(hits) != 1:
        raise LookupError(
            f"expected 1 oval d={Emu(diameter).inches:.3f}\" near "
            f"({Emu(cx).inches:.3f},{Emu(cy).inches:.3f}), found {len(hits)}"
        )
    return hits[0]


def by_id(slide, shape_id):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise LookupError(f"no shape id={shape_id}")


def ring(shape, colour, width_pt=2.25):
    shape.fill.background()
    shape.line.color.rgb = RGBColor.from_string(colour)
    shape.line.width = Pt(width_pt)


def disc(shape, colour, width_pt=0.5):
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(colour)
    shape.line.color.rgb = RGBColor.from_string(colour)
    shape.line.width = Pt(width_pt)


def clone_oval(slide, template, cx, cy, diameter, colour):
    """Copy `template`'s XML so the new marker inherits its exact styling."""
    new = copy.deepcopy(template._element)
    template._element.addnext(new)
    sh = next(s for s in slide.shapes if s._element is new)
    max_id = max(s.shape_id for s in slide.shapes if s._element is not new)
    sh._element.nvSpPr.cNvPr.set("id", str(max_id + 1))
    sh._element.nvSpPr.cNvPr.set("name", f"Marker {max_id + 1}")
    sh.width = sh.height = diameter
    sh.left = cx - diameter // 2
    sh.top = cy - diameter // 2
    disc(sh, colour)
    return sh


def resize_about_centre(shape, diameter):
    cx, cy = centre(shape)
    shape.width = shape.height = diameter
    shape.left = cx - diameter // 2
    shape.top = cy - diameter // 2


def _font_of(run):
    return dict(
        size=run.font.size,
        name=run.font.name,
        bold=run.font.bold,
        rgb=run.font.color.rgb if run.font.color and run.font.color.type is not None else None,
    )


def set_lines(shape, lines):
    """Replace a text frame's paragraphs, reusing each paragraph's own formatting.

    Paragraph i keeps the font of the original paragraph i (title bold / detail
    regular in the totals tiles); extra paragraphs inherit the last original.
    """
    tf = shape.text_frame
    fonts = [_font_of(p.runs[0]) for p in tf.paragraphs if p.runs]
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        font = fonts[min(i, len(fonts) - 1)]
        r.font.size = font["size"]
        r.font.name = font["name"]
        r.font.bold = font["bold"]
        if font["rgb"] is not None:
            r.font.color.rgb = font["rgb"]


def sub(shape, old, new):
    hit = False
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            if old in r.text:
                r.text = r.text.replace(old, new)
                hit = True
    if not hit:
        raise LookupError(f"{old!r} not found in shape id={shape.shape_id}")


def cell_text(cell, text):
    tf = cell.text_frame
    proto = tf.paragraphs[0].runs[0]
    font = dict(size=proto.font.size, name=proto.font.name, bold=proto.font.bold,
                rgb=proto.font.color.rgb if proto.font.color and proto.font.color.type is not None else None)
    tf.clear()
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = font["size"]
    r.font.name = font["name"]
    r.font.bold = font["bold"]
    if font["rgb"] is not None:
        r.font.color.rgb = font["rgb"]


# ---------------------------------------------------------------------- checks

WANT = {
    (True, True, "Duct"): (RED, BLUE),        # core out + new cable, via duct
    (True, True, "Sawcut"): (RED, ORANGE),    # core out + new cable, via sawcut
    (True, False, "Duct"): (BLUE, GREEN),     # new secondary cable only, via duct
    (True, False, "Sawcut"): (ORANGE, GREEN), # new secondary cable only, via sawcut
}


def expected(loc, row):
    if loc == "LOC-03" and row["base"]:
        return ("F9AB00", ORANGE)             # base retained + dummy plate (agreed method)
    if not row["sec"] and not row["base"]:
        return (GREY, None)                   # field verified not affected
    return WANT[(row["sec"], bool(row["base"]), row["route"] or row["route_deck"])]


def audit(prs, sheets, positions):
    """Each field-sheet light carries the marker its Yes/No columns imply, at the
    surveyed position — checked against marker_positions.json, whose coordinates
    were themselves verified against the source DXF."""
    problems = []
    for slide_ix, loc in ((1, "LOC-01"), (2, "LOC-02"), (3, "LOC-03")):
        slide = prs.slides[slide_ix]
        for row in sheets[loc]:
            light = row["light"]
            p = positions[loc].get(light)
            if p is None:
                problems.append(f"{loc} {light}: no verified position")
                continue
            cx, cy = Inches(p["cx"]), Inches(p["cy"])
            want_ring, want_inner = expected(loc, row)
            try:
                r = find_oval(slide, cx, cy, RING_D, tol=Inches(0.02))
            except LookupError as e:
                problems.append(f"{loc} {light}: {e}")
                continue
            if str(r.line.color.rgb) != want_ring:
                problems.append(f"{loc} {light}: ring {r.line.color.rgb} != {want_ring}")
            if want_inner is None:
                continue
            try:
                inner = find_oval(slide, cx, cy, INNER_D, tol=Inches(0.02))
            except LookupError as e:
                problems.append(f"{loc} {light}: inner disc {e}")
                continue
            if str(inner.fill.fore_color.rgb) != want_inner:
                problems.append(f"{loc} {light}: inner {inner.fill.fore_color.rgb} != {want_inner}")
    return problems


# ------------------------------------------------------------------------ main

def main():
    sheets = json.loads((HERE / "field_sheets.json").read_text())
    prs = Presentation(SRC)
    s2, s3, s4, s5, s6 = (prs.slides[i] for i in (1, 2, 3, 4, 5))

    # ---- LOC-01 map: restore the three secondary-cable-only markers ----------
    # Rev P06 drew these three grey ("not affected") AND re-plotted them: its
    # grey rings sit 0.15-1.16 m from the nearest surveyed light, and in the
    # TCCECH-03/035 // 03/018 cluster the nearest and second-nearest fixtures
    # are within 1.1x of each other, so those positions cannot be salvaged by
    # snapping. Rev P05's markers land on the surveyed insertion points to
    # <=0.7 mm (checked against Z1-Z2-Z3-MTA_SEGMENTATION.dxf via
    # data/registration.json — see marker_positions.json), so the P05
    # positions are used and P06's off-fixture symbols are removed.
    ring_proto = find_oval(s2, Inches(4.680) + RING_D // 2, Inches(4.277) + RING_D // 2, RING_D)
    inner_proto = find_oval(s2, *centre(ring_proto), INNER_D)

    for cx, cy, d in [(3.662, 5.827, RING_D),    # grey ring 0.145 m off fixture
                      (4.514, 5.137, RING_D),    # grey ring 0.403 m off, ambiguous
                      (4.655, 5.200, RING_D),    # grey ring 1.161 m off, ambiguous
                      (4.515, 5.152, DOT_D),     # asset dot 0.297 m off, ambiguous
                      (4.618, 5.209, DOT_D)]:    # asset dot 0.854 m off, ambiguous
        sh = find_oval(s2, Inches(cx), Inches(cy), d, tol=Inches(0.02))
        sh._element.getparent().remove(sh._element)

    positions = json.loads((HERE / "marker_positions.json").read_text())["LOC-01"]
    for light, colour in (("SBC102-02/027", BLUE),      # sec YES / base NO / duct
                          ("TCCECH-03/035", ORANGE),    # sec YES / base NO / sawcut
                          ("TCCECH-03/018", ORANGE)):   # sec YES / base NO / sawcut
        p = positions[light]
        assert p["fixture_mm"] < 1.0, f"{light}: P05 position is {p['fixture_mm']} mm off a fixture"
        cx, cy = Inches(p["cx"]), Inches(p["cy"])
        r = clone_oval(s2, ring_proto, cx, cy, RING_D, colour)
        ring(r, colour)
        clone_oval(s2, inner_proto, cx, cy, INNER_D, GREEN)

    # ---- LOC-01 map: the 03/035 // 03/018 labels printed on top of each other
    # The two fittings are 0.49 m apart, so at 1:465 both markers and both
    # labels overlap and the text renders as garble. They carry an identical
    # works action, so one shared label loses nothing and reads cleanly.
    keep = next(sh for sh in s2.shapes
                if sh.has_text_frame and sh.text_frame.text.strip() == "TCCECH-03/035")
    drop = next(sh for sh in s2.shapes
                if sh.has_text_frame and sh.text_frame.text.strip() == "TCCECH-03/018")
    drop._element.getparent().remove(drop._element)
    set_lines(keep, ["TCCECH-03/035 + 03/018  (0.49 m apart)"])
    keep.width = Inches(2.05)

    # ---- LOC-01 map: TCC103-11/021 is not on the field sheet -----------------
    # Rev P06 gave it outer red + inner green, a combination in no legend.
    # Demote to the plain existing-asset symbol, as its three siblings already are.
    stray = find_oval(s2, Inches(6.158) + RING_D // 2, Inches(6.487) + RING_D // 2,
                      RING_D, tol=Inches(0.02))
    assert str(stray.line.color.rgb) == RED
    ring(stray, GREY)
    stray_inner = find_oval(s2, *centre(stray), INNER_D)
    resize_about_centre(stray_inner, DOT_D)
    disc(stray_inner, TEAL)

    # ---- LOC-01 map: normalise the four TCC103 labels -----------------------
    # P06 left them at 7.0 pt navy where every other label is 6.5 pt dark grey, which
    # reads as emphasis — the opposite of "not in scope". TCC103-12/022's box was also
    # 8.27" wide and overran the sheet.
    for sh in list(s2.shapes):
        if not sh.has_text_frame:
            continue
        if not sh.text_frame.text.strip().startswith("TCC103-"):
            continue
        sh.width = Inches(1.25)
        for run in sh.text_frame.paragraphs[0].runs:
            run.font.size = Pt(6.5)
            run.font.bold = True
            run.font.color.rgb = RGBColor.from_string("1F2937")

    # ---- LOC-01 scope note, strictly from Document_3 ------------------------
    set_lines(by_id(s2, 312), [
        "CORE OUT SHALLOW BASE + NEW CABLE: 11 No.   (SEC. CABLE YES / SHALLOW BASE YES)",
        "   • VIA DUCT (6): SBC102-02/024, 01/027, 02/025, 01/028, 02/026, 01/029",
        "   • VIA SAWCUT (5): TCCECH-04/034, 03/036, 04/035, 03/037, 04/036",
        "NEW SECONDARY CABLE ONLY — NO CORING: 3 No.   (SEC. CABLE YES / SHALLOW BASE NO)",
        "   • VIA DUCT (1): SBC102-02/027",
        "   • VIA SAWCUT (2): TCCECH-03/035, 03/018 — 0.49 m apart, one shared label; bases sound, verify at cut line",
        "FIELD VERIFIED NOT AFFECTED: SBC102-01/026 (SEC. CABLE NO / SHALLOW BASE NO — no works)",
        "FIELD SHEET TOTAL: 15 No. — 14 secondary cables, 11 core-outs, 0 dummy plates",
        "NOT ON FIELD SHEET — CONFIRM BEFORE WORKS: TCC103-11/021, 11/126, 12/021, 12/022",
        "   • secondary routes appear to cross the cut; not listed on Document_3 — shown grey, no action assumed",
        "ISOLATE CIRCUITS: SBC102.01/.02, TCCECH.03/.04",
        "SITE RECORD 23.07.2026: 9 cored + 2 reinstatement completed",
    ])

    # ---- LOC-02: shallow-base column blank on the sheet, confirmed "NO" -----
    sub(by_id(s3, 113),
        "NO SHALLOW BASES AFFECTED — NO ASSET LIES INSIDE THE MILLING POLYGON HERE",
        "NO SHALLOW BASES AFFECTED — COLUMN BLANK ON FIELD SHEET, CONFIRMED \u2018NO\u2019 "
        "30.07.2026; NO ASSET LIES INSIDE THE MILLING POLYGON")

    # ---- LOC-03: state the field-sheet total -------------------------------
    tf = by_id(s4, 311).text_frame
    src = by_id(s4, 311)
    lines = [p.text for p in tf.paragraphs if p.text.strip()]
    lines.insert(4, "FIELD SHEET TOTAL: 12 No. — all 12 secondary-cable affected; "
                    "4 shallow-base affected (YES)")
    set_lines(src, lines)

    # ---- revision block on all three sheets --------------------------------
    for slide, rev_id, src_id, meta_id, sheet, n in (
        (s2, 307, 308, 309, "Document_3", 15),
        (s3, 108, 109, 110, "Second_milling_area rev _1 (2nd table)", 5),
        (s4, 306, 307, 308, "Second_milling_area rev _1 (1st table)", 12),
    ):
        sub(by_id(slide, 2), "REV P06", "REV P07")          # sheet header
        sub(by_id(slide, rev_id), "REV P06 (ISSUED)", "REV P07 (ISSUED)")
        sub(by_id(slide, src_id), "— FIELD CONDITION GOVERNS",
            f"({n} No. LIGHTS LISTED) — FIELD CONDITION GOVERNS")
        sub(by_id(slide, meta_id), "Issued 28.07.2026",
            "Issued 28.07.2026 · Rev P07 30.07.2026")
        note = by_id(slide, {2: 315, 3: 116, 4: 314}[list(prs.slides).index(slide) + 1])
        notes = [p.text for p in note.text_frame.paragraphs if p.text.strip()]
        extra = (" Shallow-base column blank on that sheet; confirmed \u2018NO\u2019 30.07.2026."
                 if slide is s3 else "")
        notes.append(
            "7. Existing assets carry the symbol of their as-built class — see legend. "
            "Class is resolved from the AGL asset survey against the AGL duct-layout CAD; "
            "where an asset is labelled the label governs. Where a works marker is drawn "
            "the asset symbol sits beneath it, so legend counts include assets under "
            "works markers. Civil items (handhole / manhole / pits / RRM) plot about "
            "2.1 m from the surveyed insertion point — set out from survey, not from this "
            "sheet, and confirm the individual pit before works."
        )
        notes.append(
            f"6. Rev P07: every asset on this sheet re-marked directly from field sheet "
            f"{sheet} — outer ring from the Secondary-cable / Shallow-base columns, "
            f"inner disc from Duct / Sawcut. Assets not listed on the field sheet are "
            f"shown grey with no action assumed.{extra}"
        )
        notes.sort(key=lambda t: t[:2])
        set_lines(note, notes)

    # ---- base-layer label declutter -----------------------------------------
    # LOC-03 has a dense cluster where existing-asset labels printed on top of each
    # other and rendered as garble. The shifts are computed by declutter_labels.py
    # from the rendered layout (a label's box is 1.15" wide but its text is shorter,
    # so overlap is only visible once rendered) and replayed here so this script
    # alone reproduces the issued sheet. Base-layer labels only; no symbol moves.
    offsets_path = HERE / "label_offsets.json"
    if offsets_path.exists():
        offsets = json.loads(offsets_path.read_text())
        for slide, loc in ((s2, "LOC-01"), (s3, "LOC-02"), (s4, "LOC-03")):
            for text, dy in offsets.get(loc, {}).items():
                hits = [sh for sh in slide.shapes
                        if sh.has_text_frame and sh.text_frame.text.strip() == text
                        and sh.left is not None and sh.left < Inches(11)]
                if len(hits) != 1:
                    raise LookupError(f"{loc}: {text!r} matched {len(hits)} labels, "
                                      f"cannot apply its declutter offset")
                hits[0].top = hits[0].top + Inches(dy)

    # ---- legends: rebuilt so every plotted symbol is accounted for -----------
    # P06 had dropped P05's "AGL FEED MANHOLE / HANDHOLE" row (the magenta square is
    # still drawn on all three sheets) and reused the row to relabel the grey ring
    # "Dummy Plate". The milling/cut shading was never legended. And every existing
    # asset plots as one teal dot, so the sheets now carry an as-built asset key
    # naming each label prefix. Class names come from asset_key.py, not from here.
    square_tpl = copy.deepcopy(
        next(sh for sh in s4.shapes
             if legend._kind(sh) == "RECTANGLE"
             and sh.left is not None and sh.left >= Inches(11)
             and sh.width is not None and abs(sh.width - Inches(0.21)) < Inches(0.02))._element)
    legend_problems = []
    for slide, loc in ((s2, "LOC-01"), (s3, "LOC-02"), (s4, "LOC-03")):
        # give every existing asset the symbol of its as-built class, so the sheet
        # distinguishes a centreline light from a handhole without reading the label.
        # Classes come from asset_symbols.py; unclassified symbols keep the generic dot.
        for rec in legend.load_symbols(loc):
            cx, cy = Inches(rec["x"]), Inches(rec["y"])
            hits = [sh for sh in slide.shapes
                    if legend._kind(sh) is not None and sh.left is not None
                    and sh.width is not None and sh.width <= Inches(0.11)
                    and abs(legend._centre(sh)[0] - cx) < Inches(0.006)
                    and abs(legend._centre(sh)[1] - cy) < Inches(0.006)]
            if len(hits) != 1:
                raise LookupError(f"{loc}: {len(hits)} asset symbols at "
                                  f"({rec['x']},{rec['y']}), expected 1")
            legend.draw_asset(hits[0], rec["cls"])
        declared = legend.build(slide, loc, legend.load_asset_counts(loc), square_tpl)
        legend_problems += legend.audit(slide, loc, declared)

    # ---- title sheet: record the revision and the governing field sheets ----
    s1 = prs.slides[0]
    issued = next(sh for sh in s1.shapes
                  if sh.has_text_frame and "Issued 28.07.2026" in sh.text_frame.text)
    sub(issued, "Issued 28.07.2026", "Issued 28.07.2026 · Rev P07 30.07.2026")
    subtitle = next(sh for sh in s1.shapes
                    if sh.has_text_frame and "THREE (3) MILLING LOCATIONS" in sh.text_frame.text)
    set_lines(subtitle, [p.text for p in subtitle.text_frame.paragraphs if p.text.strip()] +
              ["REV P07 — SCOPE PER FIELD SHEETS Document_3 (LOC-01) AND "
               "Second_milling_area rev _1 (LOC-02 / LOC-03)"])

    # ---- consolidated table -------------------------------------------------
    tbl = next(sh for sh in s5.shapes if sh.has_table).table
    for r in (13, 14):   # TCCECH-03/035, 03/018
        cell_text(tbl.cell(r, 2), "NEW CABLE ONLY")
        cell_text(tbl.cell(r, 4),
                  "Field sheet: sec. cable YES / shallow base NO — no coring, "
                  "no dummy plate; full cable replacement via sawcut")
    for r, m in ((15, "2.0"), (16, "2.0"), (17, "2.0"), (18, "7.8")):
        cell_text(tbl.cell(r, 2), "NOT ON FIELD SHEET — CONFIRM")
        cell_text(tbl.cell(r, 3), "—")   # no route assumed where no action is assumed
        cell_text(tbl.cell(r, 4),
                  f"Not listed on Document_3. Secondary route appears to cross the cut "
                  f"({m} m in cut) — confirm scope before works; no action assumed")

    # ---- consolidated totals ------------------------------------------------
    # Counts exclude the four TCC103 assets throughout — they are on no field sheet.
    # An earlier draft let them into the dummy-plate headline (8 = 4 + 4) while every
    # other tile excluded them; check_all.py caught the inconsistency.
    cell_text_map = {
        12: "4",
        15: "31",
        102: "11",
    }
    for sid, val in cell_text_map.items():
        set_lines(by_id(s6, sid), [val])
    set_lines(by_id(s6, 13), ["Dummy plates",
                              "LOC-03 only  ·  4 TCC103 not counted"])
    set_lines(by_id(s6, 16), ["Secondary cable runs", "per field sheets; +4 TCC103"])
    set_lines(by_id(s6, 103), [
        "Saw-cut route runs",
        "LOC-01: 7  •  LOC-03: 4  •  +4 TCC103 not on field sheet  •  "
        "13.8 m of secondary cable inside the cut",
    ])
    sub(by_id(s6, 23), "PROVISIONAL — TCC103 (4 No.)", "NOT ON FIELD SHEET — TCC103 (4 No.)")
    set_lines(by_id(s6, 110), [
        "TCC103-11/021  —  not listed; cable may cross cut (2.0 m) — confirm",
        "TCC103-11/126  —  not listed; cable may cross cut (2.0 m) — confirm",
        "TCC103-12/021  —  not listed; cable may cross cut (2.0 m) — confirm",
        "TCC103-12/022  —  not listed; cable may cross cut (7.8 m) — confirm",
    ])

    prs.save(OUT)

    # ---- verify the saved file against the field sheets ---------------------
    positions = json.loads((HERE / "marker_positions.json").read_text())
    problems = audit(Presentation(OUT), sheets, positions) + legend_problems
    if problems:
        print("AUDIT FAILED")
        for p in problems:
            print("  -", p)
        raise SystemExit(1)
    total = sum(len(sheets[k]) for k in ("LOC-01", "LOC-02", "LOC-03"))
    print(f"audit OK — {total} field-sheet lights, every marker matches its Yes/No columns")
    print("wrote", OUT)


if __name__ == "__main__":
    main()
