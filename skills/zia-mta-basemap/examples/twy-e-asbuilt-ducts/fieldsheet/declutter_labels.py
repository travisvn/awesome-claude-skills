#!/usr/bin/env python3
"""Nudge overlapping base-layer asset labels apart until they read cleanly.

LOC-03 has a dense cluster where five pairs of existing-asset labels
(TCCECH.05/.06, EL.EBASE, TCC102) print on top of each other and render as garble.
The condition predates Rev P05 and is invisible to any check that only reads the
PowerPoint geometry, because each label box is a fixed 1.15" wide while the text
inside it is much shorter — overlap depends on the *rendered* text extent.

So this works from the rendered PDF: render, read the actual text-span boxes, shift
one label of each colliding pair vertically, re-render, repeat. It converges in a few
passes or reports what it could not resolve.

Rules, so this never changes what the drawing asserts:
  - Only base-layer labels move — dot notation, e.g. "TCCECH.05.033". Works-item
    labels (dash notation with a slash, e.g. "SBC102-01/038") are never touched, and
    neither is any symbol.
  - Total displacement per label is capped at MAX_SHIFT, about 0.4 m at sheet scale,
    so a label stays adjacent to the marker it annotates.

Run after reconcile_fieldsheet.py:  python declutter_labels.py
"""

import json
import pathlib
import re
import shutil
import subprocess
import sys
import tempfile

import fitz
from pptx import Presentation
from pptx.util import Emu, Inches

HERE = pathlib.Path(__file__).parent
DECK = HERE / "TWY_E_AGL_Shop_Drawing_ZIA_P07.pptx"
OFFSETS = HERE / "label_offsets.json"
SHEETS = {1: "LOC-01", 2: "LOC-02", 3: "LOC-03"}
MAP_RIGHT = 11.05
STEP = Inches(0.075)
MAX_SHIFT = Inches(0.16)
MAX_PASSES = 6
BASE_LABEL = re.compile(r"^[A-Z_0-9]+(\.[A-Za-z0-9_]+)+$")   # dot notation only


def render(pptx, outdir):
    subprocess.run(["soffice", "--headless", "--norestore", "--convert-to", "pdf",
                    "--outdir", str(outdir), str(pptx)],
                   check=True, capture_output=True, timeout=300,
                   env={"HOME": "/tmp/lohome", "PATH": "/usr/bin:/bin"})
    return next(pathlib.Path(outdir).glob("*.pdf"))


def collisions(pdf):
    """[(page_index, textA, boxA, textB, boxB, overlap_x, overlap_y)] for map labels."""
    doc = fitz.open(pdf)
    out = []
    for pno in SHEETS:
        spans = [s for b in doc[pno].get_text("dict")["blocks"] if b["type"] == 0
                 for line in b["lines"] for s in line["spans"]
                 if s["text"].strip() and s["bbox"][2] < MAP_RIGHT * 72
                 and any(c in s["text"] for c in "./")]
        for i, a in enumerate(spans):
            for b in spans[i + 1:]:
                ox = min(a["bbox"][2], b["bbox"][2]) - max(a["bbox"][0], b["bbox"][0])
                oy = min(a["bbox"][3], b["bbox"][3]) - max(a["bbox"][1], b["bbox"][1])
                if ox > 6 and oy > 3:
                    out.append((pno, a["text"].strip(), a["bbox"],
                                b["text"].strip(), b["bbox"], ox, oy))
    return out


def label_shapes(slide):
    out = {}
    for sh in slide.shapes:
        if sh.left is None or Emu(sh.left).inches >= MAP_RIGHT or not sh.has_text_frame:
            continue
        t = sh.text_frame.text.strip()
        if t:
            out.setdefault(t, []).append(sh)
    return out


def main():
    if not DECK.exists():
        raise SystemExit(f"{DECK} not found — run reconcile_fieldsheet.py first")
    work = pathlib.Path(tempfile.mkdtemp())
    shutil.copy(DECK, work / DECK.name)
    deck = work / DECK.name
    shifted = {}

    before = collisions(render(deck, work / "r0"))
    print(f"{len(before)} overlapping label pairs before:")
    for pno, ta, _a, tb, _b, ox, oy in before:
        print(f"   {SHEETS[pno]}  {ta!r} / {tb!r}   {ox / 72:.2f}\" x {oy / 72:.2f}\"")
    if not before:
        print("nothing to do")
        return 0

    for p in range(1, MAX_PASSES + 1):
        cols = collisions(render(deck, work / f"r{p}"))
        if not cols:
            break
        prs = Presentation(deck)
        moved = 0
        for pno, ta, box_a, tb, box_b, _ox, _oy in cols:
            shapes = label_shapes(prs.slides[pno])
            # move whichever of the pair is a base-layer label and sits lower
            cands = [(t, box) for t, box in ((ta, box_a), (tb, box_b))
                     if BASE_LABEL.match(t) and t in shapes and len(shapes[t]) == 1]
            if not cands:
                continue
            cands.sort(key=lambda tb_: -tb_[1][1])       # lowest on the sheet first
            text, box = cands[0]
            other = box_b if text == ta else box_a
            direction = 1 if box[1] >= other[1] else -1  # push away from the other span
            used = shifted.get((pno, text), 0)
            if used + STEP > MAX_SHIFT:
                continue
            sh = shapes[text][0]
            sh.top = sh.top + direction * STEP
            shifted[(pno, text)] = used + STEP
            moved += 1
        if not moved:
            print(f"pass {p}: nothing left that may be moved")
            break
        prs.save(deck)
        print(f"pass {p}: nudged {moved} label(s)")

    after = collisions(render(deck, work / "final"))
    print(f"\n{len(after)} overlapping label pairs after:")
    for pno, ta, _a, tb, _b, ox, oy in after:
        print(f"   {SHEETS[pno]}  {ta!r} / {tb!r}   {ox / 72:.2f}\" x {oy / 72:.2f}\"")
    for (pno, text), amount in sorted(shifted.items()):
        print(f"   moved {text!r} on {SHEETS[pno]} by {Emu(amount).inches:+.3f}\"")

    if len(after) >= len(before):
        print("\nno improvement — leaving the deck untouched")
        return 1

    # Persist the shifts so reconcile_fieldsheet.py reproduces them and the pipeline
    # stays idempotent — otherwise the next reconcile run would undo this pass.
    saved = json.loads(OFFSETS.read_text()) if OFFSETS.exists() else {}
    for (pno, text), amount in shifted.items():
        loc = SHEETS[pno]
        saved.setdefault(loc, {})
        saved[loc][text] = round(saved[loc].get(text, 0.0) + Emu(amount).inches, 4)
    OFFSETS.write_text(json.dumps(saved, indent=1, sort_keys=True) + "\n")
    shutil.copy(deck, DECK)
    print(f"\nwrote {DECK}\nwrote {OFFSETS}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
